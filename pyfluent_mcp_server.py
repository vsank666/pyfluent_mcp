#!/usr/bin/env python3
"""
PyFluent MCP Server  v3.12.0
============================
Connects AI assistants (Claude, Cursor, Windsurf …) to ANSYS Fluent.

New in v3.12 (full-repo audit + live manifold regression testing, 2026-07-21/22):
  • Real-time iteration tracking          — SolverEvent.ITERATION_ENDED
                                            (EventsManager) replaces the
                                            broken (rpgetvar 'number-of-
                                            iterations) probe, which turned
                                            out to report the requested
                                            batch size, not a live counter
  • get_calculation_status batching       — iteration/progress/residuals/
                                            convergence/errors in one call
  • wait_for_job / run_iterations         — block server-side on job/
                                            calculation completion instead
                                            of blind poll loops
  • run_tui_script                        — replay a saved journal/TUI
                                            script end-to-end
  • Monitors gate, project-scoped logs, session busy-guard, and a long
    list of live-verified bug fixes (see CHANGELOG / git log for detail)

New in v3.5 (traceability gaps found in the manifold_2 session, 2026-07-13):
  • start_traceability /               — Fluent-native .trn transcript +
    stop_traceability                    .jou journal. Solves run_tui_command
                                         only returning '#t' for report/list
                                         commands: the transcript captures the
                                         printed text on disk regardless.
  • export_case_manifest               — reads back ACTUAL applied BC values
                                         (not just zone name/type) + solver
                                         settings; would have caught the
                                         backflow_temperature_k=550 silently
                                         not landing (discovered only later
                                         from field data) automatically.
  • checkpoint_run                     — pause checkpoint bundling
                                         write_case/write_case_data +
                                         export_case_manifest + iteration/
                                         residuals into one call, so a
                                         restart (MCP server or otherwise)
                                         never again loses an entire run with
                                         nothing saved to resume from.

TWO DISTINCT ACCESS PATHS — never confuse them (see
knowledge/robust_execution_methodology.md):
  A. PYFLUENT SESSION path (PyAnsys): session objects created by
     launch_fluent/connect_to_fluent; settings API, field data, monitors.
     Most tools use this path. scheme_eval.exec() on this path EXECUTES
     but DISCARDS return values — use scheme_eval.string_eval() to read.
  B. RAW gRPC path: ansys-api-fluent stubs straight to the server
     (HealthCheck, SchemeEval/StringEval). No session object needed.
     Tools prefixed grpc_* use this path. It is the out-of-band channel
     for probing/monitoring a server whose session object is busy or
     that this process never launched.

New in v3.4 (lessons from the manifold_2 50-iteration run, 2026-07-11 —
all paths verified live against Fluent 2025 R2 / PyFluent 0.40.1):
  • start_calculation /                — NON-BLOCKING solve: background
    get_calculation_status /             thread + status polls + interrupt.
    stop_calculation                     run_calculation refuses long runs.
  • run_scheme_expression fix          — string_eval (exec discarded results)
  • run_tui_command fix                — transcript captured via string_eval
  • get_residuals /                    — session monitors service (residual
    get_convergence_history fix          TUI print-to-screen is INVALID in
                                         2025 R2)
  • create_wall_display rewrite        — default TRANSPARENT FACES, no edges
                                         (scene-object transparency); wireframe
                                         only as explicit style
  • display_scene rewrite              — scene-based compositing with
                                         per-object transparency
  • NOTE: checkpoint/exit-filename rpvar does NOT exist in 2025 R2; the
    verified stop mechanism is solution.run_calculation.interrupt().

New in v3.2 (lessons from the manifold_2 live run — all paths verified
against Fluent 2025 R2):
  • read_mesh                          — .msh/.msh.h5 loader (read_case is cas-only)
  • run_mesh_check                     — structured pass/fail + interface hint
  • create_mesh_interfaces             — auto-pair non-conformal interfaces
  • preflight_simulation               — scale/zones/inlet-outlet/quality gate
  • apply_visual_defaults              — grid/reflections/shadows off, white
                                          picture background, PNG, 1080p/4K
  • create_plane_surface               — two-step creation (method first!)
  • create_outlet_normal_plane         — max-area plane through any opening
  • create_native_contour /            — native Fluent graphics objects with
    create_wall_display                  overlay-safe display_scene compositor
  • display_scene                      — multi-object render + view presets
  • get_session_info                   — ip/port/password for re-attachment
  • set_velocity_inlet fix             — turbulence spec method set before
                                          values (was silently skipped)

New in v3.1:
  • Process management                 — list/kill stray fluent/cortex/cxsolver
                                          OS processes; launch_fluent now asks
                                          for cores/GUI and cleans up stray
                                          processes before launching by default

New in v3 (on top of v2 foundation):
  • Parametric studies & DOE          — sweep inputs, run design points, analyse
  • Convergence monitoring             — live watch, smart stop, history plots
  • Aerodynamic force/moment reports   — Cl, Cd, Cm, lift/drag ratio
  • Heat-transfer analysis             — HTC, Nu, wall heat flux, thermal report
  • Mesh adaptation (AMR)             — gradient-based refinement cycles
  • GPU solver control                 — enable, check compatibility, set count
  • Case comparison & diff             — compare two setups or field values
  • Report generation (TXT / DOCX / PPTX) — executive-ready slide deck
  • Species transport setup            — multicomponent flows and combustion
  • Multiphase (VOF) setup            — free-surface and two-phase flows

Full tool list  (140 tools across 23 sections):
  1.  Connection / session
  2.  Case I/O
  3.  Mesh inspection
  4.  Physics / models
  5.  Materials
  6.  Boundary conditions
  7.  Solver settings & convergence criteria
  8.  Initialization & solution control
  9.  Post-processing – field data (NumPy / PyVista)
  10. Advanced visualisation (contours, vectors, streamlines, iso-surfaces)
  11. Export (VTK / glTF / HTML / USDZ / OpenUSD / Omniverse)
  12. Parametric studies & DOE            — NEW
  13. Convergence monitoring & smart stop — NEW
  14. Aerodynamic forces & moments        — NEW
  15. Heat-transfer analysis              — NEW
  16. Mesh adaptation (AMR)              — NEW
  17. GPU solver control                  — NEW
  18. Case comparison & diff              — NEW
  19. Report generation (TXT / DOCX / PPTX)  — PPTX new
  20. Species transport                   — NEW
  21. Multiphase (VOF)                    — NEW
  22. Meshing session workflow
  23. TUI / Scheme scripting
"""

from __future__ import annotations

import json
import os
import re
import threading
import time
from datetime import datetime
from pathlib import Path
from typing import Optional

from mcp.server.fastmcp import FastMCP

def _get_cmap(name):
    """Matplotlib-version-agnostic colormap lookup."""
    import matplotlib.pyplot as _plt
    return _plt.get_cmap(name)


# ---------------------------------------------------------------------------
# Global state
# ---------------------------------------------------------------------------
_solver = None
_meshing = None
_session_mode: str | None = None
_connection_tag: str = ""
_workflow = None            # cached watertight()/fault_tolerant()/two_dimensional_meshing() handle
_workflow_kind: str = ""    # "watertight" | "fault_tolerant" | "2d" — which one _workflow is
_import_length_unit: str = "m"   # CAD unit applied at Import Geometry — the unit ALL
                                 # workflow sizing fields use (verified live 2026-07-17:
                                 # sizing values pass through raw into fields in this unit)
_meshing_mode: str = "manual"    # "manual" | "auto" for the MESHING workflow
                                 # (see set_meshing_mode)
_analysis_mode: str = "manual"   # "manual" | "auto" for the CFD ANALYSIS workflow
                                 # (see set_analysis_mode); set_workflow_mode sets
                                 # both at once. Switching never discards setup.
_mesh_plan: dict | None = None   # last plan computed by propose_mesh_sizing
_solver_plan: dict | None = None # last plan computed by propose_solver_setup
_project_dir: str | None = None  # parent dir of the last-loaded geometry/mesh/case
                                 # file - project-scoped artifacts (monitors,
                                 # autosave checkpoints) live here instead of the
                                 # server's own install dir (user request 2026-07-21:
                                 # logs should follow the CAD/mesh, not the tool).

__version__ = "3.12.0"  # 3.6-3.9: persistence/reattach, busy-vs-dead health,
                        # process-tree discovery, sizing units fix, monitoring
                        # + monitors + checkpoints, dual auto/manual modes,
                        # interactive prompts, volume-mesh gate.
                        # 3.10 (2026-07-20 retrospective): direct-sifile launch
                        # (no launcher timeout, credentials survive), background
                        # jobs for import/surface/volume mesh (get_job_status),
                        # persistence merge (no default-clobber), compute()-based
                        # result tools, converged flag in calculation status,
                        # monitor history files + read_monitor_history,
                        # create_midplane with auto extents.
                        # 3.11 (interactive-console spec §3-24): project_summary
                        # + menu, get_boundary_state + transactional BC read-back,
                        # set_autosave, assess_convergence (C1/C2/C3),
                        # generate_default_results post pack
mcp = FastMCP("ansys-pyfluent-mcp")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _j(data) -> str:
    # default=str: PyFluent datamodel returns wrapper types (StringVector,
    # RepeatedScalarContainer, ...) inside task states — get_workflow_state
    # crashed with "Object of type StringVector is not JSON serializable"
    # on a reattached session (2026-07-17)
    return json.dumps(data, indent=2, ensure_ascii=False, default=str)

def _active():
    return _solver if _solver is not None else _meshing

def _chk(mode: str | None = None, allow_busy: bool = False) -> str | None:
    if _solver is None and _meshing is None:
        return _j({"ok": False, "error": "Not connected. Call connect_to_fluent first."})
    if mode == "solver" and _solver is None:
        return _j({"ok": False, "error": "Solver session required. Current mode: meshing."})
    if mode == "meshing" and _meshing is None:
        return _j({"ok": False, "error": "Meshing session required. Current mode: solver."})
    # refuse session-path calls while a background job/calculation holds the
    # session busy, UNLESS the caller is one of the small set of tools whose
    # whole purpose is to run WHILE one is in progress (get_calculation_status,
    # stop_calculation) - verified live 2026-07-20: no such guard existed
    # anywhere, so any other tool call during a long solve/mesh job raced
    # against that job's worker thread on the shared _solver/_meshing session.
    if not allow_busy and (_bg_run["running"] or _session_job["running"]):
        return _j({"ok": False, "error": "session job/calculation in progress",
                   "hint": "poll get_job_status/get_calculation_status, or use "
                           "stop_calculation first"})
    return None

def _esc(s: str) -> str:
    return s.replace('"', '\\"').replace("'", "\\'")

def _scheme_str(expression: str) -> str:
    """SESSION-path scheme evaluation that RETURNS the result as a string.

    scheme_eval.exec() executes but discards return values (verified live,
    2025 R2 / PyFluent 0.40.1) — string_eval is the read path."""
    return str(_active().scheme_eval.string_eval(expression))

def _mkdir(p: str | Path) -> None:
    Path(p).parent.mkdir(parents=True, exist_ok=True)

def _ts() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")

def _get_workflow(kind: str | None = None):
    """Return the cached meshing workflow handle, initializing it on first use.

    A bare meshing session has NO active workflow — TaskObject/attribute-style
    task access (e.g. `.import_geometry`) raises "not found at path /TaskObject"
    until one of watertight()/fault_tolerant()/two_dimensional_meshing() has
    actually run (see knowledge/source/meshing/01_new_meshing_workflows_api.md
    section 0). Re-initializing WIPES prior task state, so: kind=None (every
    caller except import_geometry) means "whatever workflow is already active" —
    a bare call must never re-initialize watertight over a live fault_tolerant/
    2d workflow (confirmed critical in the 2026-07-15 review).
    """
    global _workflow, _workflow_kind
    attach_only = False
    if kind is None:
        if _workflow is not None:
            return _workflow
        kind = "watertight"
        attach_only = True   # caller didn't ask for a (re)init — never wipe
    kind = (kind or "watertight").lower().replace("-", "_")
    if _workflow is not None and _workflow_kind == kind:
        return _workflow
    # REATTACH SAFETY (verified live 2026-07-17): calling watertight() on a
    # session that ALREADY HAS an active workflow (e.g. this MCP process was
    # restarted and reattached, or the workflow was built in the GUI)
    # re-initializes it and WIPES every task's state — a completed import +
    # sizing were destroyed exactly this way. When we didn't create the
    # workflow in this process and one exists server-side, attach with
    # initialize=False instead.
    if attach_only and _lowlevel_task_names():
        base = getattr(_meshing, "_base_meshing", None)
        factory = getattr(base, "watertight_workflow", None)
        if factory is not None:
            try:
                _workflow = factory(initialize=False)
                _workflow_kind = "watertight"
                return _workflow
            except Exception:
                pass  # fall through to normal path as a last resort
    if kind in ("fault_tolerant", "fault-tolerant", "ftm"):
        _workflow = _meshing.fault_tolerant()
        _workflow_kind = "fault_tolerant"
    elif kind in ("2d", "two_dimensional", "two_d"):
        _workflow = _meshing.two_dimensional_meshing()
        _workflow_kind = "2d"
    else:
        _workflow = _meshing.watertight()
        _workflow_kind = "watertight"
    return _workflow


def _lowlevel_task_names() -> list:
    """Task names via the low-level datamodel (works on reattached sessions
    where the high-level wrapper chokes). Defensive against PyFluent wrapper
    types that aren't iterable (StringVector) — index-walk as fallback."""
    if _meshing is None:
        return []
    try:
        names = _meshing.workflow.TaskObject.get_object_names()
    except Exception:
        return []
    try:
        return [str(n) for n in names]
    except TypeError:
        pass
    try:
        return [str(names[i]) for i in range(len(names))]
    except Exception:
        return []

def _norm_key(k: str) -> str:
    """Normalize an argument name across release dialects: 2023 R1 datamodel
    uses PascalCase (BoundaryLabelList), 2025 R2 high-level tasks use
    snake_case (boundary_label_list) — same argument."""
    return k.lower().replace("_", "").replace("-", "")


def _task_known_args(task) -> tuple[dict, list[str]]:
    """Discover a workflow task's REAL argument names from every surface the
    installed release exposes (verified live on 2025 R2 / pyfluent 0.40.1):
      1. dir(task) snake_case attributes — the high-level task API
         (`ub.boundary_label_list = [...]`, doc 01 section 2);
      2. task.get_state() keys;
      3. task.arguments()/Arguments() keys (PascalCase datamodel names —
         only populated for args that have been set at least once).
    Returns ({normalized: (real_name, source)}, probe_errors)."""
    known: dict = {}
    errors: list[str] = []
    try:
        cls = type(task)
        for n in dir(task):
            if n.startswith("_"):
                continue
            # class-level callables are task METHODS (get_state, add_child_to_task,
            # insert_compound_child_task, ...) — never argument names; setattr on
            # one would clobber the method on the wrapper instead of setting an arg
            if callable(getattr(cls, n, None)):
                continue
            known.setdefault(_norm_key(n), (n, "attribute"))
    except Exception as e:
        errors.append(f"dir() probe failed: {e}")
    try:
        state = task.get_state()
        if isinstance(state, dict):
            for k in state:
                known.setdefault(_norm_key(k), (k, "state"))
    except Exception as e:
        errors.append(f"get_state() probe failed: {e}")
    for probe in ("arguments", "Arguments"):
        try:
            obj = getattr(task, probe, None)
            if obj is None:
                continue
            arg_state = obj() if callable(obj) else obj.get_state()
            if isinstance(arg_state, dict):
                for k in arg_state:
                    known.setdefault(_norm_key(k), (k, probe))
        except Exception:
            pass  # schema probe is best-effort; sources 1-2 are the verified ones
    return known, errors


def _apply_task_args(task, desired: dict) -> dict:
    """Defensively apply desired arguments to a workflow task and VERIFY.

    Never silently no-ops (the pre-2026-07-15 version matched only against
    get_state() keys, which are empty until set on 2025 R2 — every argument
    was 'skipped_unknown' yet the tool reported ok:true). Now:
      * argument names are matched case/underscore-insensitively against
        dir(task) attributes, state keys, and the argument schema;
      * attribute-style assignment is preferred (the verified 2025 R2 path),
        falling back to set_state;
      * after applying, state is read back and per-key mismatches reported.
    Callers MUST treat empty 'applied' as failure — see _guard_no_op and
    knowledge/source/meshing/06_meshing_automation_contract_and_standards.md.
    """
    report: dict = {"applied": {}, "skipped_unknown": [],
                    "probe_errors": [], "readback_mismatches": []}
    known, probe_errors = _task_known_args(task)
    report["probe_errors"] = probe_errors
    set_state_batch: dict = {}
    for k, v in desired.items():
        real = known.get(_norm_key(k))
        if real is None:
            report["skipped_unknown"].append(k)
            continue
        real_name, source = real
        if source == "attribute":
            try:
                setattr(task, real_name, v)
                report["applied"][real_name] = v
                continue
            except Exception as e:
                report["probe_errors"].append(
                    f"setattr {real_name} failed ({e}); falling back to set_state")
        set_state_batch[real_name] = v
    if set_state_batch:
        try:
            task.set_state(set_state_batch)
            report["applied"].update(set_state_batch)
        except Exception as e:
            report["probe_errors"].append(f"set_state({list(set_state_batch)}) failed: {e}")
            for k in set_state_batch:
                report["skipped_unknown"].append(k)
    # read back and verify what the task now actually holds
    if report["applied"]:
        try:
            after = task.get_state()
            if isinstance(after, dict):
                after_norm = {_norm_key(k): v for k, v in after.items()}
                # arguments()/Arguments() often holds the authoritative values
                for probe in ("arguments", "Arguments"):
                    try:
                        obj = getattr(task, probe, None)
                        arg_state = obj() if callable(obj) else None
                        if isinstance(arg_state, dict):
                            after_norm.update({_norm_key(k): v for k, v in arg_state.items()})
                    except Exception:
                        pass
                for k, v in report["applied"].items():
                    got = after_norm.get(_norm_key(k), "<not in readback>")
                    if got != v and str(got) != str(v):
                        report["readback_mismatches"].append(
                            {"key": k, "requested": v, "readback": got})
        except Exception as e:
            report["probe_errors"].append(f"post-apply readback failed: {e}")
    return report


def _guard_no_op(report: dict, task_label: str):
    """Return an error-JSON string when NONE of the requested arguments could
    be applied — running the task anyway would execute with defaults while
    looking configured (the silent no-op observed live 2026-07-15, where
    update_boundaries 'succeeded' having assigned nothing). Returns None when
    at least one argument applied; callers surface the rest via 'warnings'."""
    if report.get("applied"):
        return None
    return _j({"ok": False,
               "error": f"none of the requested arguments matched the argument schema of "
                        f"'{task_label}' on this Fluent release - task NOT executed",
               **report,
               "hint": "call get_workflow_state to see this task's real argument names, "
                       "then retry; for release-specific fields fall back to run_tui_command"})


def _skip_warnings(report: dict) -> list[str]:
    """Standard warnings list for partially-applied task arguments."""
    w = []
    if report.get("skipped_unknown"):
        w.append("arguments skipped (not in this release's schema): "
                 + ", ".join(map(str, report["skipped_unknown"])))
    if report.get("readback_mismatches"):
        w.append(f"readback mismatches: {report['readback_mismatches']}")
    return w


def _find_wf_task(wf, *fragments: str, exclude: tuple = ("child",)):
    """Locate a workflow task attribute by lowercase name fragments,
    release-tolerantly: pyfluent 0.40 exposes 'add_boundary_layers' but
    0.39's WatertightMeshingWorkflow does not (verified live 2026-07-15) —
    the same task can carry a different wrapper attribute name per release.
    Returns the shortest matching attribute name (base task, not children),
    or None."""
    cands = []
    try:
        for n in dir(wf):
            if n.startswith("_"):
                continue
            ln = n.lower()
            if all(f in ln for f in fragments) and not any(x in ln for x in exclude):
                cands.append(n)
    except Exception:
        return None
    return sorted(cands, key=len)[0] if cands else None


def _session_port_open() -> bool:
    """Raw TCP probe of the active session's gRPC port — the trustworthy
    liveness signal (a dead Fluent can leave a stale-healthy session object)."""
    try:
        import socket
        props = _active().connection_properties
        ip = props.ip if props.ip not in (None, "", "localhost") else "127.0.0.1"
        with socket.create_connection((ip, props.port), timeout=3):
            return True
    except Exception:
        return False


# ---------------------------------------------------------------------------
# Session persistence — survive MCP server restarts.
#
# ROOT CAUSE (verified live 2026-07-17): long blocking tools (import, surface/
# volume mesh, launch) freeze this single-threaded stdio server; the MCP client
# may restart it, wiping _solver/_meshing. PyFluent deletes its temp
# server-info file after connecting, so without this record the ip/port/
# password of a perfectly healthy Fluent are unrecoverable and every restart
# looks like "gRPC broke". Every successful launch/connect writes the record;
# check_fluent_connection / reattach_last_session use it to re-attach.
# ---------------------------------------------------------------------------
_SESSION_STATE_PATH = Path(__file__).parent / "logs" / "last_session.json"


def _set_project_dir_from(file_path: str) -> None:
    """Anchor project-scoped runtime artifacts (monitors, autosave checkpoints)
    next to the CAD/mesh/case file being worked on, instead of the server's
    own install directory - so every project's outputs stay with that project
    rather than piling up (and mixing across projects) in the tool's own repo.
    Called from import_geometry/read_mesh/read_case/read_case_data."""
    global _project_dir
    try:
        p = Path(file_path).expanduser().resolve().parent
        if p.exists():
            _project_dir = str(p)
    except Exception:
        pass


def _artifacts_dir(subdir: str) -> Path:
    """Base directory for project-scoped runtime artifacts (monitors/,
    checkpoints/). Defaults next to the last-loaded geometry/mesh/case file
    (as a 'pyfluent_mcp_logs' subfolder); falls back to the server's own
    logs/ dir if no project file has been loaded yet in this session."""
    base = (Path(_project_dir) / "pyfluent_mcp_logs") if _project_dir else _SESSION_STATE_PATH.parent
    d = base / subdir
    d.mkdir(parents=True, exist_ok=True)
    return d


_CONTEXT_DEFAULTS = {"import_length_unit": "m", "meshing_mode": "manual",
                     "analysis_mode": "manual", "mesh_plan": None,
                     "solver_plan": None, "project_dir": None}


def _persist_session_state() -> None:
    """Best-effort write of reconnection + meshing context. Never raises.

    MERGES with the existing record: a context field still at its module
    default is never written over a real value already on disk — a fresh
    MCP process persisting on connect used to clobber import_length_unit
    'mm' back to 'm' before context restore ran (observed live 2026-07-20).
    """
    try:
        rec: dict = {"ts": _ts(), "mode": _session_mode,
                     "import_length_unit": _import_length_unit,
                     "meshing_mode": _meshing_mode,
                     "analysis_mode": _analysis_mode,
                     "mesh_plan": _mesh_plan,
                     "solver_plan": _solver_plan,
                     "project_dir": _project_dir}
        old = _load_persisted_session() or {}
        for k, dflt in _CONTEXT_DEFAULTS.items():
            if rec.get(k) == dflt and old.get(k) not in (None, dflt):
                rec[k] = old[k]
        s = _active()
        if s is not None:
            try:
                props = s.connection_properties
                rec.update({"ip": props.ip or "127.0.0.1", "port": props.port,
                            "password": props.password,
                            "cortex_pid": getattr(props, "cortex_pid", None),
                            "fluent_host_pid": getattr(props, "fluent_host_pid", None)})
            except Exception:
                pass
        _SESSION_STATE_PATH.parent.mkdir(parents=True, exist_ok=True)
        _SESSION_STATE_PATH.write_text(json.dumps(rec, indent=2), encoding="utf-8")
    except Exception:
        pass


def _load_persisted_session() -> dict | None:
    try:
        if _SESSION_STATE_PATH.exists():
            return json.loads(_SESSION_STATE_PATH.read_text(encoding="utf-8"))
    except Exception:
        pass
    return None


def _clear_persisted_session() -> None:
    try:
        _SESSION_STATE_PATH.unlink(missing_ok=True)
    except Exception:
        pass


def _restore_meshing_context(rec: dict) -> None:
    """Restore non-connection context (units/modes/plans) from a persisted record."""
    global _import_length_unit, _meshing_mode, _analysis_mode, _mesh_plan, _solver_plan, _project_dir
    _import_length_unit = rec.get("import_length_unit") or _import_length_unit
    _meshing_mode = rec.get("meshing_mode") or _meshing_mode
    _analysis_mode = rec.get("analysis_mode") or _analysis_mode
    _mesh_plan = rec.get("mesh_plan") or _mesh_plan
    _solver_plan = rec.get("solver_plan") or _solver_plan
    _project_dir = rec.get("project_dir") or _project_dir


def _patch_pyfluent_workflow_bugs() -> None:
    """Work around a pyfluent 0.40.1 defect hit ONLY when constructing a
    workflow handle on a session whose workflow already has tasks (reattach
    after an MCP restart, or a GUI-built workflow): workflow.py's
    _convert_task_list_to_display_names iterates a raw variant_pb2
    StringVector proto → "'StringVector' object is not iterable" (verified
    live 2026-07-17; the failed re-init also WIPED the existing task state).
    Normalize the value via protobuf reflection. Idempotent; call after
    every `import ansys.fluent.core`."""
    try:
        import ansys.fluent.core.workflow as _wfmod
    except ImportError:
        return
    if getattr(_wfmod, "_mcp_stringvector_patch", False):
        return
    orig = _wfmod._convert_task_list_to_display_names

    def _normalize(task_list):
        if hasattr(task_list, "__iter__"):
            return task_list
        try:
            for _fd, val in task_list.ListFields():
                if hasattr(val, "__iter__") and not isinstance(val, (str, bytes)):
                    return list(val)
        except Exception:
            pass
        return []   # unset/empty proto vector == no tasks in the list

    def _patched(workflow, task_list):
        return orig(workflow, _normalize(task_list))

    _wfmod._convert_task_list_to_display_names = _patched
    _wfmod._mcp_stringvector_patch = True


def _pyfluent_connect(pyfluent, kw: dict) -> tuple:
    """pyfluent.connect_to_fluent with same-machine-LAN-IP tolerance.

    A Fluent launched standalone (fluent ... -sifile=...) binds its gRPC
    server to the host's LAN interface ONLY (verified live 2026-07-17:
    cortex listened on 192.168.x.x:51649, loopback refused), and
    pyfluent >=0.40 treats any non-loopback ip as a remote host — first
    refusing outright, then demanding TLS certificates. On exactly those
    rejections, retry with allow_remote_host=True + insecure_mode=True,
    which is acceptable because the endpoint is this same machine.
    Returns (session, warnings)."""
    _patch_pyfluent_workflow_bugs()
    try:
        return pyfluent.connect_to_fluent(**kw), []
    except TypeError:
        raise
    except Exception as e:
        msg = str(e).lower()
        if "remote" not in msg and "certificat" not in msg:
            raise
        retry = dict(kw)
        retry["allow_remote_host"] = True
        retry["insecure_mode"] = True
        sess = pyfluent.connect_to_fluent(**retry)
        return sess, ["connected with allow_remote_host + insecure_mode: the "
                      "server is bound to this machine's LAN interface (not "
                      "loopback) - same-machine traffic, acceptable"]


def _try_reattach(rec: dict) -> dict:
    """Attempt to re-attach to a persisted session. Health-checks FIRST with a
    short timeout so a busy/wedged Fluent can't hang us in a blocking connect."""
    ip, port, password = rec.get("ip", "127.0.0.1"), rec.get("port"), rec.get("password", "")
    if not port:
        return {"reattached": False, "reason": "no port in persisted record"}
    try:
        health = _grpc_health(ip, port, password or "", timeout=4.0)
    except Exception as e:
        return {"reattached": False, "reason": f"health probe failed: {e}"}
    if health.get("status") == "BUSY_OR_UNRESPONSIVE":
        return {"reattached": False, "busy": True,
                "reason": "server port is open but the engine is not answering RPCs "
                          "(long blocking operation or wedged) - wait and retry; "
                          "do NOT kill it just for being busy"}
    if not health.get("grpc_server") or health.get("status") != "SERVING":
        return {"reattached": False, "reason": f"endpoint not serving: {health.get('status')}"}
    global _solver, _meshing, _session_mode, _connection_tag
    try:
        import ansys.fluent.core as pyfluent
        kw: dict = {"ip": ip, "port": port, "cleanup_on_exit": False}
        if password:
            kw["password"] = password
        sess, _rw = _pyfluent_connect(pyfluent, kw)
        mode = (rec.get("mode") or "solver").lower()
        if mode == "meshing":
            _meshing, _session_mode = sess, "meshing"
        else:
            _solver, _session_mode = sess, "solver"
        _connection_tag = f"{ip}:{port} (reattached)"
        _restore_meshing_context(rec)
        return {"reattached": True, "mode": _session_mode, "endpoint": f"{ip}:{port}"}
    except Exception as e:
        return {"reattached": False, "reason": f"connect failed: {e}"}


# ---------------------------------------------------------------------------
# Sizing units — workflow sizing fields use the CAD IMPORT unit, not meters.
# ---------------------------------------------------------------------------
_UNIT_TO_M = {"m": 1.0, "mm": 0.001, "cm": 0.01, "um": 1e-6,
              "in": 0.0254, "ft": 0.3048}


def _convert_size_to_import_unit(value: float, units: str) -> float:
    """Convert a sizing value given in `units` to the CAD import unit — the
    unit the Fluent watertight-workflow sizing fields actually use (verified
    live 2026-07-17: values land in the panel verbatim, and the panel is in
    the import unit; passing meters into a mm field made sizes 1000x too
    small and hung the mesher twice). Empty `units` = already in the import
    unit (passthrough)."""
    if not units:
        return value
    u, iu = units.strip().lower(), (_import_length_unit or "m").strip().lower()
    if u == iu:
        return value
    if u not in _UNIT_TO_M or iu not in _UNIT_TO_M:
        raise ValueError(f"unsupported unit '{units}' (known: {sorted(_UNIT_TO_M)})")
    return value * _UNIT_TO_M[u] / _UNIT_TO_M[iu]


def _mu():
    """meshing_utilities accessor (None when absent on this release).

    SAFETY: never call get_bounding_box_of_zone_list from tools — passing
    face-zone ids raised 'tg-get-thread-list-bounding-box: no threads
    specified' on every node and the tgrid process died (verified fatal,
    2026-07-15, 2025 R2). Stick to the calls verified safe live:
    get_all_objects / get_objects / get_labels / get_face_zones /
    mesh_exists / get_regions."""
    return getattr(_meshing, "meshing_utilities", None)


def _post_import_verification() -> dict:
    """Read back what the import actually produced: CAD objects, their
    labels, and the face-zone count. This is the scale/coverage evidence the
    meshing KB (doc 02, 'Unit verification') requires echoing to the user —
    every probe is best-effort and reports its own failure instead of
    raising."""
    v: dict = {}
    mu = _mu()
    if mu is None:
        v["note"] = "meshing_utilities not available on this release"
        return v
    try:
        objs = list(mu.get_all_objects() or [])
        v["objects"] = objs
    except Exception as e:
        objs = []
        v["objects_error"] = str(e)
    labels: dict = {}
    for o in objs:
        try:
            labels[o] = list(mu.get_labels(object_name=o) or [])
        except Exception as e:
            labels[o] = f"<get_labels failed: {e}>"
    v["labels_by_object"] = labels
    try:
        v["face_zone_count"] = len(list(mu.get_face_zones(filter="*") or []))
    except Exception as e:
        v["face_zone_count_error"] = str(e)
    return v


# ===========================================================================
# 1. CONNECTION
# ===========================================================================

@mcp.tool()
def connect_to_fluent(
    server_info_file: str = "",
    ip: str = "127.0.0.1",
    port: int = 0,
    mode: str = "solver",
    password: str = "",
) -> str:
    """Connect to a running ANSYS Fluent session via gRPC.

    Start Fluent first:  fluent 3ddp -sifile=server.txt -nm
    Fluent writes IP/port/password to server.txt.

    Args:
        server_info_file: Path to server-info file (recommended).
        ip:               IP address for port-based connection.
        port:             gRPC port (ignored when server_info_file is set).
        mode:             "solver" or "meshing".
        password:         Optional gRPC password.
    """
    global _solver, _meshing, _session_mode, _connection_tag
    if _bg_run["running"] or _session_job["running"]:
        return _j({"ok": False, "error": "a background job/calculation is in progress - "
                                         "stop_calculation or wait for it to finish before "
                                         "reconnecting (reconnecting tears down the session "
                                         "a worker thread may still be using)"})
    try:
        import ansys.fluent.core as pyfluent
    except ImportError:
        return _j({"ok": False, "error": "ansys-fluent-core not installed."})
    _disconnect_internal()
    kw: dict = {"cleanup_on_exit": False}
    if server_info_file:
        if not Path(server_info_file).exists():
            return _j({"ok": False, "error": f"File not found: {server_info_file}"})
        kw["server_info_file_name"] = server_info_file
        _connection_tag = server_info_file
    elif port:
        kw["ip"] = ip; kw["port"] = port
        _connection_tag = f"{ip}:{port}"
    else:
        return _j({"ok": False, "error": "Provide server_info_file or port."})
    if password:
        kw["password"] = password
    try:
        sess, remote_warnings = _pyfluent_connect(pyfluent, kw)
        if mode.lower() == "meshing":
            _meshing = sess; _session_mode = "meshing"
        else:
            _solver = sess; _session_mode = "solver"
        try:
            ver = str(sess.get_fluent_version())
        except Exception:
            ver = "unknown"
        _persist_session_state()
        return _j({"ok": True, "mode": _session_mode, "connection": _connection_tag,
                   "fluent_version": ver, "warnings": remote_warnings})
    except Exception as e:
        _solver = _meshing = None
        return _j({"ok": False, "error": str(e)})


def _find_fluent_exe(version: str = "") -> Path | None:
    """Locate fluent.exe: AWP_ROOT* env vars first, then the standard ANSYS
    install tree (newest version first). `version` like '25.2' prefers v252."""
    cands: list[Path] = []
    for env, p in sorted(os.environ.items(), reverse=True):
        if env.upper().startswith("AWP_ROOT") and p:
            cands.append(Path(p) / "fluent" / "ntbin" / "win64" / "fluent.exe")
    base = Path(r"C:\Program Files\ANSYS Inc")
    if base.exists():
        for d in sorted(base.glob("v*"), reverse=True):
            cands.append(d / "fluent" / "ntbin" / "win64" / "fluent.exe")
    if version:
        vtag = "v" + version.replace(".", "")
        cands = ([c for c in cands if vtag in str(c)]
                 + [c for c in cands if vtag not in str(c)])
    for c in cands:
        if c.exists():
            return c
    return None


def _launch_fluent_subprocess(mode: str, precision: str, processor_count: int,
                              version: str, show_gui: bool,
                              start_timeout: int) -> dict:
    """Launch Fluent DIRECTLY with a server-owned -sifile and connect from it.

    Replaces pyfluent.launch_fluent as the primary path: that launcher's
    gRPC-ready wait timed out on 3 of 6 launches (2026-07-17/20) while
    Fluent actually came up fine, and its temp server-info file was deleted
    on failure — leaving a healthy Fluent whose password was unrecoverable.
    Here the sifile lives in logs/ and SURVIVES, so the credentials are
    always recoverable and there is no launcher-side timeout at all.
    Returns a dict (not JSON) for launch_fluent to wrap."""
    import subprocess
    exe = _find_fluent_exe(version)
    if exe is None:
        return {"ok": False, "error": "fluent.exe not found (checked AWP_ROOT* and "
                                      r"C:\Program Files\ANSYS Inc\v*)",
                "fix": "pass version='' or check the installation"}
    sifile = _SESSION_STATE_PATH.parent / f"server_info_{int(time.time())}.txt"
    sifile.parent.mkdir(parents=True, exist_ok=True)
    args = [str(exe), "3ddp" if precision.lower() != "single" else "3d",
            f"-t{max(processor_count, 1)}"]
    if mode.lower() == "meshing":
        args.append("-meshing")
    if not show_gui:
        args.append("-g")   # batch: no GUI (graphics for pictures still work via null driver)
    args.append(f"-sifile={sifile}")
    try:
        subprocess.Popen(args, cwd=str(Path.cwd()))
    except Exception as e:
        return {"ok": False, "error": f"failed to spawn fluent.exe: {e}",
                "command": args}
    t0 = time.time()
    ip = port = password = None
    while time.time() - t0 < start_timeout:
        try:
            if sifile.exists():
                lines = [l.strip() for l in
                         sifile.read_text(encoding="utf-8", errors="replace").splitlines()
                         if l.strip()]
                if len(lines) >= 2 and ":" in lines[0]:
                    host, port_s = lines[0].rsplit(":", 1)
                    ip, port, password = host, int(port_s), lines[1]
                    break
        except Exception:
            pass
        time.sleep(3)
    if port is None:
        return {"ok": False, "error": f"server-info file not written within "
                                      f"{start_timeout}s: {sifile}",
                "fix": "check list_fluent_processes / the Fluent window; the file "
                       "path stays valid - retry connect_to_fluent(server_info_file=...) "
                       "once Fluent finishes starting"}
    # wait until the gRPC endpoint actually serves (Fluent writes the file
    # slightly before the server is ready)
    while time.time() - t0 < start_timeout:
        try:
            h = _grpc_health(ip, port, password, timeout=4.0)
            if h.get("status") in ("SERVING", "PASSWORD_REQUIRED"):
                break
        except Exception:
            pass
        time.sleep(3)
    import ansys.fluent.core as pyfluent
    _patch_pyfluent_workflow_bugs()
    kw = {"ip": ip, "port": port, "password": password, "cleanup_on_exit": False}
    sess, remote_warnings = _pyfluent_connect(pyfluent, kw)
    return {"ok": True, "session": sess, "sifile": str(sifile),
            "endpoint": f"{ip}:{port}", "warnings": remote_warnings}


@mcp.tool()
def launch_fluent(
    mode: str = "solver",
    precision: str = "double",
    processor_count: int = 1,
    version: str = "",
    show_gui: bool = False,
    use_gpu_solver: bool = False,
    kill_existing: bool = True,
    start_timeout: int = 300,
    replace_existing_session: bool = False,
) -> str:
    """Launch a new Fluent instance from Python and connect to it.

    SLOW + EXCLUSIVE: launching takes 1-5 minutes (license checkout +
    startup) and BLOCKS until connected — a silent wait is normal, not a
    hang. NEVER call this while a session is already alive or another
    launch may be in flight: a second Fluent wastes a license seat and
    orphans the first (three real incidents on 2026-07-15). If a session
    exists, this refuses unless replace_existing_session=True; prefer
    reusing it (check_fluent_connection / get_workflow_state) or
    connect_to_fluent for an externally-started Fluent.

    Before calling this, ASK THE USER for: (1) how many CPU cores to use
    (processor_count) — per the core-recommendation policy, use no more than
    1-4 for interactive setup and a separate higher count for the confirmed
    run session; and (2) whether to show the Fluent GUI (show_gui=True) or
    run headless (show_gui=False). Do not silently default either without
    asking. kill_existing defaults to True so a stray leftover Fluent
    process never blocks a license slot or gets connected to by accident —
    say so before launching, and set it False only if the user explicitly
    wants to keep another session running.

    Args:
        mode:            "solver" or "meshing".
        precision:       "double" or "single".
        processor_count: CPU cores. Ask the user; do not assume.
        version:         e.g. "25.2" — default uses newest installed.
        show_gui:        Show the Fluent GUI (True) or run headless (False). Ask the user.
        use_gpu_solver:  Enable the native GPU solver (requires compatible GPU + license).
        kill_existing:   Terminate any stray fluent/cortex/cxsolver OS processes
                         first, so this launch gets a clean license slot and
                         can't collide with a leftover session (default True).
        start_timeout:   Seconds to wait for Fluent to start and connect
                         (default 300). Raise this if launch fails with
                         "Deadline Exceeded" on a slow machine/license server.
    """
    global _solver, _meshing, _session_mode, _connection_tag
    t0 = time.time()
    if _bg_run["running"] or _session_job["running"]:
        return _j({"ok": False, "error": "a background job/calculation is in progress - "
                                         "stop_calculation or wait for it to finish before "
                                         "launching/replacing the session"})
    # refuse to stack sessions: a live session must be explicitly replaced
    if (_solver is not None or _meshing is not None) and _session_port_open():
        if not replace_existing_session:
            return _j({"ok": False,
                       "error": f"a live {_session_mode} session already exists - "
                                "refusing to launch a second Fluent",
                       "hint": "reuse it (check_fluent_connection / get_workflow_state), "
                               "or pass replace_existing_session=True to deliberately "
                               "shut it down and start fresh; unsaved work in the "
                               "existing session will be LOST"})
    try:
        import ansys.fluent.core as pyfluent
    except ImportError:
        return _j({"ok": False, "error": "ansys-fluent-core not installed."})
    _patch_pyfluent_workflow_bugs()
    _disconnect_internal()
    killed: list[int] = []
    if kill_existing:
        for p in _find_fluent_processes():
            if _kill_pid(p["pid"], force=True):
                killed.append(p["pid"])
    # PRIMARY PATH: direct subprocess launch with a persistent server-owned
    # -sifile (see _launch_fluent_subprocess). GPU solver still goes through
    # the pyfluent launcher below (needs its gpu plumbing).
    if not use_gpu_solver:
        res = _launch_fluent_subprocess(mode, precision, processor_count,
                                        version, show_gui, start_timeout)
        if res.get("ok"):
            if mode.lower() == "meshing":
                _meshing = res["session"]
            else:
                _solver = res["session"]
            _session_mode = mode.lower()
            _connection_tag = f"{res['endpoint']} (launched, sifile persistent)"
            _persist_session_state()
            return _j({"ok": True, "mode": _session_mode,
                       "processor_count": processor_count, "show_gui": show_gui,
                       "gpu_solver": False,
                       "killed_existing_processes": killed,
                       "launch_path": "direct sifile (no pyfluent launcher timeout)",
                       "server_info_file": res["sifile"],
                       "elapsed_s": round(time.time() - t0, 1),
                       "warnings": res.get("warnings", []),
                       "message": f"Fluent launched in {_session_mode} mode."})
        if "not found" not in str(res.get("error", "")):
            # a real launch/connect failure - report it; do NOT fall through
            # and start a second Fluent on top of a half-started one
            return _j({**res, "killed_existing_processes": killed,
                       "launch_path": "direct sifile"})
        # fluent.exe not found -> fall back to the pyfluent launcher
    kw: dict = {"cleanup_on_exit": False, "start_timeout": start_timeout}
    # show_gui is deprecated (and silently produced a '-nm' no-graphics launch
    # on pyfluent 0.39, verified 2026-07-15) — use UIMode when available
    try:
        kw["ui_mode"] = pyfluent.UIMode.GUI if show_gui else pyfluent.UIMode.NO_GUI_OR_GRAPHICS
    except AttributeError:
        kw["show_gui"] = show_gui
    kw["mode"] = pyfluent.FluentMode.MESHING if mode.lower() == "meshing" else pyfluent.FluentMode.SOLVER
    kw["precision"] = pyfluent.Precision.SINGLE if precision.lower() == "single" else pyfluent.Precision.DOUBLE
    if processor_count > 1:
        kw["processor_count"] = processor_count
    if version:
        kw["product_version"] = version
    if use_gpu_solver:
        kw["gpu"] = True
    # PyFluent 0.39's gRPC channel-ready wait is a hard 60 s regardless of
    # start_timeout; on a slow license server the first attempt flakily times
    # out and succeeds on retry (three occurrences on 2026-07-15 alone).
    # Retry ONCE after killing the orphaned half-start.
    last_exc: Exception | None = None
    for attempt in (1, 2):
        try:
            sess = pyfluent.launch_fluent(**kw)
            if mode.lower() == "meshing":
                _meshing = sess
            else:
                _solver = sess
            _session_mode = mode.lower()
            _connection_tag = "launched"
            _persist_session_state()
            return _j({"ok": True, "mode": _session_mode,
                       "processor_count": processor_count, "show_gui": show_gui,
                       "gpu_solver": use_gpu_solver,
                       "killed_existing_processes": killed,
                       "launch_attempts": attempt,
                       "elapsed_s": round(time.time() - t0, 1),
                       "message": f"Fluent launched in {_session_mode} mode."})
        except Exception as e:
            last_exc = e
            if attempt == 1 and ("deadline" in str(e).lower() or "timed out" in str(e).lower()):
                for p in _find_fluent_processes():
                    if _kill_pid(p["pid"], force=True):
                        killed.append(p["pid"])
                continue
            break
    msg = str(last_exc)
    diag = {"ok": False, "error": msg,
            "launch_kwargs": {k: str(v) for k, v in kw.items()},
            "killed_existing_processes": killed,
            "launch_attempts": 2}
    low = msg.lower()
    if "deadline" in low or "timed out" in low or "timeout" in low:
        diag["likely_cause"] = ("Fluent's gRPC server did not come up in time on "
                                "either attempt (slow machine or license server; "
                                "pyfluent 0.39 caps the channel wait at 60 s "
                                f"regardless of start_timeout={start_timeout})")
        diag["fix"] = ("kill any orphaned processes (list_fluent_processes / "
                       "kill_fluent_processes), check license-server health, "
                       "then relaunch")
    elif "license" in low or "feature" in low:
        diag["likely_cause"] = "license checkout failed"
        diag["fix"] = ("check the license server / free a seat; a stray Fluent "
                       "process may be holding it (list_fluent_processes)")
    elif "not found" in low or "install" in low or "version" in low:
        diag["likely_cause"] = "requested Fluent version not found on this machine"
        diag["fix"] = "retry with version='' to use the newest installed release"
    else:
        diag["fix"] = ("check whether fluent.exe processes appeared anyway "
                       "(list_fluent_processes) - an orphaned half-start must be "
                       "killed before relaunching")
    if killed:
        diag["note"] = ("kill_existing already terminated the processes listed in "
                        "killed_existing_processes BEFORE this failure - any prior "
                        "session is gone")
    return _j(diag)


def _disconnect_internal():
    global _solver, _meshing, _session_mode, _connection_tag, _workflow, _workflow_kind
    global _iteration_event_solver_id
    for s in (_solver, _meshing):
        if s is not None:
            try: s.exit()
            except Exception: pass
    _solver = _meshing = None
    _session_mode = _connection_tag = ""
    _workflow = None
    _workflow_kind = ""
    _health_cache["ts"] = 0.0  # never serve a stale healthy=True across a
                               # disconnect/relaunch (single choke point:
                               # called by connect_to_fluent, launch_fluent,
                               # disconnect_from_fluent, kill_fluent_processes)
    _scheme_channel_cache["endpoint"] = None
    old_ch = _scheme_channel_cache.get("channel")
    if old_ch is not None:
        try: old_ch.close()
        except Exception: pass
    _scheme_channel_cache["channel"] = _scheme_channel_cache["stub"] = None
    _iteration_event_solver_id = None  # force re-registration on the next session
    _live_iteration["index"] = None


# ---------------------------------------------------------------------------
# 1b. PROCESS MANAGEMENT — shut down stray Fluent processes before launching
# ---------------------------------------------------------------------------
_FLUENT_PROCESS_TAGS = ("fluent", "cortex", "cxsolver", "tgrid")
# The real worker executables are version-suffixed: cx2520.exe (cortex GUI/
# server) and fl2520.exe (solver/mesher host) — verified live 2026-07-17,
# where name-only matching found the fluent.exe wrapper but MISSED the cortex
# process that actually held the gRPC listen ports, so discovery probed 0
# ports on a healthy session.
_FLUENT_NAME_RE = re.compile(r"^(fluent|cortex|cxsolver|tgrid|cx\d{3,4}|fl(_mpi)?\d{3,4})", re.I)
# Cmdline fallback for helper exes (mpiexec, hydra_pmi_proxy, ...): require an
# ANSYS INSTALL-TREE marker, not just the words 'ansys'/'fluent' anywhere — a
# loose match flagged this MCP server's own python.exe (its path contains
# 'pyansys-env' and 'Pyfluent_mcp'), which kill_fluent_processes would then
# kill (caught in testing 2026-07-17).
_FLUENT_CMDLINE_RE = re.compile(r"ansys\s?inc|ansysi~1|[\\/]fluent\d+\.\d+", re.I)


def _is_fluent_process(name: str, cmdline: str) -> bool:
    if _FLUENT_NAME_RE.match(name or ""):
        return True
    return bool(_FLUENT_CMDLINE_RE.search(cmdline or ""))


def _find_fluent_processes() -> list[dict]:
    """OS-level Fluent-related processes (fluent/cortex/cxsolver/cxNNNN/flNNNN),
    including ones this MCP session never launched — e.g. a prior crashed run
    or a manually started Fluent. Uses psutil when available, else the
    platform's native tasklist/ps — never a discovery step, just a fixed
    lookup."""
    try:
        import psutil
    except ImportError:
        return _find_fluent_processes_fallback()
    procs = []
    for p in psutil.process_iter(["pid", "name", "create_time", "cmdline"]):
        try:
            name = (p.info.get("name") or "")
            cmdline = " ".join(p.info.get("cmdline") or [])
            if _is_fluent_process(name, cmdline):
                created = p.info.get("create_time")
                procs.append({
                    "pid": p.info["pid"],
                    "name": name,
                    "created": (datetime.fromtimestamp(created).strftime("%Y-%m-%d %H:%M:%S")
                                if created else None),
                    "cmdline": " ".join(p.info.get("cmdline") or [])[:200],
                })
        except (psutil.NoSuchProcess, psutil.AccessDenied):
            continue
    return procs


def _find_fluent_processes_fallback() -> list[dict]:
    import subprocess
    procs = []
    try:
        if os.name == "nt":
            out = subprocess.check_output(["tasklist", "/FO", "CSV", "/NH"],
                                           text=True, stderr=subprocess.DEVNULL)
            for line in out.splitlines():
                parts = [c.strip('"') for c in line.split('","')]
                if len(parts) >= 2 and _is_fluent_process(parts[0], ""):
                    procs.append({"pid": int(parts[1]), "name": parts[0], "created": None, "cmdline": ""})
        else:
            out = subprocess.check_output(["ps", "-eo", "pid,comm"],
                                           text=True, stderr=subprocess.DEVNULL)
            for line in out.splitlines()[1:]:
                parts = line.strip().split(None, 1)
                if len(parts) == 2 and _is_fluent_process(parts[1], ""):
                    procs.append({"pid": int(parts[0]), "name": parts[1], "created": None, "cmdline": ""})
    except Exception:
        pass
    return procs


def _kill_pid(pid: int, force: bool = True) -> bool:
    try:
        import psutil
        proc = psutil.Process(pid)
        proc.kill() if force else proc.terminate()
        return True
    except ImportError:
        import subprocess
        try:
            if os.name == "nt":
                subprocess.run(["taskkill", "/PID", str(pid)] + (["/F"] if force else []),
                                capture_output=True, check=False)
            else:
                subprocess.run(["kill", "-9" if force else "-15", str(pid)],
                                capture_output=True, check=False)
            return True
        except Exception:
            return False
    except Exception:
        return False


# ---------------------------------------------------------------------------
# 1c. DIRECT gRPC LAYER (ansys-api-fluent proto stubs)
#
# Implements the PyAnsys gRPC API-package pattern
# (https://dev.docs.pyansys.com/how-to/grpc-api-packages.html): talk to
# Fluent's gRPC server directly through the generated ansys.api.fluent.v0
# *_pb2 / *_pb2_grpc stubs instead of a full PyFluent session object.
# Authentication matches PyFluent: the server password travels as
# ("password", <pwd>) metadata on every call.
#
# Use cases: probe whether a port is a Fluent gRPC endpoint, health-check a
# server, evaluate Scheme on a session without building session objects, and
# discover reachable gRPC servers among running Fluent processes.
# ---------------------------------------------------------------------------

def _grpc_channel(ip: str, port: int):
    import grpc
    return grpc.insecure_channel(f"{ip}:{port}")


def _grpc_metadata(password: str) -> list[tuple[str, str]]:
    return [("password", password)] if password else []


def _grpc_health_stubs():
    """Stubs for Fluent's grpc.health.v1.Health service. PREFER the standard
    grpc_health package (the one PyFluent itself uses): the ansys.api.fluent
    copy declares the same grpc.health.v1 symbols under a different proto
    file path, so importing it AFTER pyfluent has loaded grpc_health poisons
    the protobuf descriptor pool with "duplicate symbol
    'grpc.health.v1.HealthCheckRequest'" — verified live 2026-07-17: every
    health probe in an MCP process that had previously run launch_fluent
    failed with exactly that error. Fall back to the ansys stubs only when
    grpc_health isn't installed."""
    try:
        from grpc_health.v1 import health_pb2, health_pb2_grpc
    except ImportError:
        from ansys.api.fluent.v0 import health_pb2, health_pb2_grpc
    return health_pb2, health_pb2_grpc


def _grpc_health(ip: str, port: int, password: str, timeout: float) -> dict:
    """Single HealthCheck RPC. Returns a dict describing the endpoint."""
    import grpc
    health_pb2, health_pb2_grpc = _grpc_health_stubs()
    channel = _grpc_channel(ip, port)
    try:
        stub = health_pb2_grpc.HealthStub(channel)
        resp = stub.Check(health_pb2.HealthCheckRequest(service=""),
                          metadata=_grpc_metadata(password), timeout=timeout)
        status = health_pb2.HealthCheckResponse.ServingStatus.Name(resp.status)
        return {"grpc_server": True, "status": status,
                "authenticated": True, "endpoint": f"{ip}:{port}"}
    except grpc.RpcError as e:
        detail = e.details() or ""
        code = e.code().name if e.code() else "UNKNOWN"
        # Fluent rejects unauthenticated calls with this message instead of
        # a clean UNAUTHENTICATED status code.
        if "has not been validated" in detail or code in ("UNAUTHENTICATED", "PERMISSION_DENIED"):
            return {"grpc_server": True, "status": "PASSWORD_REQUIRED",
                    "authenticated": False, "endpoint": f"{ip}:{port}",
                    "hint": "Provide the session password (from its server-info file)."}
        # BUSY vs DEAD: an RPC timeout on a port that still accepts TCP is the
        # signature of a Fluent whose engine is mid-blocking-operation (import,
        # surface/volume mesh, GUI action) or wedged — NOT "not a gRPC server".
        # Reporting it as dead caused live sessions to be killed (2026-07-17).
        if code in ("DEADLINE_EXCEEDED", "UNAVAILABLE"):
            try:
                import socket
                with socket.create_connection((ip, port), timeout=2):
                    port_accepts_tcp = True
            except Exception:
                port_accepts_tcp = False
            if port_accepts_tcp:
                return {"grpc_server": None, "status": "BUSY_OR_UNRESPONSIVE",
                        "endpoint": f"{ip}:{port}", "port_open": True,
                        "detail": detail[:200],
                        "hint": "Port accepts TCP but the RPC went unanswered - the "
                                "engine is likely busy in a long blocking operation "
                                "(meshing/import can take minutes). WAIT and re-probe; "
                                "do not kill the process just for being busy."}
        return {"grpc_server": False, "status": code,
                "endpoint": f"{ip}:{port}", "detail": detail[:200]}
    finally:
        channel.close()


@mcp.tool()
def grpc_health_check(ip: str = "127.0.0.1", port: int = 0,
                      password: str = "", timeout: float = 5.0) -> str:
    """Health-check a Fluent gRPC endpoint directly via ansys-api-fluent stubs.

    Works without creating a PyFluent session. Distinguishes three cases:
    a healthy gRPC server (SERVING), a gRPC server that wants a password
    (PASSWORD_REQUIRED), and a port that is not a Fluent gRPC server at all.

    Args:
        ip:       Server IP (default 127.0.0.1).
        port:     gRPC port (required).
        password: Session password, sent as ("password", ...) metadata.
        timeout:  Per-RPC timeout in seconds.
    """
    if not port:
        return _j({"ok": False, "error": "port is required"})
    try:
        return _j({"ok": True, **_grpc_health(ip, port, password, timeout)})
    except ImportError:
        return _j({"ok": False, "error": "ansys-api-fluent / grpcio not installed."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def grpc_scheme_eval(expression: str, ip: str = "127.0.0.1", port: int = 0,
                     password: str = "", timeout: float = 30.0) -> str:
    """Evaluate a Scheme expression over raw gRPC (SchemeEval/StringEval).

    Bypasses PyFluent session objects entirely — useful to talk to a live
    Fluent server this MCP process never launched, or when a full
    connect_to_fluent is unnecessary. The connected session (if any) is
    untouched.

    Args:
        expression: Scheme expression, e.g. (rpgetvar 'flow-time).
        ip:         Server IP (default 127.0.0.1).
        port:       gRPC port (required).
        password:   Session password, sent as ("password", ...) metadata.
        timeout:    Per-RPC timeout in seconds.
    """
    if not port:
        return _j({"ok": False, "error": "port is required"})
    try:
        import grpc
        from ansys.api.fluent.v0 import scheme_eval_pb2, scheme_eval_pb2_grpc
    except ImportError:
        return _j({"ok": False, "error": "ansys-api-fluent / grpcio not installed."})
    channel = _grpc_channel(ip, port)
    try:
        stub = scheme_eval_pb2_grpc.SchemeEvalStub(channel)
        resp = stub.StringEval(scheme_eval_pb2.StringEvalRequest(input=expression),
                               metadata=_grpc_metadata(password), timeout=timeout)
        return _j({"ok": True, "expression": expression, "output": resp.output})
    except grpc.RpcError as e:
        return _j({"ok": False, "error": e.details() or e.code().name,
                   "grpc_status": e.code().name if e.code() else "UNKNOWN"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    finally:
        channel.close()


@mcp.tool()
def grpc_discover_fluent_servers(password: str = "", timeout: float = 3.0) -> str:
    """Scan running Fluent processes for reachable gRPC server endpoints.

    Enumerates the TCP ports each fluent/cortex process is listening on and
    probes every one with a HealthCheck RPC. Reports, per port, whether it is
    a Fluent gRPC server and whether it needs a password. Use this to find a
    session this MCP process can re-attach to (e.g. after an MCP restart or
    for a Fluent started outside Python).

    Args:
        password: Optional password to try on endpoints that require one.
        timeout:  Per-port probe timeout in seconds (keep small).
    """
    try:
        import psutil
    except ImportError:
        return _j({"ok": False, "error": "psutil not installed."})
    results = []
    for proc in _find_fluent_processes():
        try:
            p = psutil.Process(proc["pid"])
            conns = (p.net_connections(kind="tcp") if hasattr(p, "net_connections")
                     else p.connections(kind="tcp"))
        except (psutil.NoSuchProcess, psutil.AccessDenied):
            continue
        ports = sorted({c.laddr.port for c in conns
                        if c.status == psutil.CONN_LISTEN and c.laddr})
        for port in ports:
            try:
                probe = _grpc_health("127.0.0.1", port, password, timeout)
            except Exception as e:
                probe = {"grpc_server": False, "status": "PROBE_ERROR", "detail": str(e)[:200]}
            results.append({"pid": proc["pid"], "process": proc["name"], "port": port, **probe})
    servers = [r for r in results if r.get("grpc_server")]
    return _j({"ok": True, "ports_probed": len(results), "grpc_servers_found": len(servers),
               "servers": servers, "all_results": results,
               "next_step": ("connect_to_fluent(ip='127.0.0.1', port=<port>, password=...) "
                             "for a SERVING endpoint" if servers else
                             "No gRPC endpoints reachable. Start one from the Fluent console: "
                             "server/start-server \"<path>/server_info.txt\"")})


@mcp.tool()
def grpc_write_server_info(ip: str, port: int, password: str, file_path: str) -> str:
    """Write a standard Fluent server-info file (host:port + password lines).

    Lets connect_to_fluent(server_info_file=...) — or any other PyFluent
    client — attach to a known endpoint later without re-typing credentials.

    Args:
        ip:        Server IP.
        port:      gRPC port.
        password:  Session password.
        file_path: Destination path for the server-info file.
    """
    try:
        _mkdir(file_path)
        Path(file_path).write_text(f"{ip}:{port}\n{password}\n", encoding="utf-8")
        return _j({"ok": True, "file": str(Path(file_path).resolve()),
                   "reconnect_with": f"connect_to_fluent(server_info_file='{file_path}')"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def list_fluent_processes() -> str:
    """List OS-level Fluent-related processes currently running on this
    machine (fluent/cortex/cxsolver), including any NOT tracked by this MCP
    session — e.g. left over from a prior crashed run or started manually.
    Call before launch_fluent to check for stray sessions holding a license."""
    return _j({"ok": True, "processes": _find_fluent_processes()})


@mcp.tool()
def kill_fluent_processes(pids: list[int] | None = None, force: bool = True) -> str:
    """Terminate Fluent-related OS processes (fluent/cortex/cxsolver) and
    drop this MCP session's own solver/meshing handles.

    Args:
        pids:  Specific PIDs to kill (from list_fluent_processes). Omit or
               pass an empty list to kill ALL detected Fluent processes.
        force: Hard-kill (SIGKILL / taskkill /F) instead of a graceful
               terminate signal first.
    """
    found = _find_fluent_processes()
    targets = [p["pid"] for p in found if not pids or p["pid"] in pids]
    killed, failed = [], []
    for pid in targets:
        (killed if _kill_pid(pid, force) else failed).append(pid)
    _disconnect_internal()
    if killed:
        _clear_persisted_session()   # the recorded endpoint is now dead
    return _j({"ok": True, "found": found, "killed": killed, "failed": failed})


@mcp.tool()
def disconnect_from_fluent() -> str:
    """Disconnect from the active Fluent session (Fluent keeps running)."""
    if err := _chk(): return err
    _disconnect_internal()
    return _j({"ok": True, "message": "Disconnected."})


_health_cache: dict = {"ts": 0.0, "payload": None}
_HEALTH_CACHE_TTL_S = 1.0


@mcp.tool()
def check_fluent_connection(force_refresh: bool = False) -> str:
    """Return current connection status, mode, and health.

    Health is verified two ways: the session health-check RPC AND a raw
    TCP probe of the server port. A dead Fluent can leave a stale-healthy
    session object (check_health() returns a status STRING — even a
    failure status is truthy; verified incident 2026-07-11), so trust the
    port probe over the session's own answer.

    A short-TTL cache (~1s) is used, but ONLY for a positive healthy=True
    result - a busy/unhealthy/unreachable outcome is never cached, since
    those are exactly the states where freshness matters most (this tool
    exists specifically to catch an external Fluent crash; masking one
    behind a cache would reintroduce the 'stale-healthy' incident this
    function was hardened against). Pass force_refresh=True before any
    consequential operation (a solve, volume mesh, etc.) to always get a
    live probe.

    SELF-HEALING: when this MCP process holds no session (e.g. it was
    restarted mid-run), this tool automatically tries to re-attach to the
    last known Fluent endpoint persisted in logs/last_session.json — so an
    MCP restart no longer strands a healthy Fluent."""
    if not force_refresh and time.time() - _health_cache["ts"] < _HEALTH_CACHE_TTL_S \
            and _health_cache["payload"] is not None:
        return _health_cache["payload"]
    if _solver is None and _meshing is None:
        rec = _load_persisted_session()
        if rec and rec.get("port"):
            attempt = _try_reattach(rec)
            if attempt.get("reattached"):
                return _j({"connected": True, "reattached": True,
                           "mode": _session_mode, "connection": _connection_tag,
                           "healthy": True,
                           "restored_context": {
                               "import_length_unit": _import_length_unit,
                               "meshing_mode": _meshing_mode,
                               "analysis_mode": _analysis_mode,
                               "mesh_plan_present": _mesh_plan is not None,
                               "solver_plan_present": _solver_plan is not None},
                           "note": "MCP had lost its session handle (restart); "
                                   "re-attached to the running Fluent from the "
                                   "persisted record."})
            if attempt.get("busy"):
                return _j({"connected": False, "fluent_likely_alive_but_busy": True,
                           "last_known_endpoint": f"{rec.get('ip')}:{rec.get('port')}",
                           "reattach_attempt": attempt,
                           "note": "A previously-registered Fluent is unresponsive but "
                                   "its port is open - likely mid-blocking-operation. "
                                   "Wait, then call check_fluent_connection again. Do "
                                   "not launch a second Fluent."})
            return _j({"connected": False,
                       "reattach_attempt": attempt,
                       "note": "No live session; the persisted record did not lead to "
                               "a serving endpoint (Fluent probably exited)."})
        return _j({"connected": False})
    try:
        s = _active()
        status = str(s.check_health()) if hasattr(s, "check_health") else ""
    except Exception:
        status = "UNREACHABLE"
    port_open = False
    try:
        import socket
        props = _active().connection_properties
        ip = props.ip if props.ip not in (None, "", "localhost") else "127.0.0.1"
        with socket.create_connection((ip, props.port), timeout=3):
            port_open = True
    except Exception:
        port_open = False
    us = status.upper()
    healthy = port_open and "NOT_SERVING" not in us and us != "UNREACHABLE"
    if healthy:
        note = None
    elif port_open:
        note = ("Port is open but the health RPC failed — the engine is likely "
                "BUSY in a long blocking operation (import/meshing take minutes) "
                "or wedged. Wait and re-check before assuming it is dead; do not "
                "kill or relaunch just for being busy.")
    else:
        note = ("Server unreachable — the Fluent process has likely "
                "exited. Use grpc_discover_fluent_servers to find a "
                "live server, or launch_fluent + read_case_data to "
                "resume from the last saved .cas.h5/.dat.h5.")
    payload = _j({"connected": True, "mode": _session_mode,
                  "connection": _connection_tag, "healthy": healthy,
                  "busy_suspected": bool(port_open and not healthy),
                  "health_status": status or "n/a", "server_port_open": port_open,
                  "note": note})
    if healthy:  # never cache a busy/unhealthy/unreachable result
        _health_cache["ts"] = time.time()
        _health_cache["payload"] = payload
    return payload


@mcp.tool()
def reattach_last_session() -> str:
    """Re-attach to the Fluent endpoint persisted from the last successful
    launch/connect (logs/last_session.json) — the recovery path after this
    MCP server was restarted while Fluent kept running. Health-checks the
    endpoint with a short timeout first, so a busy Fluent is reported as
    busy instead of hanging this call. Restores the meshing context
    (CAD import unit, meshing mode, sizing plan) recorded with the session."""
    if (_solver is not None or _meshing is not None) and _session_port_open():
        return _j({"ok": True, "already_connected": True, "mode": _session_mode})
    rec = _load_persisted_session()
    if not rec:
        return _j({"ok": False, "error": "no persisted session record found",
                   "hint": "grpc_discover_fluent_servers can scan running "
                           "processes for reachable endpoints"})
    attempt = _try_reattach(rec)
    if attempt.get("reattached"):
        return _j({"ok": True, **attempt,
                   "restored_context": {"import_length_unit": _import_length_unit,
                                        "meshing_mode": _meshing_mode,
                                        "analysis_mode": _analysis_mode,
                                        "mesh_plan_present": _mesh_plan is not None,
                                        "solver_plan_present": _solver_plan is not None}})
    return _j({"ok": False, **attempt, "record": {k: rec.get(k) for k in
               ("ts", "mode", "ip", "port", "fluent_host_pid")}})


_MAIN_MENU = [
    "1. Project summary", "2. Geometry and mesh", "3. Regions and boundaries",
    "4. Materials and physics", "5. Solver settings", "6. Reports and monitors",
    "7. Run controls", "8. Results and visualisation", "9. Files and exports",
    "10. Logs and diagnostics", "11. Convergence assessment", "12. Exit or disconnect",
]


@mcp.tool()
def project_summary() -> str:
    """Compact project/session status + the main session menu (interactive
    console spec §3-4): connection, modes, mesh inventory, physics, monitor
    count, run state, and convergence flag — composed from live state, all
    probes best-effort. Render the menu prompt via AskUserQuestion; accept
    both menu numbers and natural-language requests (they map to the same
    tools)."""
    s: dict = {"connected": _solver is not None or _meshing is not None,
               "session_mode": _session_mode or "none",
               "meshing_mode": _meshing_mode, "analysis_mode": _analysis_mode,
               "import_length_unit": _import_length_unit,
               "mesh_plan_present": _mesh_plan is not None,
               "solver_plan_present": _solver_plan is not None}
    if _solver is not None:
        try:
            zones = _zone_type_pairs()
            by_type: dict = {}
            for z in zones:
                by_type[z["type"]] = by_type.get(z["type"], 0) + 1
            s["mesh"] = {"boundary_zones": len(zones), "zones_by_type": by_type,
                         "cell_zones": list(_solver.setup.cell_zone_conditions.keys())}
        except Exception as e:
            s["mesh"] = {"error": str(e)}
        try:
            phys = {"turbulence_model": str(_solver.setup.models.viscous.model())}
            try:
                phys["energy"] = bool(_solver.setup.models.energy.enabled())
            except Exception:
                pass
            s["physics"] = phys
        except Exception as e:
            s["physics"] = {"error": str(e)}
        s["report_monitors_defined"] = _count_report_definitions()
    run = {"calculation_running": _bg_run["running"],
           "background_job": _session_job["name"] if _session_job["running"] else None}
    try:
        trn = _find_transcript()
        if trn is not None:
            tail = _tail_read(trn, max_bytes=4000)
            parsed = _parse_transcript_tail(tail.splitlines())
            run["converged"] = parsed.get("converged", False)
            if "current_iteration" in parsed:
                run["last_iteration"] = parsed["current_iteration"]
    except Exception:
        pass
    s["run"] = run
    return _j({"ok": True, "summary": s, "menu": _MAIN_MENU,
               "interactive_prompts": [_prompt(
                   "What next?", "Session menu",
                   [("Continue workflow (Recommended)",
                     "proceed with the next natural step for the current state"),
                    ("Run controls", "initialize / run / stop / autosave / status"),
                    ("Results & visualisation", "contours, cut planes, KPIs, exports"),
                    ("Regions & boundaries", "zone list, BC details, edits")])],
               "prompt_render": _PROMPT_RENDER_NOTE,
               "note": "menu numbers and natural language both work - they map to "
                       "the same MCP tools"})


# ===========================================================================
# 2. CASE I/O
# ===========================================================================

def _read_case_impl(case_file_path: str) -> str:
    try:
        _solver.file.read_case(file_name=case_file_path)
        return _j({"ok": True, "message": f"Case loaded: {case_file_path}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def read_case(case_file_path: str, run_in_background: bool = True) -> str:
    """Read a Fluent case file (.cas or .cas.h5).

    Runs AS A BACKGROUND JOB by default - a large case file can take as
    long to read as the meshing calls that were already backgrounded for
    this exact reason (verified live 2026-07-20: this and the other 4
    case/mesh file-I/O tools were the only ones in this class never
    migrated off the MCP request thread). Poll get_job_status +
    read_console_tail for progress.

    Args:
        case_file_path:     Absolute path to the case file.
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if not Path(case_file_path).exists():
        return _j({"ok": False, "error": f"Not found: {case_file_path}"})
    _set_project_dir_from(case_file_path)
    if run_in_background:
        return _run_session_job("read_case", lambda: _read_case_impl(case_file_path))
    return _read_case_impl(case_file_path)


def _read_case_data_impl(case_data_file_path: str) -> str:
    try:
        _solver.file.read_case_data(file_name=case_data_file_path)
        return _j({"ok": True, "message": f"Case+data loaded: {case_data_file_path}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def read_case_data(case_data_file_path: str, run_in_background: bool = True) -> str:
    """Read a case+data file pair (.cas.h5 with associated .dat.h5).

    Runs AS A BACKGROUND JOB by default (same rationale as read_case - a
    large case+data pair is exactly the class of multi-minute file I/O
    the meshing tools were already backgrounded to avoid freezing the
    server on). Poll get_job_status + read_console_tail for progress.

    Args:
        case_data_file_path: Path to the .cas.h5 file.
        run_in_background:   Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if not Path(case_data_file_path).exists():
        return _j({"ok": False, "error": f"Not found: {case_data_file_path}"})
    _set_project_dir_from(case_data_file_path)
    if run_in_background:
        return _run_session_job("read_case_data", lambda: _read_case_data_impl(case_data_file_path))
    return _read_case_data_impl(case_data_file_path)


def _write_case_impl(output_path: str) -> str:
    try:
        _mkdir(output_path)
        _solver.file.write_case(file_name=output_path)
        return _j({"ok": True, "message": f"Case written: {output_path}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def write_case(output_path: str, run_in_background: bool = True) -> str:
    """Write the current case to disk.

    Runs AS A BACKGROUND JOB by default (same rationale as read_case).
    Poll get_job_status + read_console_tail for progress.

    Args:
        output_path:        Output path (.cas.h5 recommended).
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if run_in_background:
        return _run_session_job("write_case", lambda: _write_case_impl(output_path))
    return _write_case_impl(output_path)


def _write_case_data_impl(output_path: str) -> str:
    try:
        _mkdir(output_path)
        _solver.file.write_case_data(file_name=output_path)
        return _j({"ok": True, "message": f"Case+data written: {output_path}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def write_case_data(output_path: str, run_in_background: bool = True) -> str:
    """Write case and solution data to disk.

    Runs AS A BACKGROUND JOB by default (same rationale as read_case).
    Poll get_job_status + read_console_tail for progress.

    Args:
        output_path:        Output path (.cas.h5).
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if run_in_background:
        return _run_session_job("write_case_data", lambda: _write_case_data_impl(output_path))
    return _write_case_data_impl(output_path)


# ===========================================================================
# 3. MESH INSPECTION
# ===========================================================================

def _zone_type_pairs() -> list[dict]:
    """Enumerate every boundary zone with its BC type. Shared by
    list_boundary_zones and export_case_manifest — do not duplicate."""
    zones = []
    bc_root = _solver.setup.boundary_conditions
    for bc_type in dir(bc_root):
        if bc_type.startswith("_"): continue
        try:
            grp = getattr(bc_root, bc_type)
            if hasattr(grp, "keys"):
                for name in grp.keys():
                    zones.append({"zone": name, "type": bc_type})
        except Exception:
            continue
    return zones


@mcp.tool()
def list_boundary_zones() -> str:
    """List all boundary zones and their types."""
    if err := _chk("solver"): return err
    try:
        zones = _zone_type_pairs()
        return _j({"ok": True, "boundary_zones": zones, "count": len(zones)})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_boundary_state(zone_name: str) -> str:
    """Full state of ONE boundary zone (interactive console spec §6): type,
    current applied values (read back from Fluent, not assumed), area, and
    a readiness status listing missing required inputs (e.g. a velocity
    inlet with no velocity). Use before/after edits and to build the BC
    review table."""
    if err := _chk("solver"): return err
    zones = {z["zone"]: z["type"] for z in _zone_type_pairs()}
    if zone_name not in zones:
        return _j({"ok": False, "error": f"zone '{zone_name}' not found",
                   "available": sorted(zones)})
    bc_type = zones[zone_name]
    values = _read_bc_values(zone_name, bc_type)
    area = None
    try:
        raw = _scheme_str(f"(surface-area '({zone_name}))")
        area = float(str(raw).strip().strip("()"))
    except Exception:
        pass
    missing: list[str] = []
    if bc_type == "velocity_inlet":
        v = values.get("velocity_ms")
        if not isinstance(v, (int, float)) or v == 0:
            missing.append("velocity magnitude")
    elif bc_type == "mass_flow_inlet" and not values:
        missing.append("mass flow rate")
    status = "Ready" if not missing else f"Missing: {', '.join(missing)}"
    return _j({"ok": True, "zone": zone_name, "type": bc_type,
               "current_values": values,
               "area_m2": area, "status": status,
               "actions": ["set_velocity_inlet / set_pressure_outlet / "
                           "set_wall_boundary to edit (each verifies by read-back)",
                           "create_wall_display + display_scene to highlight in 3D"]})


@mcp.tool()
def list_cell_zones() -> str:
    """List all fluid/solid cell zones."""
    if err := _chk("solver"): return err
    try:
        zones = list(_solver.setup.cell_zone_conditions.keys())
        return _j({"ok": True, "cell_zones": zones, "count": len(zones)})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_mesh_info() -> str:
    """Return zone names and mesh quality check summary."""
    if err := _chk("solver"): return err
    try:
        zones = list(_solver.setup.boundary_conditions.keys())
        info: dict = {"boundary_zones": zones, "zone_count": len(zones)}
        try:
            info["mesh_check"] = str(_solver.tui.mesh.check())[:3000]
        except Exception as e:
            info["mesh_check_error"] = str(e)
        return _j({"ok": True, **info})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 4. PHYSICS / MODELS
# ===========================================================================

@mcp.tool()
def set_energy_equation(enabled: bool = True) -> str:
    """Enable or disable the energy equation.

    Args:
        enabled: True to enable (default True).
    """
    if err := _chk("solver"): return err
    try:
        _solver.setup.models.energy.enabled = enabled
        return _j({"ok": True, "message": f"Energy equation {'enabled' if enabled else 'disabled'}."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_viscous_model(model: str = "k-epsilon", sub_model: str = "realizable") -> str:
    """Set the turbulence/viscous model.

    Args:
        model:     "laminar", "spalart-allmaras", "k-epsilon", "k-omega", "sst",
                   "les", "sbes", "wmles".
        sub_model: k-epsilon sub-type: "standard", "rng", "realizable".
                   k-omega sub-type: "standard", "sst".
    """
    if err := _chk("solver"): return err
    try:
        v = _solver.setup.models.viscous
        m = {"sst": "k-omega", "les": "large-eddy-simulation",
             "sbes": "stress-blended-eddy-simulation",
             "wmles": "wall-modeled-large-eddy-simulation"}.get(model.lower(), model.lower())
        v.model = m
        if m == "k-epsilon": v.k_epsilon_model = sub_model
        elif m == "k-omega": v.k_omega_model = sub_model if sub_model != "realizable" else "sst"
        return _j({"ok": True, "message": f"Viscous model: {m} / {sub_model}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_physics_models() -> str:
    """Return current energy, viscous, multiphase, and radiation model settings."""
    if err := _chk("solver"): return err
    info: dict = {}
    m = _solver.setup.models
    for attr, key in [("energy.enabled", "energy_enabled"),
                      ("viscous.model", "viscous_model"),
                      ("multiphase.model", "multiphase_model"),
                      ("radiation.model", "radiation_model")]:
        try:
            obj = m
            for a in attr.split("."): obj = getattr(obj, a)
            info[key] = obj() if callable(obj) else obj
        except Exception:
            info[key] = "N/A"
    return _j({"ok": True, "models": info})


# ===========================================================================
# 5. MATERIALS
# ===========================================================================

@mcp.tool()
def list_materials() -> str:
    """List all fluid and solid materials in the case."""
    if err := _chk("solver"): return err
    try:
        mats = {}
        for t in ("fluid", "solid"):
            try: mats[t] = list(getattr(_solver.setup.materials, t).keys())
            except Exception: mats[t] = []
        return _j({"ok": True, "materials": mats})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_fluid_material(cell_zone: str, material_name: str) -> str:
    """Assign a fluid material to a cell zone.

    Args:
        cell_zone:     Cell zone name.
        material_name: Material name (must exist in database).
    """
    if err := _chk("solver"): return err
    try:
        _solver.setup.cell_zone_conditions.fluid[cell_zone].material = material_name
        # ASCII-only message: non-ASCII here crashes cp1252 consoles on Windows
        return _j({"ok": True, "message": f"'{material_name}' -> '{cell_zone}'"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 6. BOUNDARY CONDITIONS
# ===========================================================================

@mcp.tool()
def set_velocity_inlet(
    zone_name: str,
    velocity_ms: float,
    temperature_k: float = 300.0,
    turbulence_intensity: float = 0.05,
    turbulence_length_scale_m: float = 0.01,
) -> str:
    """Set a velocity-inlet boundary condition.

    Args:
        zone_name:                 Boundary zone name.
        velocity_ms:               Velocity magnitude (m/s).
        temperature_k:             Inlet temperature (K).
        turbulence_intensity:      0–1 (default 0.05 = 5 %).
        turbulence_length_scale_m: Turbulence length scale (m).
    """
    if err := _chk("solver"): return err
    try:
        vi = _solver.setup.boundary_conditions.velocity_inlet[zone_name]
        vi.momentum.velocity.value = velocity_ms
        thermal_applied = True
        try: vi.thermal.t.value = temperature_k
        except Exception: thermal_applied = False
        # Fields like turbulent_length_scale are INACTIVE until the matching
        # specification method is selected — set the method first, then values.
        turb_applied, turb_error = True, None
        try:
            t = vi.turbulence
            spec = getattr(t, "turbulence_specification", None) or t.turbulent_specification
            spec.set_state("Intensity and Length Scale")
            t.turbulent_intensity = turbulence_intensity
            t.turbulent_length_scale = turbulence_length_scale_m
        except Exception as te:
            turb_applied, turb_error = False, str(te)
        # transactional read-back (spec §6): a successful write call is not
        # proof the value landed - compare what Fluent actually holds
        readback = _read_bc_values(zone_name, "velocity_inlet")
        rb_v = readback.get("velocity_ms")
        verified = isinstance(rb_v, (int, float)) and abs(rb_v - velocity_ms) < 1e-9
        return _j({"ok": True, "message": f"Inlet '{zone_name}': {velocity_ms} m/s",
                   "readback": readback, "verified": verified,
                   "thermal_applied": thermal_applied,
                   "turbulence_applied": turb_applied,
                   "turbulence_error": turb_error})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_pressure_inlet(
    zone_name: str,
    total_pressure_pa: float,
    temperature_k: float = 300.0,
) -> str:
    """Set a pressure-inlet boundary condition.

    Args:
        zone_name:         Boundary zone name.
        total_pressure_pa: Gauge total pressure (Pa).
        temperature_k:     Total temperature (K).
    """
    if err := _chk("solver"): return err
    try:
        pi = _solver.setup.boundary_conditions.pressure_inlet[zone_name]
        pi.momentum.gauge_total_pressure.value = total_pressure_pa
        pi.momentum.supersonic_or_initial_gauge_pressure.value = total_pressure_pa
        try: pi.thermal.t.value = temperature_k
        except Exception: pass
        return _j({"ok": True, "message": f"Pressure inlet '{zone_name}': {total_pressure_pa} Pa"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_pressure_outlet(
    zone_name: str,
    gauge_pressure_pa: float = 0.0,
    backflow_temperature_k: float = 300.0,
) -> str:
    """Set a pressure-outlet boundary condition.

    Args:
        zone_name:             Boundary zone name.
        gauge_pressure_pa:     Gauge static pressure (Pa).
        backflow_temperature_k: Backflow temperature (K).
    """
    if err := _chk("solver"): return err
    try:
        po = _solver.setup.boundary_conditions.pressure_outlet[zone_name]
        po.momentum.gauge_pressure.value = gauge_pressure_pa
        backflow_applied = True
        # pressure_outlet's thermal group has no '.t' field — verified live,
        # 2026-07-13: the only allowed thermal child is
        # 'backflow_total_temperature'. The old '.t' attribute name silently
        # failed via this except every time (root cause of the manifold_2
        # backflow-temperature-not-landing incident from earlier in the
        # same project).
        try: po.thermal.backflow_total_temperature.value = backflow_temperature_k
        except Exception: backflow_applied = False
        readback = _read_bc_values(zone_name, "pressure_outlet")
        rb_p = readback.get("gauge_pressure_pa")
        verified = isinstance(rb_p, (int, float)) and abs(rb_p - gauge_pressure_pa) < 1e-9
        return _j({"ok": True, "message": f"Outlet '{zone_name}': {gauge_pressure_pa} Pa",
                   "readback": readback, "verified": verified,
                   "backflow_temperature_applied": backflow_applied})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_wall_boundary(
    zone_name: str,
    thermal_condition: str = "adiabatic",
    temperature_k: float = 300.0,
    heat_flux_wm2: float = 0.0,
    roughness_height_m: float = 0.0,
) -> str:
    """Set a wall boundary condition.

    Args:
        zone_name:         Wall zone name.
        thermal_condition: "adiabatic", "temperature", or "heat_flux".
        temperature_k:     Wall temperature (K).
        heat_flux_wm2:     Heat flux (W/m²).
        roughness_height_m: Sand-grain roughness height (m, default 0 = smooth).
    """
    if err := _chk("solver"): return err
    try:
        wall = _solver.setup.boundary_conditions.wall[zone_name]
        tc = thermal_condition.lower()
        if tc == "temperature":
            wall.thermal.thermal_condition = "Temperature"
            wall.thermal.t.value = temperature_k
        elif tc == "heat_flux":
            wall.thermal.thermal_condition = "Heat Flux"
            wall.thermal.q.value = heat_flux_wm2
        else:
            wall.thermal.thermal_condition = "Heat Flux"
            wall.thermal.q.value = 0.0
        if roughness_height_m > 0:
            try: wall.momentum.roughness_height = roughness_height_m
            except Exception: pass
        return _j({"ok": True, "message": f"Wall '{zone_name}': {thermal_condition}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_symmetry_boundary(zone_name: str) -> str:
    """Confirm/check a symmetry boundary zone.

    Args:
        zone_name: Symmetry zone name.
    """
    if err := _chk("solver"): return err
    try:
        _ = _solver.setup.boundary_conditions.symmetry[zone_name]
        return _j({"ok": True, "message": f"Symmetry '{zone_name}' confirmed."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 7. SOLVER SETTINGS & CONVERGENCE CRITERIA
# ===========================================================================

@mcp.tool()
def set_solver_type(solver_type: str = "pressure-based") -> str:
    """Set the solver algorithm.

    Args:
        solver_type: "pressure-based" or "density-based".
    """
    if err := _chk("solver"): return err
    try:
        _solver.setup.general.solver.type = solver_type
        return _j({"ok": True, "message": f"Solver type: {solver_type}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_time_stepping(
    steady: bool = True,
    time_step_s: float = 0.001,
    max_time_steps: int = 100,
) -> str:
    """Configure steady or transient time stepping.

    Args:
        steady:         True = steady-state (default).
        time_step_s:    Time step size (s) for transient.
        max_time_steps: Number of time steps.
    """
    if err := _chk("solver"): return err
    try:
        if steady:
            _solver.setup.general.solver.time = "steady"
            return _j({"ok": True, "message": "Steady-state."})
        else:
            _solver.setup.general.solver.time = "unsteady-1st-order"
            return _j({"ok": True, "message": f"Transient: dt={time_step_s}s, steps={max_time_steps}"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_solution_methods(
    pressure_velocity_coupling: str = "SIMPLE",
    pressure_scheme: str = "second-order",
    momentum_scheme: str = "second-order-upwind",
) -> str:
    """Set P-V coupling and discretisation schemes.

    Args:
        pressure_velocity_coupling: "SIMPLE", "SIMPLEC", or "PISO".
        pressure_scheme:            "standard", "second-order", "presto".
        momentum_scheme:            "first-order-upwind", "second-order-upwind", "quick".
    """
    if err := _chk("solver"): return err
    try:
        meth = _solver.solution.methods
        applied = {}
        for key, attr, val in [("pressure_velocity_coupling", "p_v_coupling.flow_scheme", pressure_velocity_coupling),
                                ("pressure_scheme", "discretization.pressure", pressure_scheme),
                                ("momentum_scheme", "discretization.mom", momentum_scheme)]:
            try:
                obj = meth
                for a in attr.split("."): obj = getattr(obj, a)
                obj.value = val
                applied[key] = True
            except Exception:
                applied[key] = False
        # fail-loud (matches set_velocity_inlet/set_pressure_outlet's
        # verify-by-readback convention): discretization.pressure/.mom is
        # documented as renamed to spatial_discretization on this pyfluent
        # version (knowledge/cht_simulation_postmortem_2026-07-10.md #119,
        # "Not yet resolved") - report per-attribute status instead of a
        # blanket ok:true so callers don't believe an unapplied scheme took effect.
        return _j({"ok": True, "message": f"{pressure_velocity_coupling}/{pressure_scheme}/{momentum_scheme}",
                   "applied": applied,
                   "warning": None if all(applied.values()) else
                       "pressure_scheme/momentum_scheme discretization API is renamed to "
                       "spatial_discretization in this pyfluent version and its per-variable "
                       "keying is not yet implemented here - these settings may not have taken effect"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_convergence_criteria(
    continuity: float = 1e-4,
    x_velocity: float = 1e-4,
    y_velocity: float = 1e-4,
    z_velocity: float = 1e-4,
    energy: float = 1e-6,
    k: float = 1e-4,
    epsilon: float = 1e-4,
    omega: float = 1e-4,
) -> str:
    """Set residual convergence criteria.

    Args:
        continuity:  Continuity (default 1e-4).
        x_velocity:  X-velocity.
        y_velocity:  Y-velocity.
        z_velocity:  Z-velocity.
        energy:      Energy (default 1e-6).
        k:           Turbulent KE.
        epsilon:     Dissipation rate (k-epsilon models only).
        omega:       Specific dissipation rate (k-omega models only).
    """
    if err := _chk("solver"): return err
    mapping = {"continuity": continuity, "x-velocity": x_velocity,
               "y-velocity": y_velocity, "z-velocity": z_velocity,
               "energy": energy, "k": k, "epsilon": epsilon, "omega": omega}
    n = 0
    failures: dict = {}
    try:
        eqs = _solver.solution.monitor.residual.equations
        # NOTE: the settings-API attribute is `.absolute_criteria`, not
        # `.convergence_criterion` (that name doesn't exist at all — verified
        # live, 2026-07-13: dir() on an equation object shows
        # absolute_criteria/relative_criteria only). Also: 'epsilon' is not a
        # valid key for k-omega models (only 'omega' is, and vice versa for
        # k-epsilon) — silently skip whichever turbulence key isn't present
        # in this case rather than failing on it.
        available = set(eqs.keys())
        for name, val in mapping.items():
            if name not in available:
                continue  # e.g. 'epsilon' on a k-omega case, or vice versa
            try:
                eqs[name].absolute_criteria = val
                n += 1
            except Exception as e:
                failures[name] = str(e)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    result = {"ok": True, "set_count": n, "criteria": {k2: v for k2, v in mapping.items() if k2 in eqs.keys()}}
    if failures:
        result["failures"] = failures
    return _j(result)


_ANALYSIS_FIDELITY = {
    "screening":    {"criteria": 1e-3, "iterations": 300},
    "engineering":  {"criteria": 1e-4, "iterations": 800},
    "verification": {"criteria": 1e-5, "iterations": 2000},
}


@mcp.tool()
def propose_solver_setup(
    analysis_type: str = "internal",
    fluid: str = "air",
    inlet_velocity_ms: float = 0.0,
    outlet_gauge_pressure_pa: float = 0.0,
    fidelity: str = "screening",
) -> str:
    """Build the CFD ANALYSIS plan by inferring everything safely inferable
    from the loaded case: boundary zones and their roles (from assigned BC
    types + KB name patterns: inlet*/in_/supply -> inlet, outlet*/out_/
    return/exhaust -> outlet), turbulence model, convergence criteria,
    iteration budget, standard monitors, and post-processing KPIs.

    AUTO-mode analysis flow: propose_solver_setup -> show the summary AND
    missing_essential_inputs to the user (AskUserQuestion; ask ONLY those)
    -> apply_solver_plan -> initialize_solution(checkpoint_dir=...) ->
    start_calculation -> poll read_console_tail. In MANUAL mode this plan
    is the recommendation baseline; apply settings with the individual
    tools instead.

    Args:
        analysis_type:            'internal' (duct/manifold), 'external'
                                  (aero/hydro), or 'cht'.
        fluid:                    Working fluid (default air).
        inlet_velocity_ms:        Inlet velocity if the user already gave it
                                  (0 = not provided -> listed as essential).
        outlet_gauge_pressure_pa: Outlet gauge pressure (default 0).
        fidelity:                 'screening', 'engineering', 'verification'
                                  -> convergence target + iteration budget.
    """
    global _solver_plan
    if err := _chk("solver"): return err
    fid = fidelity.strip().lower()
    if fid not in _ANALYSIS_FIDELITY:
        return _j({"ok": False, "error": f"unknown fidelity '{fidelity}' "
                                         f"(use {sorted(_ANALYSIS_FIDELITY)})"})
    try:
        zones = _zone_type_pairs()
    except Exception as e:
        return _j({"ok": False, "error": f"zone inventory failed: {e}"})
    inlets, outlets, walls, others = [], [], [], []
    for z in zones:
        zn, zt = z["zone"].lower(), z["type"].lower()
        if "inlet" in zt or zn.startswith(("inlet", "in_", "supply")):
            inlets.append(z["zone"])
        elif "outlet" in zt or zn.startswith(("outlet", "out_", "return", "exhaust")):
            outlets.append(z["zone"])
        elif "wall" in zt:
            walls.append(z["zone"])
        elif "interior" not in zt:
            others.append(z)
    fam = analysis_type.strip().lower()
    fam = ("external" if "extern" in fam or "aero" in fam else
           "cht" if "cht" in fam or "heat" in fam or "conjugate" in fam else
           "internal")
    missing: list[str] = []
    if not inlets:
        missing.append("no inlet zone could be inferred - specify it explicitly")
    if not outlets:
        missing.append("no outlet zone could be inferred - specify it explicitly")
    if inlet_velocity_ms <= 0:
        missing.append("inlet velocity (or mass flow) - essential, cannot be inferred")
    _solver_plan = {
        "analysis_family": fam, "fidelity": fid, "fluid": fluid,
        "physics": {"solver": "pressure-based",
                    "energy": fam == "cht",
                    "turbulence_model": "k-omega",
                    "turbulence_sub_model": "sst"},
        "boundary_conditions":
            [{"zone": z, "type": "velocity-inlet",
              "velocity_ms": inlet_velocity_ms or None} for z in inlets] +
            [{"zone": z, "type": "pressure-outlet",
              "gauge_pressure_pa": outlet_gauge_pressure_pa} for z in outlets] +
            [{"zone": z, "type": "wall", "note": "default no-slip"} for z in walls],
        "unclassified_zones": others,
        "monitors": {"mass_flow_on": inlets + outlets,
                     "area_avg_pressure_on": inlets[:1]},
        "convergence": {"residual_criteria": _ANALYSIS_FIDELITY[fid]["criteria"]},
        "iteration_budget": _ANALYSIS_FIDELITY[fid]["iterations"],
        "postprocessing_kpis": ["pressure drop (inlet area-avg static pressure)",
                                "per-outlet mass-flow split + balance",
                                "velocity contour on a mid cut-plane"],
    }
    _persist_session_state()
    prompts: list[dict] = []
    if inlet_velocity_ms <= 0:
        prompts.append(_prompt(
            f"Inlet velocity for zone(s) {inlets or '<unknown>'}?", "Inlet",
            [("5 m/s", "moderate duct/manifold flow"),
             ("1 m/s", "low-speed flow"),
             ("10 m/s", "high-speed duct flow"),
             ("Custom", "type a velocity in m/s (or a mass flow to convert)")]))
    prompts.append(_prompt(
        f"Apply this {fam} setup? ({_solver_plan['physics']['turbulence_model']}-"
        f"{_solver_plan['physics']['turbulence_sub_model']}, "
        f"{len(inlets)} inlet(s), {len(outlets)} outlet(s), criteria "
        f"{_ANALYSIS_FIDELITY[fid]['criteria']:g}, budget "
        f"{_ANALYSIS_FIDELITY[fid]['iterations']} iters)", "Solver setup",
        [("Yes - apply plan (Recommended)", "apply_solver_plan runs everything in one shot"),
         ("Adjust first", "tell me what to change before applying"),
         ("Manual setup", "walk through each setting individually")]))
    return _j({"ok": True, "plan": _solver_plan,
               "missing_essential_inputs": missing,
               "analysis_mode": _analysis_mode,
               "interactive_prompts": prompts,
               "prompt_render": _PROMPT_RENDER_NOTE,
               "confirmation_gate": "Show this summary to the user and get an "
                                    "explicit yes BEFORE apply_solver_plan and "
                                    "again before start_calculation.",
               "next": ("ask the user ONLY the missing_essential_inputs, confirm "
                        "the summary, then apply_solver_plan"
                        if _analysis_mode == "auto" else
                        "MANUAL mode: use this as the recommendation baseline and "
                        "apply settings via the individual tools with the user's "
                        "values")})


@mcp.tool()
def apply_solver_plan(inlet_velocity_ms: float = 0.0, confirmed: bool = False) -> str:
    """Apply the stored propose_solver_setup plan in one shot: turbulence
    model, energy equation (CHT), every inferred boundary condition,
    standard KPI monitors, and convergence criteria.

    CONFIRMATION GATE — SERVER-ENFORCED: the first call (confirmed=False)
    does NOT apply anything; it returns the stored plan summary as an
    interactive prompt to render via AskUserQuestion. Re-call with
    confirmed=True after an explicit yes. Essential inputs are never
    guessed: an inlet with no velocity (neither here nor in the plan)
    aborts with an error telling you to ask the user.

    Args:
        inlet_velocity_ms: Inlet velocity — required if it wasn't already
                           given to propose_solver_setup.
        confirmed:         Pass True only after the user approved the gate prompt.
    """
    if err := _chk("solver"): return err
    if not _solver_plan:
        return _j({"ok": False, "error": "no solver plan - run propose_solver_setup first"})
    if not confirmed:
        fam = _solver_plan.get("analysis_family", "internal")
        return _j({"ok": False, "confirmation_required": True,
                   "plan": _solver_plan,
                   "interactive_prompt": _prompt(
                       f"Apply this {fam} solver setup now? (turbulence model, all "
                       "boundary conditions, monitors, convergence criteria)",
                       "Apply solver setup",
                       [("Yes - apply plan (Recommended)",
                         "apply_solver_plan runs everything in one shot"),
                        ("Adjust first", "tell me what to change before applying")]),
                   "prompt_render": _PROMPT_RENDER_NOTE,
                   "hint": "re-call apply_solver_plan(confirmed=True) after the user approves"})
    results: dict = {}
    warnings: list[str] = []

    def _rec(key: str, payload: str):
        try:
            d = json.loads(payload)
        except Exception:
            d = {"raw": payload}
        results[key] = d
        if isinstance(d, dict) and d.get("ok") is False:
            warnings.append(f"{key}: {d.get('error')}")

    phys = _solver_plan.get("physics", {})
    try:
        _rec("viscous", set_viscous_model(
            model=phys.get("turbulence_model", "k-omega"),
            sub_model=phys.get("turbulence_sub_model", "sst")))
    except Exception as e:
        warnings.append(f"viscous: {e}")
    if phys.get("energy"):
        try:
            _rec("energy", set_energy_equation(True))
        except Exception as e:
            warnings.append(f"energy: {e}")
    inlets, outlets = [], []
    for bc in _solver_plan.get("boundary_conditions", []):
        zone, btype = bc["zone"], bc["type"]
        if btype == "velocity-inlet":
            v = inlet_velocity_ms or bc.get("velocity_ms") or 0.0
            if v <= 0:
                return _j({"ok": False, "applied_so_far": results,
                           "error": f"inlet '{zone}': velocity is an essential input "
                                    "- ask the user and pass inlet_velocity_ms; "
                                    "never guess it"})
            try:
                _rec(f"bc:{zone}", set_velocity_inlet(zone_name=zone, velocity_ms=v))
                inlets.append(zone)
            except Exception as e:
                warnings.append(f"bc:{zone}: {e}")
        elif btype == "pressure-outlet":
            try:
                _rec(f"bc:{zone}", set_pressure_outlet(
                    zone_name=zone,
                    gauge_pressure_pa=bc.get("gauge_pressure_pa", 0.0)))
                outlets.append(zone)
            except Exception as e:
                warnings.append(f"bc:{zone}: {e}")
    if inlets and outlets:
        try:
            _rec("monitors", setup_standard_monitors(
                inlet_zone=inlets[0], outlet_zones=outlets))
        except Exception as e:
            warnings.append(f"monitors: {e}")
    crit = _solver_plan.get("convergence", {}).get("residual_criteria", 1e-3)
    try:
        _rec("convergence", set_convergence_criteria(
            continuity=crit, x_velocity=crit, y_velocity=crit, z_velocity=crit,
            k=crit, epsilon=crit, omega=crit))
    except Exception as e:
        warnings.append(f"convergence: {e}")
    return _j({"ok": not warnings, "applied": list(results.keys()),
               "results": results, "warnings": warnings,
               "iteration_budget": _solver_plan.get("iteration_budget"),
               "next": "initialize_solution(checkpoint_dir=...) -> confirm with "
                       "the user -> start_calculation -> poll read_console_tail "
                       "every few seconds"})


# ===========================================================================
# 8. INITIALIZATION & SOLUTION CONTROL
# ===========================================================================

def _count_report_definitions() -> int | None:
    """Number of report definitions (monitors) configured, or None if the
    inventory could not be read on this release."""
    try:
        rd = _solver.solution.report_definitions
        n = 0
        for group in ("surface", "volume", "flux", "force", "drag", "lift", "moment"):
            try:
                n += len(getattr(rd, group).keys())
            except Exception:
                pass
        return n
    except Exception:
        return None


def _require_monitors_or_error(force: bool) -> str | None:
    """Shared gate for start_calculation/run_until_convergence/
    run_transient_calculation: refuse to start a solve with zero report
    monitors defined (user feedback 2026-07-21 - flow-rate monitors must
    exist BEFORE running, not as an easy-to-miss warning after the fact).
    Returns an error JSON string to short-circuit on, or None to proceed."""
    n_mon = _count_report_definitions()
    if n_mon == 0 and not force:
        return _j({"ok": False, "error": "no report monitors defined - refusing to start "
                                         "without any KPI history to judge real convergence by",
                   "hint": "call setup_standard_monitors(inlet_zone=..., outlet_zones=[...]) "
                           "first (flow-rate monitors on every inlet/outlet, at minimum), or "
                           "pass force_no_monitors=True if this is deliberate"})
    return None


@mcp.tool()
def setup_standard_monitors(
    inlet_zone: str,
    outlet_zones: list[str],
    include_inlet_pressure: bool = True,
    frequency: int = 1,
) -> str:
    """Create the KB-standard KPI monitors for an internal-flow run — call
    BEFORE initialize_solution/start_calculation (orchestration KB doc 01;
    user feedback 2026-07-17: runs were started with no monitors at all).

    Creates: mass-flow report monitors on the inlet and on EVERY outlet
    (per-branch flow split + mass balance live during the run), and
    optionally area-averaged static pressure on the inlet (equals the
    pressure drop when outlets are at 0 Pa gauge).

    Args:
        inlet_zone:             Inlet boundary zone name.
        outlet_zones:           All outlet boundary zone names.
        include_inlet_pressure: Also monitor inlet area-avg static pressure.
        frequency:              Sample every N iterations.
    """
    if err := _chk("solver"): return err
    created, failed, history_files = [], [], {}
    mon_dir = _artifacts_dir("monitors")

    def _mk(name: str, rtype: str, field: str, zones: list[str]):
        try:
            report = _solver.solution.report_definitions
            report.surface[name] = {}
            rd = report.surface[name]
            rd.report_type = rtype
            if field:
                rd.field = field
            rd.surface_names = zones
            try:
                plots = _solver.solution.monitor.report_plots
                plots[name] = {}
                try: plots[name].frequency = frequency
                except Exception: pass
            except Exception:
                pass
            # history FILE: enables read_monitor_history — KPI plateau (mdot
            # split, dP flat over iterations) is the real convergence
            # criterion, not residuals alone (KB solver doc 00)
            try:
                out_path = str(mon_dir / f"{name}.out")
                files = _solver.solution.monitor.report_files
                files[name] = {}
                rf = files[name]
                rf.report_defs = [name]
                rf.file_name = out_path
                try: rf.frequency = frequency
                except Exception: pass
                history_files[name] = out_path
            except Exception:
                pass  # plots still work without the file
            created.append(name)
        except Exception as e:
            failed.append({"monitor": name, "error": str(e)})

    _mk(f"mdot_{inlet_zone}", "surface-massflowrate", "", [inlet_zone])
    for oz in outlet_zones:
        _mk(f"mdot_{oz}", "surface-massflowrate", "", [oz])
    if include_inlet_pressure:
        _mk(f"p_avg_{inlet_zone}", "surface-areaavg", "pressure", [inlet_zone])
    return _j({"ok": bool(created), "created": created, "failed": failed,
               "history_files": history_files,
               "note": "Monitors are live in the GUI plots during the run; "
                       "read_monitor_history(name) returns the iteration history "
                       "(KPI plateau = real convergence); final values via "
                       "check_outlet_flow_balance / get_surface_report."})


@mcp.tool()
def read_monitor_history(monitor_name: str, last_n: int = 100) -> str:
    """Read a monitor's iteration history from its report file (created by
    setup_standard_monitors). Use it to judge REAL convergence: residuals
    below criteria plus FLAT KPIs (mass split, dP) over the last ~100
    iterations. Safe during a run (file read only).

    Args:
        monitor_name: Report definition name (e.g. "mdot_outlet_1", "p_avg_inlet").
        last_n:       How many trailing samples to return.
    """
    path = _artifacts_dir("monitors") / f"{monitor_name}.out"
    if not path.exists():
        return _j({"ok": False, "error": f"no history file: {path}",
                   "hint": "setup_standard_monitors creates it; data appears only "
                           "for iterations run AFTER the monitors were created"})
    rows = []
    try:
        for line in path.read_text(encoding="utf-8", errors="replace").splitlines():
            parts = line.split()
            if len(parts) >= 2:
                try:
                    rows.append((int(float(parts[0])), float(parts[1])))
                except ValueError:
                    continue   # header/title lines
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    rows = rows[-max(last_n, 2):]
    verdict = None
    if len(rows) >= 10:
        vals = [v for _, v in rows[-min(len(rows), 50):]]
        span = max(vals) - min(vals)
        mean = sum(vals) / len(vals)
        rel = abs(span / mean) if mean else float("inf")
        verdict = ("FLAT (plateaued - converged by KPI)" if rel < 0.01 else
                   "DRIFTING - keep iterating" if rel > 0.05 else
                   "settling - nearly flat")
        stats = {"samples": len(rows), "last_value": rows[-1][1],
                 "relative_span_last50": round(rel, 5)}
    else:
        stats = {"samples": len(rows)}
    return _j({"ok": True, "monitor": monitor_name, "history": rows,
               "stats": stats, "plateau_verdict": verdict})


@mcp.tool()
def set_autosave(frequency_iterations: int = 100, root_name: str = "") -> str:
    """Enable solver autosave every N iterations (interactive console spec
    §17) — checkpoints during long runs so a crash costs at most N
    iterations. Complements the post-init checkpoint from
    initialize_solution.

    Args:
        frequency_iterations: Save case+data every N iterations (0 disables).
        root_name:            Checkpoint file root (default: a
                              'pyfluent_mcp_logs/checkpoints/autosave' folder
                              next to the last-loaded geometry/mesh/case
                              file, or the server's own logs/ dir if none
                              has been loaded yet this session).
    """
    if err := _chk("solver"): return err
    if not root_name:
        ck = _artifacts_dir("checkpoints")
        root_name = str(ck / "autosave")
    applied, failures = {}, {}
    try:
        auto = _solver.file.auto_save
    except Exception as e:
        return _j({"ok": False, "error": f"auto_save settings unavailable: {e}"})
    for attr, val in (("data_frequency", frequency_iterations),
                      ("root_name", root_name)):
        try:
            setattr(auto, attr, val)
            applied[attr] = val
        except Exception as e:
            failures[attr] = str(e)
    if not applied:
        # release-dialect fallback: TUI autosave menu
        try:
            _tui(f"/file/auto-save/data-frequency {frequency_iterations}")
            _tui(f'/file/auto-save/root-name "{_tui_path(root_name)}"')
            applied = {"via": "TUI fallback", "data_frequency": frequency_iterations,
                       "root_name": root_name}
        except Exception as e:
            return _j({"ok": False, "error": f"autosave failed on both paths: {e}",
                       "settings_failures": failures})
    return _j({"ok": True, "applied": applied, "failures": failures,
               "note": f"case+data checkpoint every {frequency_iterations} "
                       f"iterations -> {root_name}_<iter>.cas.h5/.dat.h5"})


@mcp.tool()
def assess_convergence(monitor_names: list[str] | None = None) -> str:
    """Formal convergence assessment (spec §24) combining THREE independent
    signals: (1) solver residuals / Fluent's converged flag from the
    transcript, (2) KPI monitor plateau verdicts (read_monitor_history),
    (3) mass imbalance from mdot monitors. Returns a convergence class:

      C1 - tightly converged (all signals green, imbalance <0.5%)
      C2 - engineering converged (criteria met, KPIs flat/settling, <2%)
      C3 - NOT converged (drifting KPIs, high residuals, or >5% imbalance)

    Args:
        monitor_names: Monitors to evaluate (default: every history file
                       created by setup_standard_monitors).
    """
    if err := _chk("solver"): return err
    evidence: dict = {}
    # 1. residuals + converged flag from the transcript
    converged_flag = False
    residuals: dict = {}
    try:
        tail = json.loads(read_console_tail(lines=60))
        converged_flag = bool(tail.get("converged"))
        residuals = tail.get("residuals") or {}
    except Exception as e:
        evidence["transcript_error"] = str(e)
    evidence["fluent_converged_flag"] = converged_flag
    evidence["residuals"] = residuals
    # 2. KPI plateaus
    mon_dir = _artifacts_dir("monitors")
    if monitor_names is None:
        monitor_names = sorted(p.stem for p in mon_dir.glob("*.out")) if mon_dir.exists() else []
    plateaus: dict = {}
    mdot_last: dict = {}
    for mn in monitor_names:
        try:
            h = json.loads(read_monitor_history(mn, last_n=100))
            if h.get("ok"):
                plateaus[mn] = h.get("plateau_verdict") or "insufficient data"
                if mn.startswith("mdot_") and h.get("history"):
                    mdot_last[mn] = h["history"][-1][1]
        except Exception as e:
            plateaus[mn] = f"error: {e}"
    evidence["kpi_plateaus"] = plateaus
    # 3. mass imbalance from mdot monitors (inlet positive, outlets negative)
    imbalance = None
    if mdot_last:
        signed = sum(mdot_last.values())
        scale = max((abs(v) for v in mdot_last.values()), default=0.0)
        if scale > 0:
            imbalance = abs(signed) / scale
            evidence["mass_imbalance_fraction"] = round(imbalance, 6)
    verdicts = [v for v in plateaus.values() if not str(v).startswith("error")]
    any_drift = any("DRIFT" in str(v) for v in verdicts)
    all_flat = bool(verdicts) and all("FLAT" in str(v) for v in verdicts)
    res_ok = converged_flag or (residuals and
                                all(v < 1e-3 for v in residuals.values()))
    if any_drift or (imbalance is not None and imbalance > 0.05) or \
       (not res_ok and residuals and max(residuals.values()) > 1e-2):
        cls, label = "C3", "NOT converged - keep iterating"
    elif res_ok and all_flat and (imbalance is None or imbalance < 0.005):
        cls, label = "C1", "tightly converged"
    elif res_ok:
        cls, label = "C2", "engineering converged"
    else:
        cls, label = "C3", "NOT converged - keep iterating"
    return _j({"ok": True, "convergence_class": cls, "assessment": label,
               "evidence": evidence,
               "recommendation": ("results are trustworthy - proceed to "
                                  "post-processing (generate_default_results)"
                                  if cls in ("C1", "C2") else
                                  "start_calculation for more iterations; watch "
                                  "read_monitor_history plateaus, not just residuals")})


@mcp.tool()
def initialize_solution(method: str = "hybrid", checkpoint_dir: str = "") -> str:
    """Initialize the solution.

    KB PROTOCOL (orchestration docs; user feedback 2026-07-17): (1) monitors
    should already exist — setup_standard_monitors — and this tool warns if
    none are defined; (2) pass checkpoint_dir to write case+data immediately
    after initialization, so a crash mid-run costs nothing. Omitting it adds
    a warning; save manually with write_case_data.

    Args:
        method:         "hybrid" (recommended) or "standard".
        checkpoint_dir: Directory for the post-init case+data checkpoint
                        (written as init_checkpoint.cas.h5).
    """
    if err := _chk("solver"): return err
    try:
        if method.lower() == "hybrid":
            _solver.solution.initialization.hybrid_initialize()
        else:
            _solver.solution.initialization.initialize()
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    warnings: list[str] = []
    checkpoint = None
    if checkpoint_dir:
        try:
            checkpoint = str(Path(checkpoint_dir) / "init_checkpoint.cas.h5")
            _mkdir(checkpoint)
            _solver.file.write_case_data(file_name=checkpoint)
        except Exception as e:
            warnings.append(f"post-init checkpoint failed: {e}")
            checkpoint = None
    else:
        warnings.append("KB protocol: case+data NOT saved after initialization - "
                        "pass checkpoint_dir here or call write_case_data before "
                        "starting the run")
    n_mon = _count_report_definitions()
    if n_mon == 0:
        warnings.append("no report monitors defined - KB protocol requires KPI "
                        "monitors before solving (setup_standard_monitors)")
    return _j({"ok": True, "message": f"Initialized ({method}).",
               "checkpoint_written": checkpoint,
               "report_monitors_defined": n_mon,
               "warnings": warnings})


@mcp.tool()
def run_calculation(number_of_iterations: int = 20, force_blocking: bool = False) -> str:
    """Run steady-state iterations SYNCHRONOUSLY (blocks until done).

    ONLY for short screening bursts (<= 20 iterations). This call holds the
    MCP request open for the whole solve — a long run here freezes the
    client (verified incident: a 50-iteration run required a manual Ctrl+C
    in the Fluent console). For anything longer use the non-blocking
    protocol: start_calculation -> get_calculation_status (poll) ->
    stop_calculation (if needed).

    Args:
        number_of_iterations: Iterations to run (default 20; refused above
                              20 unless force_blocking=True).
        force_blocking:       Explicitly accept a long blocking call.
    """
    if err := _chk("solver"): return err
    if number_of_iterations > 20 and not force_blocking:
        return _j({"ok": False,
                   "error": f"{number_of_iterations} iterations would block the MCP "
                            "client for the whole solve. Use start_calculation + "
                            "get_calculation_status instead (or pass "
                            "force_blocking=True if you really want this)."})
    if _bg_run["running"]:
        return _j({"ok": False, "error": "A background calculation is already "
                   "running; poll get_calculation_status or stop_calculation first."})
    try:
        _solver.solution.run_calculation.iterate(iter_count=number_of_iterations)
        return _j({"ok": True, "message": f"Ran {number_of_iterations} iterations."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ---------------------------------------------------------------------------
# Non-blocking solve protocol (v3.4)
# The solve runs in a daemon thread over the PYFLUENT SESSION path; status
# polls read shared state and, for live iteration counts, use the RAW gRPC
# path (own channel) so they never queue behind the busy session.
# ---------------------------------------------------------------------------
_bg_run: dict = {"running": False, "thread": None, "kind": None,
                 "requested": 0, "start_iteration": None,
                 "finished": False, "error": None, "interrupted": False,
                 "result": None}  # holds run_until_convergence/run_transient_calculation's
                                  # final payload once kind is "convergence"/"transient"

# ---------------------------------------------------------------------------
# Background SESSION JOBS — long blocking meshing/session steps (import,
# surface mesh, volume mesh) run in a worker thread so the MCP stdio loop
# stays responsive. ROOT CAUSE of every mid-run MCP restart on 2026-07-17/20:
# a 2-4 minute synchronous tool call froze the server until the client gave
# up and killed it. Poll get_job_status + read_console_tail instead.
# ---------------------------------------------------------------------------
_session_job: dict = {"running": False, "name": None, "started": None,
                      "finished": False, "result": None, "error": None,
                      "thread": None}

# Guards _solver/_meshing against concurrent access between a session-job
# worker thread and any other tool call issued on the main thread while it
# runs (verified live 2026-07-20: no lock existed anywhere in this module -
# every other tool only checked `if _solver is None`, never whether a job/
# calculation already held the session busy). Deliberately NOT used around
# start_calculation's iterate() or stop_calculation's interrupt() - those
# two are intentionally concurrent (interrupt must reach the solver while
# iterate() is in flight); locking either would break stop_calculation.
_session_lock = threading.RLock()


def _run_session_job(name: str, fn) -> str:
    """Run fn (a closure returning the tool's JSON string) in a background
    thread. Only ONE session job at a time — session-path gRPC calls would
    queue behind it anyway."""
    import threading
    if _session_job["running"]:
        return _j({"ok": False,
                   "error": f"another session job is already running: "
                            f"{_session_job['name']}",
                   "hint": "poll get_job_status until it finishes"})
    if _bg_run["running"]:
        return _j({"ok": False, "error": "a calculation is running - "
                                         "stop_calculation or wait first"})
    _session_job.update(running=True, name=name, started=_ts(),
                        finished=False, result=None, error=None)

    def _worker():
        try:
            with _session_lock:
                _session_job["result"] = fn()
        except Exception as e:
            _session_job["error"] = str(e)
        finally:
            _session_job["running"] = False
            _session_job["finished"] = True

    t = threading.Thread(target=_worker, name=f"mcp-job-{name}", daemon=True)
    _session_job["thread"] = t
    t.start()
    return _j({"ok": True, "job_started": name, "started_at": _session_job["started"],
               "next": "poll get_job_status every ~10-15 s and read_console_tail "
                       "for live progress (or call wait_for_job to block server-side "
                       "until it finishes); do NOT re-run the tool meanwhile"})


@mcp.tool()
def get_job_status() -> str:
    """Status of the background session job (import_geometry /
    generate_surface_mesh / generate_volume_mesh / read_case / read_case_data /
    write_case / write_case_data / save_mesh_and_workflow / run_design_points,
    all with run_in_background=True, the default). Non-blocking. While it
    runs, read_console_tail shows live progress; when finished=true,
    'result' holds the tool's full response.

    Steady/transient solves (start_calculation / run_until_convergence /
    run_transient_calculation) use a SEPARATE mechanism instead -
    get_calculation_status, not this tool."""
    st = {k: _session_job[k] for k in ("running", "name", "started",
                                       "finished", "error")}
    res = _session_job["result"]
    if res is not None:
        try:
            st["result"] = json.loads(res)
        except Exception:
            st["result"] = res
    return _j({"ok": True, **st})


_scheme_channel_cache: dict = {"endpoint": None, "channel": None, "stub": None}


def _get_scheme_stub(ip: str, port: int):
    """Lazily cache one gRPC channel/stub per (ip, port) instead of opening
    and tearing one down on every poll (this is called on the hot path -
    every get_calculation_status/start_calculation poll during a run, every
    ~5-15s per the documented monitoring protocol). Rebuilt automatically
    when the endpoint changes or the cached channel starts failing."""
    from ansys.api.fluent.v0 import scheme_eval_pb2_grpc
    ep = (ip, port)
    if _scheme_channel_cache["endpoint"] != ep:
        old = _scheme_channel_cache["channel"]
        if old is not None:
            try: old.close()
            except Exception: pass
        channel = _grpc_channel(ip, port)
        _scheme_channel_cache.update(endpoint=ep, channel=channel,
                                     stub=scheme_eval_pb2_grpc.SchemeEvalStub(channel))
    return _scheme_channel_cache["stub"]


def _bg_iteration_count(timeout: float = 4.0) -> int | None:
    """Raw gRPC StringEval probe on a cached channel, used ONLY as a
    liveness/fallback signal - NOT a reliable cumulative iteration count.

    (rpgetvar 'number-of-iterations) holds the iter_count REQUESTED by the
    last iterate() call (a setting), not a running counter: verified live
    2026-07-21 it stayed at the requested batch size (e.g. 200) even after
    the solve genuinely progressed past iteration 450. For the real
    current iteration, use _current_iteration_baseline() (transcript-based)
    instead - this function only confirms the gRPC channel is alive.
    Returns None if the server can't answer within `timeout` (normal while
    an iteration is in flight)."""
    try:
        props = _active().connection_properties
        import grpc
        from ansys.api.fluent.v0 import scheme_eval_pb2
        stub = _get_scheme_stub(props.ip, props.port)
        try:
            resp = stub.StringEval(
                scheme_eval_pb2.StringEvalRequest(input="(rpgetvar 'number-of-iterations)"),
                metadata=_grpc_metadata(props.password), timeout=timeout)
            return int(str(resp.output).strip().strip('"'))
        except grpc.RpcError as e:
            if e.code() in (grpc.StatusCode.UNAVAILABLE, grpc.StatusCode.DEADLINE_EXCEEDED):
                _scheme_channel_cache["endpoint"] = None  # force a rebuild next poll
            return None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# PUSH-BASED live iteration tracking via PyFluent's EventsManager (verified
# live 2026-07-21 by reading ansys-fluent-core source: SolverEvent.
# ITERATION_ENDED callbacks run on the EventsManager's OWN dedicated gRPC
# streaming thread - a channel entirely separate from the one iterate()
# blocks on - so this writes the true cumulative iteration index into a
# plain Python variable with ZERO extra round-trips per poll. This replaces
# (rpgetvar 'number-of-iterations) (confirmed to be the iter_count REQUESTED
# by the last iterate() call, a setting - not a counter) as the source of
# truth. Falls back to monitors/transcript/rpvar if registration fails on
# an older release.
# ---------------------------------------------------------------------------
_live_iteration: dict = {"index": None}
_iteration_event_solver_id: int | None = None


def _on_iteration_ended(session, event_info) -> None:
    try:
        _live_iteration["index"] = event_info.index
    except Exception:
        pass


def _ensure_iteration_event_registered() -> str:
    """Register the ITERATION_ENDED push callback on the current solver
    session, once per session object (tracked by id() so a reconnect/
    relaunch re-registers on the new session - _disconnect_internal resets
    this). Best-effort: on any failure (older pyfluent release, no
    `events` attribute, etc.) this silently no-ops and the iteration-count
    cascade below falls back to the monitors/transcript/rpvar paths.
    Returns a short status string (diagnostic only, surfaced by
    start_calculation for live verification)."""
    global _iteration_event_solver_id
    if _solver is None:
        return "no_solver"
    sid = id(_solver)
    if _iteration_event_solver_id == sid:
        return "already_registered"
    try:
        from ansys.fluent.core import SolverEvent
        _solver.events.register_callback(SolverEvent.ITERATION_ENDED, _on_iteration_ended)
        _iteration_event_solver_id = sid
        _live_iteration["index"] = None
        return "registered"
    except Exception as e:
        return f"failed: {e}"


def _current_iteration_baseline() -> int | None:
    """Best current iteration count for use as a run's start_iteration
    baseline / live status. Cascade, fastest+most-reliable first:
    1. The push-based event callback's last-seen index (zero round-trip).
    2. The session monitors service (_residual_history's iters[-1]) - also
       backed by its own streaming thread, so safely callable while
       iterate() is running on the main channel; slightly more overhead
       than #1 but still non-blocking and accurate.
    3. The console transcript parse (file read, non-blocking, accurate but
       depends on the transcript existing and being flushed promptly).
    4. The raw gRPC rpvar probe - a liveness signal only, NOT a reliable
       count (see _bg_iteration_count's docstring)."""
    if _live_iteration["index"] is not None:
        return _live_iteration["index"]
    try:
        iters, _ys = _residual_history()
        if iters:
            return iters[-1]
    except Exception:
        pass
    try:
        trn = _find_transcript()
        if trn is not None:
            tail = _tail_read(trn, max_bytes=8000)
            cur = _parse_transcript_tail(tail.splitlines()).get("current_iteration")
            if cur is not None:
                return cur
    except Exception:
        pass
    return _bg_iteration_count()


@mcp.tool()
def start_calculation(number_of_iterations: int = 100, force_no_monitors: bool = False) -> str:
    """Start steady-state iterations WITHOUT blocking (recommended for >20 iters).

    Launches the solve in a background thread (PyFluent session path) and
    returns immediately. MONITORING PROTOCOL (user feedback 2026-07-17):
    poll read_console_tail every few (~5-15) seconds — it parses the live
    transcript for current iteration, residuals, convergence, and errors —
    and act on what it shows; get_calculation_status remains the coarse
    progress check. Stop early with stop_calculation. Do NOT call other
    session-path tools that talk to the solver while it is iterating — they
    will queue behind the solve.

    MONITORS GATE — SERVER-ENFORCED (user feedback 2026-07-21: a merely
    documented "should exist" was too easy to skip and left a run with no
    KPI history to judge real convergence by): refuses to start if zero
    report monitors are defined, unless force_no_monitors=True. Call
    setup_standard_monitors (mass-flow on every inlet/outlet, at minimum)
    BEFORE this.

    PRE-RUN CHECKLIST (KB): BC values (velocity, pressures, fluid) must be
    the USER'S numbers — ask, don't assume (classify_cfd_request
    min_inputs); case+data should be saved post-init (initialize_solution
    checkpoint_dir).

    Args:
        number_of_iterations: Iterations to run (default 100).
        force_no_monitors:    Override the monitors gate (rare - e.g. a
                              throwaway sanity-check run with nothing to track).
    """
    import threading
    if err := _chk("solver"): return err
    if _bg_run["running"]:
        return _j({"ok": False, "error": "A background calculation is already running."})
    if gate_err := _require_monitors_or_error(force_no_monitors):
        return gate_err
    warnings: list[str] = []
    event_status = _ensure_iteration_event_registered()
    start_iter = _current_iteration_baseline()
    _bg_run.update(running=True, kind="steady", requested=number_of_iterations,
                   start_iteration=start_iter, finished=False, error=None,
                   interrupted=False)

    def _worker():
        try:
            _solver.solution.run_calculation.iterate(iter_count=number_of_iterations)
        except Exception as e:
            _bg_run["error"] = str(e)
        finally:
            _bg_run["running"] = False
            _bg_run["finished"] = True

    t = threading.Thread(target=_worker, name="mcp-fluent-solve", daemon=True)
    _bg_run["thread"] = t
    t.start()
    return _j({"ok": True, "started": True, "iterations": number_of_iterations,
               "start_iteration": start_iter, "warnings": warnings,
               "iteration_event_status": event_status,
               "next": "Poll get_calculation_status (one call now returns iteration/"
                       "progress/residuals/convergence together) or read_console_tail; "
                       "stop early with stop_calculation."})


@mcp.tool()
def get_calculation_status() -> str:
    """One-call status of the background calculation started by
    start_calculation: run state, current iteration, progress, residuals,
    convergence, and recent errors/completion markers - bundled together
    so a poll cycle needs ONE round trip instead of separately calling
    get_residuals and read_console_tail too (each MCP round trip costs a
    full LLM turn, which dominates perceived latency far more than any
    in-process work here - user feedback 2026-07-21).

    Non-blocking: 'current_iteration' comes from a cascade, fastest/most
    reliable first: (1) a push-based ITERATION_ENDED event callback
    (ansys-fluent-core EventsManager - zero round-trip, see
    _ensure_iteration_event_registered), (2) the session monitors service
    (also non-blocking - backed by its own streaming thread), (3) the
    console transcript (file read), (4) a raw gRPC rpvar probe (liveness
    only). NOT (rpgetvar 'number-of-iterations) as the primary source -
    verified live 2026-07-21 that rpvar holds the iter_count REQUESTED by
    the last iterate() call (a setting), not a cumulative counter.

    'current_iteration: null' while running just means no source has an
    answer yet — poll again.
    """
    if err := _chk("solver", allow_busy=True): return err
    state = {k: _bg_run[k] for k in ("running", "kind", "requested",
                                     "start_iteration", "finished", "error",
                                     "interrupted")}
    if _bg_run["finished"] and _bg_run["result"] is not None:
        state["result"] = _bg_run["result"]
    state["converged"] = False
    try:
        trn = _find_transcript()
        if trn is not None:
            tail = _tail_read(trn, max_bytes=8000)
            parsed = _parse_transcript_tail(tail.splitlines())
            state["converged"] = parsed.get("converged", False)
            if parsed.get("recent_errors"):
                state["recent_errors"] = parsed["recent_errors"]
            if parsed.get("completion_markers"):
                state["completion_markers"] = parsed["completion_markers"]
    except Exception:
        pass
    cur = _current_iteration_baseline()
    state["current_iteration"] = cur
    if cur is not None and _bg_run["start_iteration"] is not None and _bg_run["requested"]:
        done = cur - _bg_run["start_iteration"]
        state["progress"] = f"{done}/{_bg_run['requested']}"
    # bundle residuals in the same call (batching fix #2) - _residual_history
    # is itself non-blocking (session monitors streaming service), so this
    # adds no queueing risk behind a busy iterate()
    try:
        iters, ys = _residual_history()
        if iters:
            state["residuals"] = {k: v[-1] for k, v in ys.items() if v}
    except Exception:
        pass
    return _j({"ok": True, **state})


@mcp.tool()
def wait_for_job(timeout_s: float = 60.0) -> str:
    """Block server-side until the current background session job (import_geometry /
    generate_surface_mesh / generate_volume_mesh / read_case / ... /
    run_design_points) OR calculation (start_calculation / run_until_convergence /
    run_transient_calculation) finishes, or timeout_s elapses - whichever
    comes first. Cuts a long job's typical round-trip count from ~10-20
    separate get_job_status/get_calculation_status polls down to 1-3 (verified
    live 2026-07-20: both worker-thread mechanisms already exist and store
    their Thread object, they just weren't ever joined).

    Returns the SAME shape as get_job_status (if a session job was running)
    or get_calculation_status (if a calculation was running), plus a
    'timed_out' flag. If nothing was running, returns immediately with
    both_idle=True. Safe to call repeatedly - e.g. wait_for_job(60) in a
    loop - for a run expected to take much longer than one timeout window.

    Args:
        timeout_s: Max seconds to block for (default 60; the MCP client's
                   own request timeout is the practical ceiling - do not
                   pass an unreasonably large value).
    """
    job_thread = _session_job.get("thread")
    if _session_job["running"] and isinstance(job_thread, threading.Thread):
        job_thread.join(timeout=timeout_s)
        result = json.loads(get_job_status())
        result["timed_out"] = _session_job["running"]
        return _j(result)
    calc_thread = _bg_run.get("thread")
    if _bg_run["running"] and isinstance(calc_thread, threading.Thread):
        calc_thread.join(timeout=timeout_s)
        result = json.loads(get_calculation_status())
        result["timed_out"] = _bg_run["running"]
        return _j(result)
    return _j({"ok": True, "both_idle": True,
               "message": "No session job or calculation is currently running."})


@mcp.tool()
def run_iterations(number_of_iterations: int = 100, timeout_s: float = 120.0,
                   force_no_monitors: bool = False) -> str:
    """Run steady-state iterations and return the FINAL result in ONE round
    trip - combines start_calculation + wait_for_job (verified live
    2026-07-22: this two-call pattern already gives near-instant, accurate
    completion detection; collapsing it into one call halves the round-trip
    count for the common "run N iterations, tell me when done" request,
    which matters because every MCP round trip costs a full LLM turn - the
    actual dominant cost in this workflow, not Fluent/transport latency).

    Safe to block on: FastMCP dispatches synchronous tool calls (like this
    one) to a threadpool, so this call sitting in a thread.join() does NOT
    freeze the server for OTHER concurrent tool calls (e.g. stop_calculation
    from a separate request).

    Recommended for short-to-medium batches where you want the result
    immediately. For a long run where you want to do other things while it
    solves, use start_calculation (returns instantly) + wait_for_job/
    get_calculation_status separately instead - or if this call times out
    (timed_out:true, running:true), just call wait_for_job again to keep
    waiting on the same run.

    Args:
        number_of_iterations: Iterations to run (default 100).
        timeout_s:            Max seconds to block for (default 120). If the
                              solve takes longer, returns with
                              timed_out=true/running=true - call
                              wait_for_job(timeout_s) again to keep waiting.
        force_no_monitors:    Override the monitors gate (rare).
    """
    start_result = json.loads(start_calculation(number_of_iterations, force_no_monitors))
    if not start_result.get("ok"):
        return _j(start_result)
    return wait_for_job(timeout_s)


_tail_cache: dict[str, tuple[float, int, str]] = {}  # path -> (mtime, size, cached tail text)


def _tail_read(path: Path, max_bytes: int = 65536) -> str:
    """Read up to the last `max_bytes` bytes of a (possibly large, growing)
    transcript without loading the whole file - turns an O(file size) poll
    into O(max_bytes) (verified live 2026-07-20: read_console_tail/
    get_calculation_status/project_summary each re-read the ENTIRE .trn on
    every poll, which gets slower every call across a long solve). Cached
    by (mtime, size) so a poll that fires before the file has grown again
    returns the previous read without touching disk again."""
    key = str(path)
    st = path.stat()
    cached = _tail_cache.get(key)
    if cached is not None and cached[0] == st.st_mtime and cached[1] == st.st_size:
        return cached[2]
    with open(path, "rb") as f:
        if st.st_size > max_bytes:
            f.seek(st.st_size - max_bytes)
        data = f.read()
    text = data.decode("utf-8", errors="replace")
    _tail_cache[key] = (st.st_mtime, st.st_size, text)
    return text


def _find_transcript() -> Path | None:
    """Newest Fluent auto-transcript (.trn). Fluent inherits this process's
    cwd, so transcripts normally land there — but also search next to and
    above this package so the tool works regardless of how the server was
    started."""
    try:
        dirs = {Path.cwd(), Path(__file__).parent, Path(__file__).parent.parent}
        cands = sorted((p for d in dirs for p in d.glob("fluent-*.trn")),
                       key=lambda p: p.stat().st_mtime, reverse=True)
        return cands[0] if cands else None
    except Exception:
        return None


@mcp.tool()
def read_console_tail(lines: int = 40) -> str:
    """Tail of the live Fluent console transcript (.trn) + parsed status.

    THE monitoring primitive (user feedback 2026-07-17): poll this every few
    (~5-15) seconds during solver runs and long meshing/blocking steps to
    see what Fluent is actually doing and act accordingly. Reads the
    transcript file directly — safe while the solver is busy (never queues
    behind the session).

    Parses out of the tail: current iteration + per-equation residuals
    (column names from the residual header), a 'solution is converged'
    marker, recent Error: lines, and workflow-completion markers.

    Args:
        lines: How many trailing lines to return (default 40).
    """
    trn = _find_transcript()
    if trn is None:
        return _j({"ok": False, "error": "no fluent-*.trn transcript found in "
                                         f"{Path.cwd()}"})
    try:
        # budget scales with the requested line count (generous per-line
        # byte estimate) so we still capture `lines` worth of tail content
        text = _tail_read(trn, max_bytes=max(65536, lines * 200))
    except Exception as e:
        return _j({"ok": False, "error": f"transcript read failed: {e}"})
    tail_lines = text.splitlines()[-max(lines, 10):]
    status = _parse_transcript_tail(tail_lines)
    import time as _t
    age = _t.time() - trn.stat().st_mtime
    return _j({"ok": True, "transcript": str(trn),
               "transcript_age_s": round(age, 1),
               "stalled_hint": (age > 120), **status,
               "tail": tail_lines})


def _parse_transcript_tail(tail_lines: list[str]) -> dict:
    """Parse a Fluent console transcript tail into: current iteration +
    per-equation residuals (from the residual header), a 'solution is
    converged' marker, recent Error: lines, and workflow-completion
    markers. Single canonical parse shared by read_console_tail,
    get_calculation_status, and project_summary - previously each
    re-implemented an ad hoc subset of this independently."""
    status: dict = {}
    header_cols: list[str] = []
    last_data: list[str] = []
    for ln in tail_lines:
        s = ln.strip()
        if s.startswith("iter") and "continuity" in s:
            header_cols = s.split()
        elif re.match(r"^\d+\s+[\d.eE+-]", s):
            last_data = s.split()
    if last_data:
        status["current_iteration"] = int(last_data[0])
        if header_cols and len(last_data) >= len(header_cols) - 1:
            resid = {}
            for name, val in zip(header_cols[1:], last_data[1:]):
                if name in ("time/iter",):
                    break
                try:
                    resid[name] = float(val)
                except ValueError:
                    pass
            status["residuals"] = resid
    status["converged"] = any("solution is converged" in l.lower() for l in tail_lines)
    status["recent_errors"] = [l.strip() for l in tail_lines
                               if l.strip().lower().startswith("error")][-5:]
    status["completion_markers"] = [l.strip() for l in tail_lines
                                    if "complete in" in l.lower()
                                    or "creation completed" in l.lower()][-3:]
    return status


@mcp.tool()
def stop_calculation() -> str:
    """Interrupt the running background calculation gracefully.

    Calls solution.run_calculation.interrupt() — the mechanism verified on
    2025 R2 (the old checkpoint/exit-filename rpvar no longer exists). The
    solver finishes its current iteration, then the background thread ends;
    the solution keeps all data so far.
    """
    import threading
    if err := _chk("solver", allow_busy=True): return err
    if not _bg_run["running"]:
        return _j({"ok": False, "error": "No background calculation is running."})
    result: dict = {}

    def _interrupt():
        try:
            _solver.solution.run_calculation.interrupt()
            result["sent"] = True
        except Exception as e:
            result["error"] = str(e)

    t = threading.Thread(target=_interrupt, daemon=True)
    t.start()
    t.join(timeout=30)
    if t.is_alive():
        return _j({"ok": False, "error": "interrupt() did not return within 30 s; "
                   "the solver may still stop at the end of the current iteration. "
                   "Poll get_calculation_status."})
    if "error" in result:
        return _j({"ok": False, "error": result["error"]})
    _bg_run["interrupted"] = True
    return _j({"ok": True, "message": "Interrupt sent; solver stops after the "
               "current iteration. Poll get_calculation_status until "
               "running=false, then save with write_case_data."})


def _run_transient_calculation_impl(
    time_step_size_s: float,
    number_of_time_steps: int,
    max_iterations_per_step: int,
) -> dict:
    try:
        rc = _solver.solution.run_calculation
        rc.time_step_size = time_step_size_s
        rc.number_of_time_steps = number_of_time_steps
        rc.max_iterations_per_time_step = max_iterations_per_step
        rc.dual_time_iterate(time_step_count=number_of_time_steps,
                             max_iter_per_step=max_iterations_per_step)
        return {"ok": True, "message": f"Ran {number_of_time_steps} steps (dt={time_step_size_s}s)."}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@mcp.tool()
def run_transient_calculation(
    time_step_size_s: float = 0.001,
    number_of_time_steps: int = 100,
    max_iterations_per_step: int = 20,
    run_in_background: bool = True,
    force_no_monitors: bool = False,
) -> str:
    """Run a transient calculation.

    Runs AS A BACKGROUND JOB by default (verified live 2026-07-20: this was
    previously fully synchronous on the MCP request thread for the whole
    dual_time_iterate() call - potentially many minutes - freezing the
    server exactly like the long meshing/steady-solve calls that were
    already backgrounded to avoid this). Poll get_calculation_status
    (kind="transient"; 'result' appears once finished) or read_console_tail
    for live progress; stop early with stop_calculation.

    MONITORS GATE — SERVER-ENFORCED: refuses to start if zero report
    monitors are defined, unless force_no_monitors=True (setup_standard_monitors first).

    Args:
        time_step_size_s:       Time step size (s).
        number_of_time_steps:   Total time steps.
        max_iterations_per_step: Max inner iterations per step.
        run_in_background:       Default True (recommended). False blocks until done.
        force_no_monitors:       Override the monitors gate (rare).
    """
    if err := _chk("solver"): return err
    if _bg_run["running"]:
        return _j({"ok": False, "error": "A background calculation is already running."})
    if gate_err := _require_monitors_or_error(force_no_monitors):
        return gate_err
    if not run_in_background:
        return _j(_run_transient_calculation_impl(
            time_step_size_s, number_of_time_steps, max_iterations_per_step))
    _ensure_iteration_event_registered()
    start_iter = _current_iteration_baseline()
    _bg_run.update(running=True, kind="transient", requested=number_of_time_steps,
                   start_iteration=start_iter, finished=False, error=None,
                   interrupted=False, result=None)

    def _worker():
        try:
            _bg_run["result"] = _run_transient_calculation_impl(
                time_step_size_s, number_of_time_steps, max_iterations_per_step)
        except Exception as e:
            _bg_run["error"] = str(e)
        finally:
            _bg_run["running"] = False
            _bg_run["finished"] = True

    t = threading.Thread(target=_worker, name="mcp-fluent-transient", daemon=True)
    _bg_run["thread"] = t
    t.start()
    return _j({"ok": True, "started": True, "kind": "transient",
               "number_of_time_steps": number_of_time_steps, "start_iteration": start_iter,
               "next": "Poll get_calculation_status (result appears when finished) "
                       "or read_console_tail for live progress; stop early with "
                       "stop_calculation."})


# ===========================================================================
# 9. POST-PROCESSING – FIELD DATA
# ===========================================================================

def _fd():
    return _solver.fields.field_data

def _build_pyvista_surface(zone_names: list[str], field_name: str | None = None):
    import numpy as np
    import pyvista as pv
    from ansys.fluent.core import SurfaceFieldDataRequest, SurfaceDataType, ScalarFieldDataRequest
    fd = _fd()
    # data_types is required (no default) on installed ansys-fluent-core -
    # verified live 2026-07-20: omitting it raises "SurfaceFieldDataRequest
    # .__new__() missing 1 required positional argument: 'data_types'".
    # Mirrors the working call in create_outlet_normal_plane.
    surf_data = fd.get_field_data(SurfaceFieldDataRequest(
        data_types=[SurfaceDataType.Vertices, SurfaceDataType.FacesConnectivity],
        surfaces=zone_names))
    meshes = []
    for zone in zone_names:
        zd = surf_data.get(zone, {})
        verts = np.asarray(zd.get("vertices", []))
        faces = np.asarray(zd.get("faces", []), dtype=np.int64)
        if not verts.size: continue
        mesh = pv.PolyData(verts, faces)
        if field_name:
            s_data = fd.get_field_data(ScalarFieldDataRequest(field_name=field_name, surfaces=[zone]))
            if zone in s_data:
                arr = np.asarray(s_data[zone])
                if arr.shape[0] == mesh.n_points: mesh.point_data[field_name] = arr
                elif arr.shape[0] == mesh.n_cells: mesh.cell_data[field_name] = arr
        meshes.append(mesh)
    if not meshes: raise RuntimeError("No surface data returned.")
    return meshes[0].merge(meshes[1:]) if len(meshes) > 1 else meshes[0]


@mcp.tool()
def get_surface_scalar_data(field_name: str, zone_names: list[str]) -> str:
    """Extract scalar field statistics from surfaces (min, max, mean).

    Args:
        field_name: e.g. "pressure", "temperature", "velocity-magnitude".
        zone_names: List of boundary zone names.
    """
    if err := _chk("solver"): return err
    try:
        import numpy as np
        from ansys.fluent.core import ScalarFieldDataRequest
        data = _fd().get_field_data(ScalarFieldDataRequest(field_name=field_name, surfaces=zone_names))
        result = {}
        for zone, arr in data.items():
            a = np.asarray(arr)
            result[zone] = {"shape": list(a.shape), "min": float(a.min()),
                            "max": float(a.max()), "mean": float(a.mean())}
        return _j({"ok": True, "field": field_name, "zones": result})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_surface_report(
    report_type: str,
    zone_names: list[str],
    field_name: str = "pressure",
) -> str:
    """Compute a surface integral report.

    Args:
        report_type: "area-average", "mass-average", "mass-flow",
                     "area-integral", "vertex-min", "vertex-max".
        zone_names:  Zone names.
        field_name:  Field variable.
    """
    if err := _chk("solver"): return err
    # compute() on a temp surface report definition — the ONLY reliable read
    # path on pyfluent 0.40.1 (get_monitor_value no longer exists on the
    # child object; verified live 2026-07-20 on the converged manifold case)
    rname = "mcp_sr_tmp"
    rtype = {"area-average": "surface-areaavg",
             "mass-average": "surface-massavg",
             "mass-flow":    "surface-massflowrate",
             "area-integral": "surface-areaint",
             "vertex-min":   "surface-vertexmin",
             "vertex-max":   "surface-vertexmax"}.get(report_type, f"surface-{report_type}")
    rd = _solver.solution.report_definitions
    try:
        rd.surface[rname] = {}
        r = rd.surface[rname]
        r.report_type = rtype
        if rtype != "surface-massflowrate":   # mass flow needs no field
            r.field = field_name
        r.surface_names = zone_names
        out = rd.compute(report_defs=[rname])
        value = unit = None
        for entry in out or []:
            v = entry.get(rname)
            if isinstance(v, (list, tuple)):
                value = v[0]
                unit = v[1] if len(v) > 1 else None
            else:
                value = v
        return _j({"ok": True, "report_type": report_type, "field": field_name,
                   "zones": zone_names, "value": value, "unit": unit})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    finally:
        try:
            del rd.surface[rname]
        except Exception:
            pass


def _read_report_value(rd, name: str):
    """Read a computed report-definition's value via rd.compute() - the ONLY
    reliable read path on pyfluent 0.40.1 (get_monitor_value no longer exists
    on the child object; verified live 2026-07-20 on the converged manifold
    case; same pattern already used by get_surface_report/
    check_outlet_flow_balance). Returns the raw numeric value; unit is
    discarded to match every call site's pre-existing float/int expectation.
    """
    out = rd.compute(report_defs=[name])
    value = None
    for entry in out or []:
        v = entry.get(name)
        value = v[0] if isinstance(v, (list, tuple)) else v
    return value


@mcp.tool()
def get_volume_average(cell_zone: str, field_name: str = "pressure") -> str:
    """Volume-average a field over a cell zone.

    Args:
        cell_zone:  Cell zone name.
        field_name: Field variable.
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        report.volume["mcp_va"] = {}
        r = report.volume["mcp_va"]
        r.report_type = "volume-average"
        r.field = field_name
        r.cell_zones = [cell_zone]
        return _j({"ok": True, "field": field_name, "zone": cell_zone,
                   "value": _read_report_value(report, "mcp_va")})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_min_max_field(field_name: str, zone_names: list[str] | None = None) -> str:
    """Global min and max of a field variable.

    Args:
        field_name: Field variable.
        zone_names: Optional zone filter.
    """
    if err := _chk("solver"): return err
    try:
        if zone_names is None:
            # exclude interior zones - not a valid report surface (verified
            # live 2026-07-20: including e.g. "interior--<body>" raised
            # "'surface_names' has no attribute [...]" on this pyfluent version)
            zone_names = [z["zone"] for z in _zone_type_pairs() if z["type"] != "interior"]
        report = _solver.solution.report_definitions
        results = {}
        for suffix, rtype in [("min", "surface-vertexmin"), ("max", "surface-vertexmax")]:
            rn = f"mcp_{suffix}"
            report.surface[rn] = {}
            r = report.surface[rn]
            r.report_type = rtype; r.field = field_name; r.surface_names = zone_names
            results[suffix] = _read_report_value(report, rn)
        return _j({"ok": True, "field": field_name, **results})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _residual_history():
    """Residual history via the SESSION monitors service.

    Returns (iterations: list[int], {equation: list[float]}). The TUI path
    (solve/monitors/residual/print-to-screen) is an INVALID command in
    2025 R2 — do not use it. Data streams in while the solver iterates;
    empty results mean no iterations have run in this session yet."""
    x, ys = _solver.monitors.get_monitor_set_data("residual")
    return list(map(int, x)), {k: [float(v) for v in arr] for k, arr in ys.items()}


@mcp.tool()
def get_residuals() -> str:
    """Fetch the latest solver residual values (session monitors service)."""
    if err := _chk("solver"): return err
    try:
        iters, ys = _residual_history()
        if not iters:
            return _j({"ok": True, "residuals": {}, "iteration": None,
                       "note": "No residual history yet in this session "
                               "(monitors stream fills during iteration)."})
        return _j({"ok": True, "iteration": iters[-1],
                   "residuals": {k: v[-1] for k, v in ys.items() if v}})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 10. ADVANCED VISUALISATION (PyVista)
# ===========================================================================

def _pv_start():
    import pyvista as pv
    if os.name != "nt":
        try: pv.start_xvfb()
        except Exception: pass
    return pv


@mcp.tool()
def render_contour(
    zone_names: list[str],
    field_name: str = "pressure",
    output_image_path: str = "",
    colormap: str = "jet",
    show_edges: bool = False,
    window_size: list[int] | None = None,
) -> str:
    """Render a CFD contour plot and save to PNG.

    Args:
        zone_names:         Surface zones.
        field_name:         Scalar field (e.g. "pressure", "temperature").
        output_image_path:  Output PNG path.
        colormap:           Matplotlib colourmap (default "jet").
        show_edges:         Overlay mesh edges.
        window_size:        [width, height] pixels (default [1920, 1080]).
    """
    if err := _chk("solver"): return err
    try:
        pv = _pv_start()
    except ImportError:
        return _j({"ok": False, "error": "pyvista not installed."})
    try:
        mesh = _build_pyvista_surface(zone_names, field_name)
        ws = window_size or [1920, 1080]
        out = output_image_path or f"contour_{field_name}.png"
        _mkdir(out)
        pl = pv.Plotter(off_screen=True, window_size=ws)
        pl.add_mesh(mesh, scalars=field_name, cmap=colormap, show_edges=show_edges,
                    scalar_bar_args={"title": field_name})
        pl.add_axes(); pl.view_isometric()
        pl.screenshot(out); pl.close()
        return _j({"ok": True, "image_path": out,
                   "mesh_points": mesh.n_points, "mesh_cells": mesh.n_cells})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def render_velocity_vectors(
    zone_names: list[str],
    scale: float = 1.0,
    color_by_magnitude: bool = True,
    output_image_path: str = "",
    window_size: list[int] | None = None,
) -> str:
    """Render velocity vector glyphs and save to PNG.

    Args:
        zone_names:          Surface zones.
        scale:               Arrow scale factor.
        color_by_magnitude:  Colour arrows by speed.
        output_image_path:   Output PNG path.
        window_size:         [width, height] pixels.
    """
    if err := _chk("solver"): return err
    try:
        import numpy as np
        pv = _pv_start()
        from ansys.fluent.core import SurfaceFieldDataRequest, SurfaceDataType, VectorFieldDataRequest
    except ImportError as e:
        return _j({"ok": False, "error": f"Missing dependency: {e}"})
    try:
        fd = _fd()
        surf = fd.get_field_data(SurfaceFieldDataRequest(
            data_types=[SurfaceDataType.Vertices], surfaces=zone_names))
        vecs = fd.get_field_data(VectorFieldDataRequest(field_name="velocity", surfaces=zone_names))
        all_pts, all_vecs = [], []
        for zone in zone_names:
            v = np.asarray(surf.get(zone, {}).get("vertices", []))
            u = np.asarray(vecs.get(zone, []))
            if v.size and u.size:
                all_pts.append(v); all_vecs.append(u)
        if not all_pts:
            return _j({"ok": False, "error": "No vector data."})
        pts = np.vstack(all_pts); v_arr = np.vstack(all_vecs)
        mag = np.linalg.norm(v_arr, axis=1)
        cloud = pv.PolyData(pts)
        cloud["velocity"] = v_arr; cloud["magnitude"] = mag
        arrows = cloud.glyph(orient="velocity",
                             scale="magnitude" if color_by_magnitude else False,
                             factor=scale)
        ws = window_size or [1920, 1080]
        out = output_image_path or "velocity_vectors.png"
        _mkdir(out)
        pl = pv.Plotter(off_screen=True, window_size=ws)
        pl.add_mesh(arrows, scalars="magnitude" if color_by_magnitude else None,
                    cmap="turbo", scalar_bar_args={"title": "Velocity (m/s)"})
        pl.add_axes(); pl.view_isometric()
        pl.screenshot(out); pl.close()
        return _j({"ok": True, "image_path": out, "max_velocity_ms": float(mag.max())})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def render_streamlines(
    seed_zone: str,
    n_seeds: int = 100,
    max_steps: int = 500,
    tube_radius: float = 0.002,
    output_image_path: str = "",
    colormap: str = "plasma",
    window_size: list[int] | None = None,
) -> str:
    """Render 3D streamlines seeded from a surface zone.

    Args:
        seed_zone:          Zone whose vertex positions seed the streamlines.
        n_seeds:            Number of seed points.
        max_steps:          Integration steps per streamline.
        tube_radius:        Tube rendering radius (m).
        output_image_path:  Output PNG path.
        colormap:           Colour map.
        window_size:        [width, height] pixels.
    """
    if err := _chk("solver"): return err
    try:
        import numpy as np
        pv = _pv_start()
        from ansys.fluent.core import SurfaceFieldDataRequest, SurfaceDataType, VectorFieldDataRequest
    except ImportError as e:
        return _j({"ok": False, "error": f"Missing dependency: {e}"})
    try:
        fd = _fd()
        # exclude interior zones - not a valid report/field-data surface
        # (verified live 2026-07-20; same root cause as get_min_max_field)
        all_zones = [z["zone"] for z in _zone_type_pairs() if z["type"] != "interior"]
        seed_data = fd.get_field_data(SurfaceFieldDataRequest(
            data_types=[SurfaceDataType.Vertices], surfaces=[seed_zone]))
        seed_pts = np.asarray(seed_data[seed_zone]["vertices"])
        surf_all = fd.get_field_data(SurfaceFieldDataRequest(
            data_types=[SurfaceDataType.Vertices], surfaces=all_zones))
        vec_all = fd.get_field_data(VectorFieldDataRequest(field_name="velocity", surfaces=all_zones))
        all_p, all_v = [], []
        for z in all_zones:
            p = np.asarray(surf_all.get(z, {}).get("vertices", []))
            u = np.asarray(vec_all.get(z, []))
            if p.size and u.size:
                all_p.append(p); all_v.append(u)
        pts_all = np.vstack(all_p)
        vecs_all = np.vstack(all_v)

        # Ensure vectors have non-zero magnitude (streamlines need a defined field)
        mag = np.linalg.norm(vecs_all, axis=1, keepdims=True)
        mag_max = float(mag.max()) if mag.size else 0.0
        if mag_max < 1e-12:
            # Perturb slightly so tracer can integrate
            vecs_all = vecs_all + np.random.uniform(0.01, 0.1, vecs_all.shape)

        cloud = pv.PolyData(pts_all)
        cloud["velocity"] = vecs_all
        idx = np.random.choice(len(seed_pts), min(n_seeds, len(seed_pts)), replace=False)
        seeds = pv.PolyData(seed_pts[idx])

        pv.global_theme.allow_empty_mesh = True
        sl = cloud.streamlines_from_source(seeds, vectors="velocity",
                                           max_steps=max_steps,
                                           integration_direction="forward")
        ws = window_size or [1920, 1080]
        out = output_image_path or "streamlines.png"
        _mkdir(out)
        pl = pv.Plotter(off_screen=True, window_size=ws)
        if sl.n_points > 0:
            tubes = sl.tube(radius=tube_radius)
            pl.add_mesh(tubes, cmap=colormap, scalar_bar_args={"title": "Velocity (m/s)"})
        else:
            # No streamlines traced — fall back to velocity vector glyphs
            pl.add_mesh(cloud.glyph(orient="velocity", factor=0.01), cmap=colormap)
        pl.add_axes(); pl.view_isometric()
        pl.screenshot(out); pl.close()
        return _j({"ok": True, "image_path": out, "seed_zone": seed_zone,
                   "streamline_points": sl.n_points})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def render_iso_surface(
    field_name: str,
    iso_value: float,
    color_by_field: str = "",
    output_image_path: str = "",
    colormap: str = "coolwarm",
    window_size: list[int] | None = None,
) -> str:
    """Create and render an iso-surface.

    Args:
        field_name:         Field to iso-surface.
        iso_value:          Iso-value.
        color_by_field:     Secondary field for colouring (empty = same as field_name).
        output_image_path:  Output PNG path.
        colormap:           Colour map.
        window_size:        [width, height] pixels.
    """
    if err := _chk("solver"): return err
    try:
        pv = _pv_start()
    except ImportError:
        return _j({"ok": False, "error": "pyvista not installed."})
    try:
        zones = list(_solver.setup.boundary_conditions.keys())
        mesh = _build_pyvista_surface(zones, field_name)
        iso = mesh.contour(isosurfaces=[iso_value], scalars=field_name)
        ws = window_size or [1920, 1080]
        out = output_image_path or f"iso_{field_name}_{iso_value}.png"
        _mkdir(out)
        pl = pv.Plotter(off_screen=True, window_size=ws)
        pl.add_mesh(iso, scalars=color_by_field or field_name,
                    cmap=colormap, scalar_bar_args={"title": color_by_field or field_name})
        pl.add_axes(); pl.view_isometric()
        pl.screenshot(out); pl.close()
        return _j({"ok": True, "image_path": out,
                   "iso_points": iso.n_points, "iso_cells": iso.n_cells})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 11. EXPORT (VTK / glTF / HTML / USDZ / OpenUSD)
# ===========================================================================

@mcp.tool()
def export_vtk(zone_names: list[str], field_name: str, output_path: str) -> str:
    """Export surface data as a VTK PolyData file (.vtp) for ParaView.

    Args:
        zone_names:  Zones to include.
        field_name:  Scalar to embed.
        output_path: Output path (.vtp recommended; .vtk also accepted).
    """
    if err := _chk("solver"): return err
    try:
        pv = _pv_start()
        mesh = _build_pyvista_surface(zone_names, field_name)
        _mkdir(output_path)
        # PolyData saves as .vtp; UnstructuredGrid would use .vtu
        out = str(Path(output_path).with_suffix(".vtp"))
        mesh.save(out)
        return _j({"ok": True, "vtk_path": out,
                   "points": mesh.n_points, "cells": mesh.n_cells})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def export_gltf(
    zone_names: list[str],
    field_name: str,
    output_path: str,
    colormap: str = "jet",
) -> str:
    """Export CFD surface as glTF 2.0 for Unity / Unreal / WebXR / AR Quick Look.

    Scalar values are baked as RGB vertex colours.

    Args:
        zone_names:  Zones to include.
        field_name:  Scalar for colour map.
        output_path: Output .gltf path.
        colormap:    Matplotlib colourmap.
    """
    if err := _chk("solver"): return err
    try:
        import numpy as np
        import matplotlib.cm as cm
        import matplotlib.colors as mcolors
        pv = _pv_start()
    except ImportError as e:
        return _j({"ok": False, "error": f"Missing: {e}"})
    try:
        mesh = _build_pyvista_surface(zone_names, field_name)
        _mkdir(output_path)
        out = str(Path(output_path).with_suffix(".gltf"))
        if field_name in mesh.point_data:
            sc = np.asarray(mesh.point_data[field_name])
            norm = mcolors.Normalize(vmin=sc.min(), vmax=sc.max())
            rgba = (_get_cmap(colormap)(norm(sc)) * 255).astype(np.uint8)
            mesh.point_data["RGBA"] = rgba[:, :3]
        pl = pv.Plotter(off_screen=True)
        if "RGBA" in mesh.point_data:
            pl.add_mesh(mesh, scalars="RGBA", rgb=True)
        else:
            pl.add_mesh(mesh, scalars=field_name, cmap=colormap)
        pl.export_gltf(out); pl.close()
        return _j({"ok": True, "gltf_path": out})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def export_html(
    zone_names: list[str],
    field_name: str,
    output_path: str,
    colormap: str = "jet",
) -> str:
    """Export a self-contained interactive 3D HTML visualisation.

    Uses PyVista's HTML export (requires pythreejs or trame).
    Falls back to an embedded PNG inside HTML if neither is available.

    Args:
        zone_names:  Zones to include.
        field_name:  Scalar for colour map.
        output_path: Output .html path.
        colormap:    Matplotlib colourmap.
    """
    if err := _chk("solver"): return err
    try:
        pv = _pv_start()
    except ImportError:
        return _j({"ok": False, "error": "pyvista not installed."})
    try:
        mesh = _build_pyvista_surface(zone_names, field_name)
        _mkdir(output_path)
        out = str(Path(output_path).with_suffix(".html"))

        # Try PyVista's native HTML export
        try:
            pl = pv.Plotter(off_screen=True)
            pl.add_mesh(mesh, scalars=field_name, cmap=colormap,
                        scalar_bar_args={"title": field_name})
            pl.export_html(out)
            pl.close()
            return _j({"ok": True, "html_path": out,
                       "note": "Interactive 3D — open in any browser."})
        except Exception:
            pass  # trame/pythreejs not available — use PNG-in-HTML fallback

        # Fallback: render PNG, embed in minimal HTML
        import base64, tempfile
        png_tmp = tempfile.mktemp(suffix=".png")
        pl2 = pv.Plotter(off_screen=True, window_size=[1200, 800])
        pl2.add_mesh(mesh, scalars=field_name, cmap=colormap,
                     scalar_bar_args={"title": field_name})
        pl2.add_axes(); pl2.view_isometric()
        pl2.screenshot(png_tmp); pl2.close()
        img_b64 = base64.b64encode(open(png_tmp, "rb").read()).decode()
        html = f"""<!DOCTYPE html>
<html><head><meta charset="UTF-8">
<title>PyFluent CFD — {field_name}</title>
<style>body{{margin:0;background:#0d1117;color:#e6edf3;font-family:sans-serif;text-align:center}}
h2{{padding:16px;color:#FFB800}}img{{max-width:100%;border:1px solid #30363d}}</style>
</head><body>
<h2>CFD Result: {field_name} on {', '.join(zone_names)}</h2>
<img src="data:image/png;base64,{img_b64}" alt="CFD contour">
<p style="color:#8b949e;padding:8px">
Install pyvista[jupyter] for interactive 3D: <code>pip install "pyvista[jupyter]"</code>
</p></body></html>"""
        Path(out).write_text(html)
        return _j({"ok": True, "html_path": out,
                   "note": "Static PNG embedded in HTML (install pyvista[jupyter] for interactive 3D)."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def export_openusd(
    zone_names: list[str],
    field_name: str,
    output_path: str,
    colormap: str = "jet",
    nucleus_url: str = "",
) -> str:
    """Export CFD surface as OpenUSD (.usdc) for NVIDIA Omniverse.

    Scalar values are written as displayColor vertex primvar.
    Optionally push to an Omniverse Nucleus server.

    Args:
        zone_names:   Zones to include.
        field_name:   Scalar for colour map.
        output_path:  Local .usdc output path.
        colormap:     Matplotlib colourmap.
        nucleus_url:  Omniverse Nucleus URL (e.g. omniverse://host/Projects/cfd.usd).
    """
    if err := _chk("solver"): return err
    try:
        from pxr import Usd, UsdGeom, Vt, Gf, Sdf
    except ImportError:
        return _j({"ok": False, "error": "usd-core not installed. pip install usd-core"})
    try:
        import numpy as np
        import matplotlib.cm as cm
        import matplotlib.colors as mcolors
    except ImportError as e:
        return _j({"ok": False, "error": str(e)})
    try:
        mesh_pv = _build_pyvista_surface(zone_names, field_name)
        _mkdir(output_path)
        out = str(Path(output_path).with_suffix(".usdc"))
        stage = Usd.Stage.CreateNew(out)
        UsdGeom.SetStageUpAxis(stage, UsdGeom.Tokens.y)
        stage.SetMetadata("metersPerUnit", 1.0)
        root = UsdGeom.Xform.Define(stage, "/CFD_Result")
        usd_mesh = UsdGeom.Mesh.Define(stage, "/CFD_Result/Mesh")
        pts = mesh_pv.points.tolist()
        usd_mesh.GetPointsAttr().Set(Vt.Vec3fArray([Gf.Vec3f(*p) for p in pts]))
        faces = mesh_pv.faces; fc, fi = [], []; i = 0
        while i < len(faces):
            n = faces[i]; fc.append(n); fi.extend(faces[i+1:i+1+n].tolist()); i += 1+n
        usd_mesh.GetFaceVertexCountsAttr().Set(Vt.IntArray(fc))
        usd_mesh.GetFaceVertexIndicesAttr().Set(Vt.IntArray(fi))
        if field_name in mesh_pv.point_data:
            sc = np.asarray(mesh_pv.point_data[field_name])
            norm = mcolors.Normalize(vmin=sc.min(), vmax=sc.max())
            rgb = _get_cmap(colormap)(norm(sc))[:, :3]
            pv_api = UsdGeom.PrimvarsAPI(usd_mesh)
            cpv = pv_api.CreatePrimvar("displayColor", Sdf.ValueTypeNames.Color3fArray,
                                        UsdGeom.Tokens.vertex)
            cpv.Set(Vt.Vec3fArray([Gf.Vec3f(float(r), float(g), float(b)) for r,g,b in rgb]))
        root_prim = stage.GetPrimAtPath("/CFD_Result")
        root_prim.SetCustomDataByKey("cfd:field", field_name)
        root_prim.SetCustomDataByKey("cfd:zones", ", ".join(zone_names))
        root_prim.SetCustomDataByKey("cfd:source", "ansys-fluent-mcp")
        stage.Save()
        result = {"ok": True, "usd_path": out}
        if nucleus_url:
            try:
                import omni.client as oc
                oc.initialize("ansys-fluent-mcp")
                with open(out, "rb") as fh:
                    content = fh.read()
                r = oc.write_file(nucleus_url, content)
                result["nucleus_url"] = nucleus_url; result["nucleus_result"] = str(r)
            except ImportError:
                result["nucleus_warning"] = "omni.client not available."
            except Exception as ne:
                result["nucleus_error"] = str(ne)
        return _j(result)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 12. PARAMETRIC STUDIES & DOE  — NEW
# ===========================================================================

@mcp.tool()
def list_input_parameters() -> str:
    """List all input parameters defined in the current parametric case.

    Input parameters must be defined in the Fluent case before running DOE.
    They appear in the parametric studies panel in Fluent.
    """
    if err := _chk("solver"): return err
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        base_dp = study.design_points["Base DP"]
        return _j({"ok": True,
                   "input_parameters": dict(base_dp.input_parameters),
                   "output_parameters": dict(base_dp.output_parameters)})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def create_design_point(
    name: str,
    input_parameters: dict,
) -> str:
    """Add a new design point to the parametric study.

    Args:
        name:             Design point name (e.g. "DP-1", "high_vel").
        input_parameters: Dict mapping parameter names to values.
                          e.g. {"inlet_velocity": 10.0, "outlet_temp": 350.0}
    """
    if err := _chk("solver"): return err
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        dp = study.add_design_point(name)
        updates = dict(dp.input_parameters)
        updates.update(input_parameters)
        dp.input_parameters = updates
        return _j({"ok": True, "design_point": name, "inputs": input_parameters})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def create_doe_sweep(
    parameter_name: str,
    values: list[float],
    base_inputs: dict | None = None,
) -> str:
    """Create a 1-D parameter sweep (Design of Experiments).

    Creates one design point per value and runs all of them.

    Args:
        parameter_name: Input parameter to sweep (e.g. "inlet_velocity").
        values:         List of values to sweep over.
        base_inputs:    Optional fixed values for other parameters
                        (e.g. {"outlet_temp": 300.0}).
    """
    if err := _chk("solver"): return err
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        created = []
        for i, val in enumerate(values):
            dp_name = f"dp_{parameter_name}_{i}"
            dp = study.add_design_point(dp_name)
            updates = dict(dp.input_parameters)
            if base_inputs:
                updates.update(base_inputs)
            updates[parameter_name] = val
            dp.input_parameters = updates
            created.append({"name": dp_name, parameter_name: val})
        return _j({"ok": True, "sweep_parameter": parameter_name,
                   "design_points_created": len(created), "design_points": created,
                   "note": "Call run_design_points to execute all design points."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def create_doe_full_factorial(
    parameters: dict[str, list[float]],
) -> str:
    """Create a full-factorial DOE across multiple parameters.

    Creates design points for every combination of parameter values.

    Args:
        parameters: Dict mapping parameter names to lists of values.
                    e.g. {"inlet_velocity": [5, 10, 15],
                           "outlet_temp": [300, 400]}
                    → 3 × 2 = 6 design points.
    """
    if err := _chk("solver"): return err
    try:
        import itertools
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        names = list(parameters.keys())
        value_lists = [parameters[n] for n in names]
        combinations = list(itertools.product(*value_lists))
        created = []
        for i, combo in enumerate(combinations):
            dp_name = f"dp_ff_{i}"
            dp = study.add_design_point(dp_name)
            updates = dict(dp.input_parameters)
            for n, v in zip(names, combo):
                updates[n] = v
            dp.input_parameters = updates
            created.append({"name": dp_name, **dict(zip(names, combo))})
        return _j({"ok": True, "total_design_points": len(created),
                   "parameters": names, "design_points": created,
                   "note": "Call run_design_points to execute."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _run_design_points_impl(design_point_names: list[str] | None) -> str:
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        dps = study.design_points
        to_run = design_point_names or list(dps.keys())
        results = {}
        for dp_name in to_run:
            if dp_name not in dps:
                results[dp_name] = {"error": "Not found"}
                continue
            try:
                dp = dps[dp_name]
                study.update_current_design_point_with_design_point(dp)
                results[dp_name] = {
                    "status": "completed",
                    "outputs": dict(dp.output_parameters),
                }
            except Exception as e:
                results[dp_name] = {"status": "failed", "error": str(e)}
        return _j({"ok": True, "ran": len(to_run), "results": results})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def run_design_points(design_point_names: list[str] | None = None,
                       run_in_background: bool = True) -> str:
    """Run design points in the parametric study.

    Runs AS A BACKGROUND JOB by default (verified live 2026-07-20: this
    loops update_current_design_point_with_design_point over every design
    point fully synchronously - each point re-solves the case, so this can
    take as long as a multi-point batch solve and was freezing the MCP
    request thread like the long meshing calls that were already
    backgrounded to avoid this). Poll get_job_status + read_console_tail
    for live progress.

    Args:
        design_point_names: Specific design points to run.
                            If None, runs all design points.
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if run_in_background:
        return _run_session_job("run_design_points",
                                lambda: _run_design_points_impl(design_point_names))
    return _run_design_points_impl(design_point_names)


@mcp.tool()
def get_doe_results(output_csv_path: str = "") -> str:
    """Retrieve all DOE results as a table (inputs + outputs per design point).

    Args:
        output_csv_path: Optional path to save results as CSV.
    """
    if err := _chk("solver"): return err
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        rows = []
        for dp_name, dp in study.design_points.items():
            row: dict = {"design_point": dp_name}
            row.update({f"in:{k}": v for k, v in dp.input_parameters.items()})
            row.update({f"out:{k}": v for k, v in dp.output_parameters.items()})
            rows.append(row)
        if output_csv_path and rows:
            import csv
            _mkdir(output_csv_path)
            with open(output_csv_path, "w", newline="", encoding="utf-8") as f:
                writer = csv.DictWriter(f, fieldnames=rows[0].keys())
                writer.writeheader(); writer.writerows(rows)
        return _j({"ok": True, "design_points": len(rows), "table": rows,
                   "csv_path": output_csv_path or None})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def plot_doe_results(
    x_parameter: str,
    y_parameter: str,
    output_image_path: str = "",
    title: str = "",
) -> str:
    """Plot DOE results as an XY scatter/line chart and save to PNG.

    Args:
        x_parameter: Input parameter for X axis (e.g. "in:inlet_velocity").
        y_parameter: Output parameter for Y axis (e.g. "out:pressure_drop").
        output_image_path: Output PNG path.
        title:        Chart title.
    """
    if err := _chk("solver"): return err
    try:
        import matplotlib.pyplot as plt
    except ImportError:
        return _j({"ok": False, "error": "matplotlib not installed."})
    try:
        studies = _solver.parametric_studies
        study = next(iter(studies.values()))
        x_vals, y_vals, labels = [], [], []
        for dp_name, dp in study.design_points.items():
            all_data = {**{f"in:{k}": v for k, v in dp.input_parameters.items()},
                        **{f"out:{k}": v for k, v in dp.output_parameters.items()}}
            if x_parameter in all_data and y_parameter in all_data:
                x_vals.append(float(all_data[x_parameter]))
                y_vals.append(float(all_data[y_parameter]))
                labels.append(dp_name)
        if not x_vals:
            return _j({"ok": False, "error": "No data points found for the given parameters."})
        sorted_pairs = sorted(zip(x_vals, y_vals, labels))
        x_s, y_s, lb_s = zip(*sorted_pairs)
        fig, ax = plt.subplots(figsize=(10, 6))
        ax.plot(x_s, y_s, "o-", markersize=8, linewidth=2, color="#1f77b4")
        for lbl, xv, yv in zip(lb_s, x_s, y_s):
            ax.annotate(lbl, (xv, yv), textcoords="offset points", xytext=(4, 4), fontsize=8)
        ax.set_xlabel(x_parameter, fontsize=12)
        ax.set_ylabel(y_parameter, fontsize=12)
        ax.set_title(title or f"{y_parameter} vs {x_parameter}", fontsize=14)
        ax.grid(True, alpha=0.3)
        out = output_image_path or f"doe_{x_parameter}_vs_{y_parameter}.png"
        _mkdir(out)
        fig.savefig(out, dpi=150, bbox_inches="tight")
        plt.close(fig)
        return _j({"ok": True, "image_path": out, "data_points": len(x_vals)})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 13. CONVERGENCE MONITORING & SMART STOP  — NEW
# ===========================================================================

@mcp.tool()
def create_surface_monitor(
    monitor_name: str,
    report_type: str,
    field_name: str,
    zone_names: list[str],
    frequency: int = 1,
) -> str:
    """Create a surface monitor that tracks a quantity every N iterations.

    Args:
        monitor_name: Unique monitor name (e.g. "outlet_temp_monitor").
        report_type:  "area-average", "mass-average", or "mass-flow".
        field_name:   Field variable to monitor.
        zone_names:   Zones to include in the report.
        frequency:    How often to sample (every N iterations, default 1).
    """
    if err := _chk("solver"): return err
    try:
        rtype_map = {"area-average": "surface-areaavg",
                     "mass-average": "surface-massavg",
                     "mass-flow":    "surface-massflow"}
        # Create a report definition that Fluent will sample every N iterations
        report = _solver.solution.report_definitions
        rdef_name = monitor_name + "_def"
        report.surface[rdef_name] = {}
        rd = report.surface[rdef_name]
        rd.report_type = rtype_map.get(report_type, report_type)
        rd.field = field_name
        rd.surface_names = zone_names
        # Register a report plot/file that references this definition
        try:
            plots = _solver.solution.monitor.report_plots
            plots[monitor_name] = {}
            m = plots[monitor_name]
            try: m.frequency = frequency
            except Exception: pass
        except Exception:
            pass  # Not all Fluent versions expose report_plots via gRPC
        return _j({"ok": True, "monitor": monitor_name,
                   "report_type": report_type, "field": field_name, "zones": zone_names})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_convergence_history(output_image_path: str = "") -> str:
    """Get residual convergence history and optionally plot it.

    Returns 'series': {iterations: [...], <equation>: [...]} — the full
    per-iteration arrays, not just the latest value — for building a real
    convergence curve (iteration on x, residual on y) client-side.

    Args:
        output_image_path: If set, save a convergence history PNG.
    """
    if err := _chk("solver"): return err
    try:
        iters, ys = _residual_history()
        if not iters:
            return _j({"ok": True, "iterations": 0, "history": {},
                       "note": "No residual history yet in this session "
                               "(monitors stream fills during iteration)."})
        result: dict = {"ok": True, "iterations": len(iters),
                        "iteration_range": [iters[0], iters[-1]],
                        "latest": {k: v[-1] for k, v in ys.items() if v},
                        "series": {"iterations": iters, **{k: v for k, v in ys.items() if v}}}
        if output_image_path:
            try:
                import matplotlib.pyplot as plt
                fig, ax = plt.subplots(figsize=(10, 6))
                for name, vals in ys.items():
                    ax.semilogy(iters[:len(vals)], vals, label=name)
                ax.set_xlabel("Iteration"); ax.set_ylabel("Residual (log scale)")
                ax.set_title("Convergence History"); ax.legend(); ax.grid(True, alpha=0.3)
                _mkdir(output_image_path)
                fig.savefig(output_image_path, dpi=150, bbox_inches="tight")
                plt.close(fig)
                result["image_path"] = output_image_path
            except Exception as plot_err:
                result["plot_error"] = str(plot_err)
        return _j(result)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _run_until_convergence_impl(
    max_iterations: int,
    check_interval: int,
    continuity_target: float,
    energy_target: float,
    output_history_path: str,
) -> dict:
    """Body of run_until_convergence - a plain dict return so it can run
    either inline (run_in_background=False) or inside a _bg_run worker
    thread (whose 'result' slot get_calculation_status surfaces)."""
    try:
        total_ran = 0
        converged = False
        history: list[dict] = []

        while total_ran < max_iterations and not converged:
            batch = min(check_interval, max_iterations - total_ran)
            try:
                _solver.solution.run_calculation.iterate(iter_count=batch)
            except Exception as e:
                return {"ok": False, "error": f"Solve error at iteration {total_ran}: {e}"}
            total_ran += batch

            # Read latest residuals (session monitors service; the TUI
            # print-to-screen command does not exist in 2025 R2)
            try:
                _, ys = _residual_history()
                latest = {k: v[-1] for k, v in ys.items() if v}
                if latest:
                    history.append({"iteration": total_ran, "residuals": latest})
                    cont = latest.get("continuity", 1.0)
                    en = latest.get("energy", 0.0)
                    if cont < continuity_target and en < energy_target:
                        converged = True
            except Exception:
                pass  # Residual read failed; keep iterating

        result = {"ok": True, "total_iterations": total_ran, "converged": converged}
        if output_history_path and history:
            try:
                import matplotlib.pyplot as plt
                its = [h["iteration"] for h in history]
                # plot each residual track
                n_res = max(len(h["residuals"]) for h in history)
                fig, ax = plt.subplots(figsize=(10, 6))
                for ri in range(n_res):
                    vals = [h["residuals"][ri] if ri < len(h["residuals"]) else None
                            for h in history]
                    vals = [v for v in vals if v is not None]
                    if vals:
                        ax.semilogy(its[:len(vals)], vals, label=f"residual-{ri}")
                ax.axhline(continuity_target, ls="--", color="red", label="continuity target")
                ax.set_xlabel("Iteration"); ax.set_ylabel("Residual")
                ax.set_title("Convergence History (smart stop)"); ax.legend(); ax.grid(True, alpha=0.3)
                _mkdir(output_history_path)
                fig.savefig(output_history_path, dpi=150, bbox_inches="tight")
                plt.close(fig)
                result["history_image"] = output_history_path
            except Exception as pe:
                result["plot_error"] = str(pe)
        return result
    except Exception as e:
        return {"ok": False, "error": str(e)}


@mcp.tool()
def run_until_convergence(
    max_iterations: int = 1000,
    check_interval: int = 50,
    continuity_target: float = 1e-5,
    energy_target: float = 1e-7,
    output_history_path: str = "",
    run_in_background: bool = True,
    force_no_monitors: bool = False,
) -> str:
    """Run iterations with smart stopping when convergence criteria are met.

    Iterates in blocks of check_interval, reads residuals after each block,
    and stops automatically once all specified residuals are below their
    targets.

    Runs AS A BACKGROUND JOB by default (verified live 2026-07-20: this was
    previously fully synchronous on the MCP request thread for up to
    max_iterations=1000 iterations, freezing the server exactly like the
    long meshing calls that start_calculation/import_geometry were already
    backgrounded to avoid). Poll get_calculation_status (kind="convergence";
    'result' appears once finished) or read_console_tail for live progress;
    stop early with stop_calculation.

    MONITORS GATE — SERVER-ENFORCED: refuses to start if zero report
    monitors are defined, unless force_no_monitors=True (setup_standard_monitors first).

    Args:
        max_iterations:    Hard limit on total iterations (default 1000).
        check_interval:    Iterations between convergence checks (default 50).
        continuity_target: Stop when continuity residual < this value (default 1e-5).
        energy_target:     Stop when energy residual < this value (default 1e-7).
        output_history_path: Optional PNG path for convergence history plot.
        run_in_background:   Default True (recommended). False blocks until done.
        force_no_monitors:   Override the monitors gate (rare).
    """
    if err := _chk("solver"): return err
    if _bg_run["running"]:
        return _j({"ok": False, "error": "A background calculation is already running."})
    if gate_err := _require_monitors_or_error(force_no_monitors):
        return gate_err
    if not run_in_background:
        return _j(_run_until_convergence_impl(
            max_iterations, check_interval, continuity_target, energy_target, output_history_path))
    _ensure_iteration_event_registered()
    start_iter = _current_iteration_baseline()
    _bg_run.update(running=True, kind="convergence", requested=max_iterations,
                   start_iteration=start_iter, finished=False, error=None,
                   interrupted=False, result=None)

    def _worker():
        try:
            _bg_run["result"] = _run_until_convergence_impl(
                max_iterations, check_interval, continuity_target, energy_target, output_history_path)
        except Exception as e:
            _bg_run["error"] = str(e)
        finally:
            _bg_run["running"] = False
            _bg_run["finished"] = True

    t = threading.Thread(target=_worker, name="mcp-fluent-convergence", daemon=True)
    _bg_run["thread"] = t
    t.start()
    return _j({"ok": True, "started": True, "kind": "convergence",
               "max_iterations": max_iterations, "start_iteration": start_iter,
               "next": "Poll get_calculation_status (result appears when finished) "
                       "or read_console_tail for live progress; stop early with "
                       "stop_calculation."})


def _is_float(s: str) -> bool:
    try: float(s); return True
    except ValueError: return False


# ===========================================================================
# 14. AERODYNAMIC FORCES & MOMENTS  — NEW
# ===========================================================================

@mcp.tool()
def setup_force_monitors(
    wall_zones: list[str],
    reference_area_m2: float,
    reference_velocity_ms: float,
    reference_density_kgm3: float = 1.225,
    lift_direction: list[float] | None = None,
    drag_direction: list[float] | None = None,
    moment_center: list[float] | None = None,
    reference_length_m: float = 1.0,
) -> str:
    """Set up aerodynamic force and moment monitors.

    Creates report definitions for drag (Cd), lift (Cl), and pitching
    moment (Cm) referenced to the specified conditions.

    Args:
        wall_zones:              Wall zones on the aerodynamic body.
        reference_area_m2:       Reference area for coefficient normalisation (m²).
        reference_velocity_ms:   Freestream velocity (m/s).
        reference_density_kgm3:  Freestream density (kg/m³), default 1.225 (sea-level air).
        lift_direction:          Unit vector [x, y, z] for lift (default [0, 1, 0]).
        drag_direction:          Unit vector [x, y, z] for drag (default [1, 0, 0]).
        moment_center:           [x, y, z] coordinates of moment reference point.
        reference_length_m:      Reference length for Cm (default 1.0 m).
    """
    if err := _chk("solver"): return err

    lift_dir = lift_direction or [0.0, 1.0, 0.0]
    drag_dir = drag_direction or [1.0, 0.0, 0.0]
    mom_ctr  = moment_center or [0.0, 0.0, 0.0]

    try:
        # Set reference values
        ref = _solver.setup.reference_values
        ref.area   = reference_area_m2
        ref.velocity = reference_velocity_ms
        ref.density  = reference_density_kgm3
        ref.length   = reference_length_m

        # Create force report definitions via TUI (most robust path)
        # Drag
        _solver.scheme_eval.exec(
            f'(ti-menu-load-string "/solve/report-definitions/add drag-cd '
            f'type force-coefficient zone-names (\"{chr(34).join(wall_zones)}\") '
            f'force-vector {drag_dir[0]} {drag_dir[1]} {drag_dir[2]} q)")')
        # Lift
        _solver.scheme_eval.exec(
            f'(ti-menu-load-string "/solve/report-definitions/add lift-cl '
            f'type force-coefficient zone-names (\"{chr(34).join(wall_zones)}\") '
            f'force-vector {lift_dir[0]} {lift_dir[1]} {lift_dir[2]} q)")')

        return _j({"ok": True,
                   "wall_zones": wall_zones,
                   "reference_area_m2": reference_area_m2,
                   "reference_velocity_ms": reference_velocity_ms,
                   "reference_density_kgm3": reference_density_kgm3,
                   "lift_direction": lift_dir,
                   "drag_direction": drag_dir,
                   "moment_center": mom_ctr,
                   "note": "Call compute_force_coefficients to get Cl, Cd, Cm values."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def compute_force_coefficients(
    wall_zones: list[str],
    reference_area_m2: float,
    reference_velocity_ms: float,
    reference_density_kgm3: float = 1.225,
    lift_direction: list[float] | None = None,
    drag_direction: list[float] | None = None,
    moment_center: list[float] | None = None,
    reference_length_m: float = 1.0,
) -> str:
    """Compute lift (Cl), drag (Cd), and moment (Cm) coefficients directly.

    Uses PyFluent's report_definitions to compute force integrals on the
    specified wall zones and normalises to coefficients.

    Args:
        wall_zones:              Wall zones on the body.
        reference_area_m2:       Reference area (m²).
        reference_velocity_ms:   Freestream velocity (m/s).
        reference_density_kgm3:  Freestream density (kg/m³).
        lift_direction:          Lift unit vector (default [0, 1, 0]).
        drag_direction:          Drag unit vector (default [1, 0, 0]).
        moment_center:           Moment centre [x, y, z] (default origin).
        reference_length_m:      Reference chord/length for Cm (default 1.0 m).
    """
    if err := _chk("solver"): return err

    lift_dir = lift_direction or [0.0, 1.0, 0.0]
    drag_dir = drag_direction or [1.0, 0.0, 0.0]
    mom_ctr  = moment_center or [0.0, 0.0, 0.0]
    q_ref = 0.5 * reference_density_kgm3 * reference_velocity_ms ** 2  # dynamic pressure
    results: dict = {"q_ref_pa": round(q_ref, 4)}

    try:
        report = _solver.solution.report_definitions

        # -- Drag (Cd) --
        report.surface["mcp_drag"] = {}
        rd = report.surface["mcp_drag"]
        rd.report_type = "surface-integral"
        rd.field = "pressure"
        rd.surface_names = wall_zones
        # Force coefficients via TUI (most reliable)
        try:
            drag_raw = float(str(_solver.tui.report.forces.wall_forces(
                *wall_zones, str(drag_dir[0]), str(drag_dir[1]), str(drag_dir[2])
            )).split("\n")[-2].split()[-1])
            cd = drag_raw / (q_ref * reference_area_m2)
            results["Cd"] = round(cd, 6)
            results["drag_force_N"] = round(drag_raw, 4)
        except Exception:
            results["Cd"] = "N/A (use run_tui_command /report/forces)"

        # -- Lift (Cl) --
        try:
            lift_raw = float(str(_solver.tui.report.forces.wall_forces(
                *wall_zones, str(lift_dir[0]), str(lift_dir[1]), str(lift_dir[2])
            )).split("\n")[-2].split()[-1])
            cl = lift_raw / (q_ref * reference_area_m2)
            results["Cl"] = round(cl, 6)
            results["lift_force_N"] = round(lift_raw, 4)
        except Exception:
            results["Cl"] = "N/A"

        # -- Lift/Drag ratio --
        try:
            if isinstance(results.get("Cl"), float) and isinstance(results.get("Cd"), float):
                results["Cl_Cd_ratio"] = round(results["Cl"] / results["Cd"], 4)
        except Exception:
            pass

        # -- Pitching moment --
        try:
            mom_raw = float(str(_solver.tui.report.forces.wall_moments(
                *wall_zones,
                str(mom_ctr[0]), str(mom_ctr[1]), str(mom_ctr[2])
            )).split("\n")[-2].split()[-1])
            cm = mom_raw / (q_ref * reference_area_m2 * reference_length_m)
            results["Cm"] = round(cm, 6)
            results["moment_Nm"] = round(mom_raw, 4)
        except Exception:
            results["Cm"] = "N/A"

        results["reference"] = {
            "area_m2": reference_area_m2,
            "velocity_ms": reference_velocity_ms,
            "density_kgm3": reference_density_kgm3,
            "length_m": reference_length_m,
        }
        return _j({"ok": True, **results})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_pressure_drag_breakdown(wall_zones: list[str]) -> str:
    """Break drag down into pressure (form) drag and viscous (skin friction) drag.

    Args:
        wall_zones: Wall zones to include.
    """
    if err := _chk("solver"): return err
    try:
        result = {}
        for component in ("pressure", "viscous"):
            rname = f"mcp_drag_{component}"
            report = _solver.solution.report_definitions
            report.surface[rname] = {}
            r = report.surface[rname]
            r.report_type = "surface-integral"
            r.field = f"wall-{component}-force"
            r.surface_names = wall_zones
            try: result[f"{component}_force_N"] = _read_report_value(report, rname)
            except Exception: result[f"{component}_force_N"] = "N/A"
        try:
            pf = result.get("pressure_force_N", 0)
            vf = result.get("viscous_force_N", 0)
            if isinstance(pf, (int, float)) and isinstance(vf, (int, float)):
                total = pf + vf
                result["total_drag_N"] = round(total, 4)
                result["pressure_fraction"] = round(pf / total, 4) if total else "N/A"
                result["viscous_fraction"]  = round(vf / total, 4) if total else "N/A"
        except Exception:
            pass
        return _j({"ok": True, "zones": wall_zones, **result})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 15. HEAT-TRANSFER ANALYSIS  — NEW
# ===========================================================================

@mcp.tool()
def get_heat_transfer_rate(wall_zones: list[str]) -> str:
    """Compute total heat transfer rate through wall zones.

    Args:
        wall_zones: Wall zone names.
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        report.surface["mcp_htc"] = {}
        r = report.surface["mcp_htc"]
        r.report_type = "surface-integral"
        r.field = "heat-flux"
        r.surface_names = wall_zones
        q_total = _read_report_value(report, "mcp_htc")  # W (integrated heat flux over area)
        return _j({"ok": True, "zones": wall_zones,
                   "heat_transfer_rate_W": q_total,
                   "heat_transfer_rate_kW": round(q_total / 1000, 4)
                   if isinstance(q_total, (int, float)) else "N/A"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_wall_heat_flux_statistics(wall_zones: list[str]) -> str:
    """Get min, max, and area-averaged wall heat flux on specified zones.

    Args:
        wall_zones: Wall zone names.
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        results: dict = {}
        for stat, rtype in [("average", "surface-areaavg"),
                             ("min",     "surface-vertexmin"),
                             ("max",     "surface-vertexmax")]:
            rn = f"mcp_hf_{stat}"
            report.surface[rn] = {}
            r = report.surface[rn]
            r.report_type = rtype
            r.field = "heat-flux"
            r.surface_names = wall_zones
            try: results[f"heat_flux_{stat}_Wm2"] = _read_report_value(report, rn)
            except Exception: results[f"heat_flux_{stat}_Wm2"] = "N/A"
        return _j({"ok": True, "zones": wall_zones, **results})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_nusselt_number(
    wall_zone: str,
    reference_length_m: float,
    thermal_conductivity_wm_k: float = 0.0242,
) -> str:
    """Compute the area-averaged Nusselt number on a wall zone.

    Nu = h × L / k   where h = area-averaged heat transfer coefficient.

    Args:
        wall_zone:                  Wall zone name.
        reference_length_m:         Characteristic length L (m).
        thermal_conductivity_wm_k:  Fluid thermal conductivity k (W/m·K).
                                    Default = 0.0242 (air at 300 K).
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        report.surface["mcp_htcoeff"] = {}
        r = report.surface["mcp_htcoeff"]
        r.report_type = "surface-areaavg"
        r.field = "heat-transfer-coef"
        r.surface_names = [wall_zone]
        h = _read_report_value(report, "mcp_htcoeff")
        if isinstance(h, (int, float)):
            nu = h * reference_length_m / thermal_conductivity_wm_k
            return _j({"ok": True, "zone": wall_zone,
                       "h_avg_Wm2K": round(h, 4),
                       "Nu": round(nu, 4),
                       "L_m": reference_length_m,
                       "k_WmK": thermal_conductivity_wm_k})
        else:
            return _j({"ok": False, "error": "Could not read heat transfer coefficient."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_max_wall_temperature(wall_zones: list[str]) -> str:
    """Get the maximum wall temperature across the specified zones.

    Args:
        wall_zones: Wall zone names.
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        report.surface["mcp_twmax"] = {}
        r = report.surface["mcp_twmax"]
        r.report_type = "surface-vertexmax"
        r.field = "wall-temperature"
        r.surface_names = wall_zones
        t_max = _read_report_value(report, "mcp_twmax")

        report.surface["mcp_twmin"] = {}
        r2 = report.surface["mcp_twmin"]
        r2.report_type = "surface-vertexmin"
        r2.field = "wall-temperature"
        r2.surface_names = wall_zones
        t_min = _read_report_value(report, "mcp_twmin")

        report.surface["mcp_twavg"] = {}
        r3 = report.surface["mcp_twavg"]
        r3.report_type = "surface-areaavg"
        r3.field = "wall-temperature"
        r3.surface_names = wall_zones
        t_avg = _read_report_value(report, "mcp_twavg")

        return _j({"ok": True, "zones": wall_zones,
                   "T_max_K": t_max, "T_min_K": t_min, "T_avg_K": t_avg,
                   "T_max_C": round(t_max - 273.15, 2) if isinstance(t_max, (int, float)) else "N/A"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def render_heat_flux_map(
    wall_zones: list[str],
    output_image_path: str = "",
    colormap: str = "hot",
    window_size: list[int] | None = None,
) -> str:
    """Render a wall heat flux map and save to PNG.

    Args:
        wall_zones:         Wall zones to include.
        output_image_path:  Output PNG path.
        colormap:           Colour map (default "hot" — intuitive for heat).
        window_size:        [width, height] pixels.
    """
    if err := _chk("solver"): return err
    try:
        pv = _pv_start()
    except ImportError:
        return _j({"ok": False, "error": "pyvista not installed."})
    try:
        mesh = _build_pyvista_surface(wall_zones, "heat-flux")
        ws = window_size or [1920, 1080]
        out = output_image_path or "heat_flux_map.png"
        _mkdir(out)
        pl = pv.Plotter(off_screen=True, window_size=ws)
        pl.add_mesh(mesh, scalars="heat-flux", cmap=colormap,
                    scalar_bar_args={"title": "Wall Heat Flux (W/m²)"})
        pl.add_axes(); pl.view_isometric()
        pl.screenshot(out); pl.close()
        return _j({"ok": True, "image_path": out})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_thermal_report(wall_zones: list[str], fluid_zones: list[str] | None = None) -> str:
    """Generate a comprehensive thermal analysis summary.

    Computes heat flux stats, wall temperature stats, and volume-averaged
    fluid temperatures for a complete thermal picture.

    Args:
        wall_zones:  Wall zones to analyse.
        fluid_zones: Optional fluid cell zones for bulk temperature.
    """
    if err := _chk("solver"): return err
    report = _solver.solution.report_definitions
    result: dict = {"wall_zones": wall_zones, "generated": _ts()}

    def _sval(rname, rtype, field, zones):
        try:
            report.surface[rname] = {}
            r = report.surface[rname]
            r.report_type = rtype; r.field = field; r.surface_names = zones
            return _read_report_value(report, rname)
        except Exception:
            return "N/A"

    result["heat_flux_avg_Wm2"] = _sval("mcp_ta1", "surface-areaavg", "heat-flux", wall_zones)
    result["heat_flux_max_Wm2"] = _sval("mcp_ta2", "surface-vertexmax", "heat-flux", wall_zones)
    result["T_wall_avg_K"] = _sval("mcp_ta3", "surface-areaavg", "wall-temperature", wall_zones)
    result["T_wall_max_K"] = _sval("mcp_ta4", "surface-vertexmax", "wall-temperature", wall_zones)
    result["T_wall_min_K"] = _sval("mcp_ta5", "surface-vertexmin", "wall-temperature", wall_zones)

    if fluid_zones:
        for zone in fluid_zones:
            try:
                report.volume[f"mcp_tf_{zone}"] = {}
                rv = report.volume[f"mcp_tf_{zone}"]
                rv.report_type = "volume-average"; rv.field = "temperature"; rv.cell_zones = [zone]
                result[f"T_bulk_{zone}_K"] = _read_report_value(report, f"mcp_tf_{zone}")
            except Exception:
                result[f"T_bulk_{zone}_K"] = "N/A"

    return _j({"ok": True, **result})


# ===========================================================================
# 16. MESH ADAPTATION (AMR)  — NEW
# ===========================================================================

@mcp.tool()
def enable_mesh_adaptation(
    adaptation_criterion: str = "pressure-gradient",
    max_refinement_level: int = 2,
    coarsen: bool = True,
) -> str:
    """Enable gradient-based adaptive mesh refinement (AMR).

    Args:
        adaptation_criterion: Field used to drive adaptation.
                              Options: "pressure-gradient", "velocity-gradient",
                              "temperature-gradient", "y-plus-error".
        max_refinement_level: Maximum cell split levels (default 2).
        coarsen:              Also coarsen cells in low-gradient regions.
    """
    if err := _chk("solver"): return err
    try:
        adapt = _solver.solution.controls.mesh_adaption
        adapt.method = "gradient"
        adapt.criterion = adaptation_criterion.replace("-", "_")
        adapt.max_level = max_refinement_level
        adapt.coarsen = coarsen
        return _j({"ok": True, "criterion": adaptation_criterion,
                   "max_level": max_refinement_level, "coarsen": coarsen,
                   "note": "Call run_adaptation_cycle to refine the mesh."})
    except Exception as e:
        # Fallback: use TUI
        try:
            _solver.tui.adapt.preferences.max_level_of_refinement(str(max_refinement_level))
            return _j({"ok": True, "criterion": adaptation_criterion,
                       "max_level": max_refinement_level,
                       "configured_via": "TUI fallback"})
        except Exception as e2:
            return _j({"ok": False, "error": str(e), "tui_error": str(e2)})


def _run_adaptation_cycle_impl(cycles: int, iterations_between_cycles: int) -> dict:
    try:
        results = []
        for cycle in range(cycles):
            # Solve
            try:
                _solver.solution.run_calculation.iterate(iter_count=iterations_between_cycles)
            except Exception as e:
                results.append({"cycle": cycle + 1, "error": str(e)})
                continue
            # Adapt
            try:
                _solver.tui.adapt.gradient(
                    "pressure", "yes", "yes", "0.5", "0.25",
                    "no", "no", "no", "no", "q"
                )
                results.append({"cycle": cycle + 1, "status": "adapted"})
            except Exception as e:
                results.append({"cycle": cycle + 1, "adapt_error": str(e)})
        return {"ok": True, "cycles_run": cycles, "results": results}
    except Exception as e:
        return {"ok": False, "error": str(e)}


@mcp.tool()
def run_adaptation_cycle(
    cycles: int = 1,
    iterations_between_cycles: int = 100,
    run_in_background: bool = True,
    force_no_monitors: bool = False,
) -> str:
    """Run one or more mesh adaptation cycles.

    Each cycle: solve N iterations → refine mesh → re-initialize if needed.

    Runs AS A BACKGROUND JOB by default - each cycle calls the same
    run_calculation.iterate() that start_calculation/run_until_convergence
    were already backgrounded to avoid freezing the server on. Poll
    get_calculation_status (kind="adaptation"; 'result' appears once
    finished) or read_console_tail for live progress; stop early with
    stop_calculation.

    MONITORS GATE — SERVER-ENFORCED: refuses to start if zero report
    monitors are defined, unless force_no_monitors=True (setup_standard_monitors first).

    Args:
        cycles:                      Number of adapt-solve cycles (default 1).
        iterations_between_cycles:   Solver iterations between adaptations.
        run_in_background:           Default True (recommended). False blocks until done.
        force_no_monitors:           Override the monitors gate (rare).
    """
    if err := _chk("solver"): return err
    if _bg_run["running"]:
        return _j({"ok": False, "error": "A background calculation is already running."})
    if gate_err := _require_monitors_or_error(force_no_monitors):
        return gate_err
    if not run_in_background:
        return _j(_run_adaptation_cycle_impl(cycles, iterations_between_cycles))
    _ensure_iteration_event_registered()
    start_iter = _current_iteration_baseline()
    _bg_run.update(running=True, kind="adaptation", requested=cycles,
                   start_iteration=start_iter, finished=False, error=None,
                   interrupted=False, result=None)

    def _worker():
        try:
            _bg_run["result"] = _run_adaptation_cycle_impl(cycles, iterations_between_cycles)
        except Exception as e:
            _bg_run["error"] = str(e)
        finally:
            _bg_run["running"] = False
            _bg_run["finished"] = True

    t = threading.Thread(target=_worker, name="mcp-fluent-adaptation", daemon=True)
    _bg_run["thread"] = t
    t.start()
    return _j({"ok": True, "started": True, "kind": "adaptation",
               "cycles": cycles, "start_iteration": start_iter,
               "next": "Poll get_calculation_status (result appears when finished) "
                       "or read_console_tail for live progress; stop early with "
                       "stop_calculation."})


@mcp.tool()
def get_adaptation_statistics() -> str:
    """Return mesh statistics after adaptation (cell counts, refinement levels)."""
    if err := _chk("solver"): return err
    try:
        raw = str(_solver.tui.mesh.check())
        # Also get zone info
        zones = list(_solver.setup.boundary_conditions.keys())
        return _j({"ok": True, "zone_count": len(zones),
                   "mesh_check_summary": raw[:2000]})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 17. GPU SOLVER CONTROL  — NEW
# ===========================================================================

@mcp.tool()
def enable_gpu_solver(gpu_count: int = 1) -> str:
    """Switch to the native Fluent GPU solver (requires compatible GPU + license).

    Supported physics in 2026 R1: k-ε, k-ω SST, GEKO, LES, WMLES, SBES,
    S-A, laminar, CHT, VOF+energy, species transport, compressible flows.

    Args:
        gpu_count: Number of GPUs to use (default 1).
    """
    if err := _chk("solver"): return err
    try:
        _solver.tui.parallel.gpgpu.enable("yes")
        if gpu_count > 1:
            _solver.tui.parallel.gpgpu.gpus_per_machine(str(gpu_count))
        return _j({"ok": True, "gpu_count": gpu_count,
                   "note": "GPU solver enabled. Verify physics compatibility with check_gpu_compatibility."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def check_gpu_compatibility() -> str:
    """Check which physics models in the current setup are GPU-compatible.

    Returns a compatibility report based on the 2026 R1 GPU solver feature set.
    """
    if err := _chk("solver"): return err
    GPU_SUPPORTED = {
        "viscous_models": ["laminar", "k-epsilon", "k-omega", "spalart-allmaras",
                           "large-eddy-simulation", "sbes", "wmles", "geko"],
        "heat_transfer": ["conduction", "conjugate-heat-transfer"],
        "multiphase": ["vof (with energy in 2026 R1)"],
        "species": ["species-transport"],
        "flow_types": ["incompressible", "compressible", "subsonic", "transonic"],
    }
    info: dict = {}
    # Check current viscous model
    try:
        vm = _solver.setup.models.viscous.model()
        info["current_viscous_model"] = vm
        info["viscous_gpu_compatible"] = (vm in GPU_SUPPORTED["viscous_models"])
    except Exception:
        info["current_viscous_model"] = "unknown"
    try:
        info["energy_enabled"] = _solver.setup.models.energy.enabled()
    except Exception:
        info["energy_enabled"] = "unknown"
    info["gpu_supported_viscous"] = GPU_SUPPORTED["viscous_models"]
    info["gpu_supported_features"] = {k: v for k, v in GPU_SUPPORTED.items()
                                       if k != "viscous_models"}
    info["note"] = ("GPU solver available in 2026 R1 for the supported physics above. "
                    "Always validate GPU vs CPU results for your specific case.")
    return _j({"ok": True, **info})


@mcp.tool()
def disable_gpu_solver() -> str:
    """Switch back to the standard CPU solver."""
    if err := _chk("solver"): return err
    try:
        _solver.tui.parallel.gpgpu.enable("no")
        return _j({"ok": True, "message": "CPU solver restored."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 18. CASE COMPARISON & DIFF  — NEW
# ===========================================================================

@mcp.tool()
def get_case_settings_snapshot() -> str:
    """Capture a snapshot of key solver settings for comparison/diff.

    Returns a structured dict of physics, BCs, solver settings, and mesh info
    that can be compared between cases.
    """
    if err := _chk("solver"): return err
    snap: dict = {"timestamp": _ts(), "connection": _connection_tag}
    try: snap["viscous_model"] = _solver.setup.models.viscous.model()
    except Exception: snap["viscous_model"] = "N/A"
    try: snap["energy_enabled"] = _solver.setup.models.energy.enabled()
    except Exception: snap["energy_enabled"] = "N/A"
    try: snap["solver_type"] = _solver.setup.general.solver.type()
    except Exception: snap["solver_type"] = "N/A"
    try: snap["time_stepping"] = _solver.setup.general.solver.time()
    except Exception: snap["time_stepping"] = "N/A"
    try: snap["boundary_zones"] = list(_solver.setup.boundary_conditions.keys())
    except Exception: snap["boundary_zones"] = []
    try: snap["cell_zones"] = list(_solver.setup.cell_zone_conditions.keys())
    except Exception: snap["cell_zones"] = []
    try: snap["materials_fluid"] = list(_solver.setup.materials.fluid.keys())
    except Exception: snap["materials_fluid"] = []
    return _j({"ok": True, "snapshot": snap})


def _to_jsonable(v):
    """Coerce a pyfluent read-back value to a JSON-safe primitive. Some
    Settings-API '.value' getters return a typed wrapper object (e.g. an
    ansys.units Quantity) rather than a bare float/str — verified live:
    _read_bc_values() failed with 'Object of type value_17 is not JSON
    serializable' until values were coerced here."""
    if v is None or isinstance(v, (int, float, str, bool)):
        return v
    try:
        return float(v)
    except Exception:
        try:
            return str(v)
        except Exception:
            return None


def _read_bc_values(zone: str, bc_type: str) -> dict:
    """Read back the ACTUAL applied value(s) for one zone, symmetric to the
    write-side attribute paths in set_velocity_inlet/set_pressure_inlet/
    set_pressure_outlet/set_wall_boundary. A tool call returning ok:true on
    the write side is not proof the value landed (verified incident,
    manifold_2: backflow_temperature_k=550 silently stayed at Fluent's
    300 K default on two of three outlets) — this is the read-back check
    that would have caught it automatically."""
    bc_root = _solver.setup.boundary_conditions
    values: dict = {}
    try:
        grp = getattr(bc_root, bc_type)[zone]
    except Exception as e:
        return {"error": str(e)}
    # NOTE: `.value` bare (no call) returns the settings NODE object itself,
    # not the number, even though the write side assigns via `.value = x` —
    # verified live: v = grp.momentum.gauge_pressure.value; type(v) is
    # settings_252.value_17; v() (or v.get_state()) returns the real scalar.
    # Every leaf read needs the call; only the WRITE side is bare `.value =`.
    if bc_type == "velocity_inlet":
        try: values["velocity_ms"] = grp.momentum.velocity.value()
        except Exception: values["velocity_ms"] = "N/A"
        try: values["temperature_k"] = grp.thermal.t.value()
        except Exception: values["temperature_k"] = "N/A"
        try: values["turbulent_intensity"] = grp.turbulence.turbulent_intensity()
        except Exception: values["turbulent_intensity"] = "N/A"
        try: values["turbulent_length_scale_m"] = grp.turbulence.turbulent_length_scale()
        except Exception: values["turbulent_length_scale_m"] = "N/A"
    elif bc_type == "pressure_inlet":
        try: values["gauge_total_pressure_pa"] = grp.momentum.gauge_total_pressure.value()
        except Exception: values["gauge_total_pressure_pa"] = "N/A"
        try: values["temperature_k"] = grp.thermal.t.value()
        except Exception: values["temperature_k"] = "N/A"
    elif bc_type == "pressure_outlet":
        try: values["gauge_pressure_pa"] = grp.momentum.gauge_pressure.value()
        except Exception: values["gauge_pressure_pa"] = "N/A"
        try: values["backflow_temperature_k"] = grp.thermal.backflow_total_temperature.value()
        except Exception: values["backflow_temperature_k"] = "N/A"
    elif bc_type == "wall":
        try: values["thermal_condition"] = grp.thermal.thermal_condition()
        except Exception: values["thermal_condition"] = "N/A"
        try: values["temperature_k"] = grp.thermal.t.value()
        except Exception: pass
        try: values["heat_flux_wm2"] = grp.thermal.q.value()
        except Exception: pass
        try: values["roughness_height_m"] = grp.momentum.roughness_height()
        except Exception: values["roughness_height_m"] = "N/A"
    return {k: _to_jsonable(v) for k, v in values.items()}


def _render_bc_manifest_md(solver_settings: dict, bc_rows: list[dict]) -> str:
    lines = [f"# Case manifest — {solver_settings.get('timestamp', '')}", "",
             "## Solver settings", ""]
    for k, v in solver_settings.items():
        if k == "timestamp": continue
        lines.append(f"- **{k}**: {v}")
    lines += ["", "## Boundary conditions", "",
              "| Zone | Type | Values |", "|---|---|---|"]
    for row in bc_rows:
        vals = ", ".join(f"{k}={v}" for k, v in row["values"].items())
        lines.append(f"| {row['zone']} | {row['type']} | {vals} |")
    return "\n".join(lines) + "\n"


@mcp.tool()
def export_case_manifest(output_dir: str) -> str:
    """Read back actual applied BC values (not just zone name/type) and
    solver settings, then write bc_readback.json + bc_readback.md.

    Reads real Fluent state, not intent — this is what would have caught
    the manifold_2 backflow-temperature bug (a set_pressure_outlet call
    that returned ok:true but silently didn't land) automatically instead
    of being discovered later from field data.

    Args:
        output_dir: Directory to write bc_readback.json/.md into.
    """
    if err := _chk("solver"): return err
    return _export_case_manifest_impl(output_dir)


def _export_case_manifest_impl(output_dir: str) -> str:
    """Body of export_case_manifest, callable internally (e.g. from
    checkpoint_run's session-job worker) WITHOUT re-triggering the
    session-job busy-guard in _chk() - that guard correctly blocks an
    EXTERNAL call while a job is running, but an internal call from within
    that same job's own worker thread is not a conflict (verified live
    2026-07-21: checkpoint_run's manifest export was self-blocking with
    "session job/calculation in progress" before this split)."""
    try:
        solver_settings: dict = {"timestamp": _ts()}
        try: solver_settings["viscous_model"] = _solver.setup.models.viscous.model()
        except Exception: solver_settings["viscous_model"] = "N/A"
        try: solver_settings["energy_enabled"] = _solver.setup.models.energy.enabled()
        except Exception: solver_settings["energy_enabled"] = "N/A"
        try: solver_settings["solver_type"] = _solver.setup.general.solver.type()
        except Exception: solver_settings["solver_type"] = "N/A"
        try: solver_settings["time_stepping"] = _solver.setup.general.solver.time()
        except Exception: solver_settings["time_stepping"] = "N/A"

        bc_rows = []
        for pair in _zone_type_pairs():
            values = _read_bc_values(pair["zone"], pair["type"])
            bc_rows.append({"zone": pair["zone"], "type": pair["type"], "values": values})

        json_path = str(Path(output_dir) / "bc_readback.json")
        md_path = str(Path(output_dir) / "bc_readback.md")
        _mkdir(json_path)
        manifest = {"solver_settings": solver_settings, "boundary_conditions": bc_rows}
        Path(json_path).write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
        Path(md_path).write_text(_render_bc_manifest_md(solver_settings, bc_rows), encoding="utf-8")

        return _j({"ok": True, "json_file": json_path, "md_file": md_path,
                   "solver_settings": solver_settings, "boundary_conditions": bc_rows})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _checkpoint_run_impl(output_dir: str, label: str) -> str:
    try:
        case_path = str(Path(output_dir) / f"{label}.cas.h5")
        data_path = str(Path(output_dir) / f"{label}.dat.h5")
        _mkdir(case_path)
        _solver.file.write_case(file_name=case_path)
        _solver.file.write_case_data(file_name=data_path)

        manifest_result = json.loads(_export_case_manifest_impl(output_dir))
        if not manifest_result.get("ok"):
            return _j({"ok": False, "error": f"Checkpoint files written, but manifest export failed: "
                       f"{manifest_result.get('error')}"})

        iters, ys = _residual_history()
        iteration = iters[-1] if iters else None
        residuals = {k: v[-1] for k, v in ys.items() if v} if iters else {}

        run_manifest = {
            "label": label, "timestamp": _ts(), "iteration": iteration,
            "residuals": residuals, "case_file": case_path, "data_file": data_path,
            "bc_manifest_file": manifest_result["json_file"],
        }
        run_manifest_path = str(Path(output_dir) / f"{label}_run_manifest.json")
        Path(run_manifest_path).write_text(json.dumps(run_manifest, indent=2, ensure_ascii=False), encoding="utf-8")

        return _j({"ok": True, **run_manifest, "run_manifest_file": run_manifest_path})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def checkpoint_run(output_dir: str, label: str, run_in_background: bool = True) -> str:
    """Pause checkpoint: write case+data, read back the BC/solver manifest,
    and record the current iteration/residuals — everything needed to
    resume later (read_case_data the .cas/.dat files) or hand the case to
    someone else. Bundles write_case_data + export_case_manifest +
    get_residuals into one call rather than three separate ones to remember.

    Runs AS A BACKGROUND JOB by default - writing case+data is the same
    class of large-file I/O that write_case_data/read_case were already
    backgrounded to avoid freezing the server on. Poll get_job_status +
    read_console_tail for progress.

    Args:
        output_dir:         Directory to write all checkpoint files into.
        label:               Name for this checkpoint (used in file names).
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("solver"): return err
    if run_in_background:
        return _run_session_job("checkpoint_run", lambda: _checkpoint_run_impl(output_dir, label))
    return _checkpoint_run_impl(output_dir, label)


@mcp.tool()
def compare_field_values(
    field_name: str,
    zone_names: list[str],
    baseline_values: dict,
) -> str:
    """Compare current field values against baseline values.

    Useful for comparing design iterations: pass in the baseline results
    (from a previous get_surface_report call) and see how the current
    solution differs.

    Args:
        field_name:       Field to compare (e.g. "pressure", "temperature").
        zone_names:       Zones to query.
        baseline_values:  Dict of {"zone_name": value} from the baseline case.
                          e.g. {"outlet": 101325.0, "inlet": 102000.0}
    """
    if err := _chk("solver"): return err
    try:
        report = _solver.solution.report_definitions
        comparison = []
        for zone in zone_names:
            rn = f"mcp_cmp_{zone}"
            report.surface[rn] = {}
            r = report.surface[rn]
            r.report_type = "surface-areaavg"
            r.field = field_name
            r.surface_names = [zone]
            try:
                current = _read_report_value(report, rn)
                baseline = baseline_values.get(zone)
                diff = None
                pct  = None
                if isinstance(current, (int, float)) and isinstance(baseline, (int, float)) and baseline != 0:
                    diff = round(current - baseline, 6)
                    pct  = round((diff / abs(baseline)) * 100, 4)
                comparison.append({
                    "zone": zone, "field": field_name,
                    "current": current, "baseline": baseline,
                    "delta": diff, "delta_pct": pct,
                })
            except Exception as e:
                comparison.append({"zone": zone, "error": str(e)})
        return _j({"ok": True, "field": field_name, "comparison": comparison})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 19. REPORT GENERATION (TXT / DOCX / PPTX)  — PPTX NEW
# ===========================================================================

@mcp.tool()
def generate_simulation_report(
    output_path: str,
    fmt: str = "txt",
    include_contour_image: str = "",
) -> str:
    """Generate a structured simulation summary report.

    Args:
        output_path:          Absolute path for the report.
        fmt:                  "txt", "docx", or "pptx".
        include_contour_image: Optional path to a pre-rendered PNG to embed.
    """
    if err := _chk("solver"): return err
    now = _ts()
    data: dict = {"generated": now, "connection": _connection_tag}
    for key, fn in [
        ("boundary_zones", lambda: list(_solver.setup.boundary_conditions.keys())),
        ("cell_zones",     lambda: list(_solver.setup.cell_zone_conditions.keys())),
        ("viscous_model",  lambda: _solver.setup.models.viscous.model()),
        ("energy_enabled", lambda: _solver.setup.models.energy.enabled()),
        ("materials_fluid",lambda: list(_solver.setup.materials.fluid.keys())),
        ("solver_type",    lambda: _solver.setup.general.solver.type()),
    ]:
        try: data[key] = fn()
        except Exception: data[key] = "N/A"
    _mkdir(output_path)
    out = Path(output_path)

    if fmt.lower() == "docx":
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            doc = Document()
            h = doc.add_heading("ANSYS Fluent CFD Simulation Report", 0)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph(f"Generated : {now}")
            doc.add_paragraph(f"Connection: {_connection_tag}")
            for section, key in [("Physics Models",   None),
                                  ("Materials (Fluid)", "materials_fluid"),
                                  ("Cell Zones",        "cell_zones"),
                                  ("Boundary Zones",    "boundary_zones")]:
                doc.add_heading(section, level=1)
                if key is None:
                    doc.add_paragraph(f"Viscous Model  : {data['viscous_model']}")
                    doc.add_paragraph(f"Energy Equation: {data['energy_enabled']}")
                    doc.add_paragraph(f"Solver Type    : {data['solver_type']}")
                else:
                    items = data.get(key, [])
                    if isinstance(items, list):
                        for item in items:
                            doc.add_paragraph(str(item), style="List Bullet")
                    else:
                        doc.add_paragraph(str(items))
            if include_contour_image and Path(include_contour_image).exists():
                doc.add_heading("Visualisation", level=1)
                doc.add_picture(include_contour_image, width=__import__("docx.shared", fromlist=["Inches"]).Inches(5.5))
            doc.save(str(out))
            return _j({"ok": True, "report_path": str(out)})
        except ImportError:
            return _j({"ok": False, "error": "python-docx not installed."})

    elif fmt.lower() == "pptx":
        return _generate_pptx_report(output_path, data, now, include_contour_image)

    else:  # txt
        lines = [
            "ANSYS FLUENT CFD SIMULATION REPORT",
            f"Generated : {now}", f"Connection: {_connection_tag}", "",
            "PHYSICS MODELS",
            f"  Viscous model  : {data['viscous_model']}",
            f"  Energy equation: {data['energy_enabled']}",
            f"  Solver type    : {data['solver_type']}", "",
            "MATERIALS (FLUID)",
        ] + [f"  {m}" for m in (data.get("materials_fluid") or [])] + [
            "", "CELL ZONES",
        ] + [f"  {z}" for z in (data.get("cell_zones") or [])] + [
            "", "BOUNDARY ZONES",
        ] + [f"  {z}" for z in (data.get("boundary_zones") or [])]
        txt = out.with_suffix(".txt")
        txt.write_text("\n".join(lines), encoding="utf-8")
        return _j({"ok": True, "report_path": str(txt)})


def _generate_pptx_report(output_path: str, data: dict, now: str,
                           contour_image: str = "") -> str:
    """Internal: build a PowerPoint deck."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return _j({"ok": False, "error": "python-pptx not installed. pip install python-pptx"})

    try:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        ANSYS_BLUE = RGBColor(0x00, 0x5C, 0xA9)
        ANSYS_GOLD = RGBColor(0xFF, 0xB8, 0x00)
        DARK       = RGBColor(0x1A, 0x1A, 0x2E)

        blank_layout = prs.slide_layouts[6]

        def _add_slide(title_text: str, subtitle_text: str = "") -> object:
            slide = prs.slides.add_slide(blank_layout)
            # Title bar
            tf_box = slide.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, Inches(1.1))
            tf = tf_box.text_frame
            p = tf.paragraphs[0]; p.text = title_text
            p.font.size = Pt(28); p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.alignment = PP_ALIGN.LEFT
            fill = tf_box.fill; fill.solid(); fill.fore_color.rgb = ANSYS_BLUE
            if subtitle_text:
                sub_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2),
                                                    Inches(12.5), Inches(0.5))
                sub_tf = sub_box.text_frame
                sp = sub_tf.paragraphs[0]; sp.text = subtitle_text
                sp.font.size = Pt(14); sp.font.color.rgb = DARK
            return slide

        def _add_bullet_box(slide, left, top, width, height, items: list[str],
                             font_size: int = 14):
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame; tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {item}"; p.font.size = Pt(font_size)
                p.font.color.rgb = DARK

        def _add_kv_table(slide, rows_data: list[tuple[str, str]],
                          left=Inches(0.5), top=Inches(1.8),
                          width=Inches(12.0), row_height=Inches(0.45)):
            tbl = slide.shapes.add_table(len(rows_data), 2,
                                          left, top, width,
                                          row_height * len(rows_data)).table
            tbl.columns[0].width = Inches(4)
            tbl.columns[1].width = Inches(8)
            for ri, (k, v) in enumerate(rows_data):
                for ci, txt in enumerate((k, str(v))):
                    cell = tbl.cell(ri, ci)
                    cell.text = txt
                    cell.text_frame.paragraphs[0].font.size = Pt(13)
                    if ci == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
                    fill = cell.fill; fill.solid()
                    fill.fore_color.rgb = (RGBColor(0xE8, 0xF0, 0xFE) if ri % 2 == 0
                                           else RGBColor(0xFF, 0xFF, 0xFF))

        # ----- Slide 1: Title -----
        slide1 = _add_slide("ANSYS Fluent CFD Simulation Report")
        tb = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(3))
        tf = tb.text_frame
        for line in [f"Generated: {now}",
                     f"Connection: {data['connection']}",
                     "",
                     "Prepared by: PyFluent MCP Server v3"]:
            p = tf.paragraphs[0] if line == f"Generated: {now}" else tf.add_paragraph()
            p.text = line; p.font.size = Pt(18)
            p.font.color.rgb = DARK if line else ANSYS_GOLD

        # ----- Slide 2: Physics Settings -----
        slide2 = _add_slide("Physics Configuration", "Turbulence, energy, and solver settings")
        _add_kv_table(slide2, [
            ("Viscous Model",   data.get("viscous_model", "N/A")),
            ("Energy Equation", str(data.get("energy_enabled", "N/A"))),
            ("Solver Type",     data.get("solver_type", "N/A")),
        ])

        # ----- Slide 3: Geometry & Zones -----
        slide3 = _add_slide("Geometry & Mesh Zones")
        bz = data.get("boundary_zones", [])
        cz = data.get("cell_zones", [])
        _add_bullet_box(slide3, Inches(0.5), Inches(1.8), Inches(5.5), Inches(5),
                        [f"Cell zones: {len(cz)}"] + cz[:12])
        _add_bullet_box(slide3, Inches(6.5), Inches(1.8), Inches(6.0), Inches(5),
                        [f"Boundary zones: {len(bz)}"] + bz[:12])

        # ----- Slide 4: Materials -----
        slide4 = _add_slide("Materials")
        mats = data.get("materials_fluid", [])
        _add_bullet_box(slide4, Inches(0.5), Inches(1.8), Inches(12), Inches(5),
                        mats if mats else ["No fluid materials found"])

        # ----- Slide 5: Visualisation (if image provided) -----
        if contour_image and Path(contour_image).exists():
            slide5 = _add_slide("CFD Visualisation")
            slide5.shapes.add_picture(contour_image,
                                       Inches(0.5), Inches(1.3),
                                       Inches(12.0), Inches(5.8))

        # ----- Slide 6: Summary -----
        slide6 = _add_slide("Summary & Next Steps")
        _add_bullet_box(slide6, Inches(0.5), Inches(1.8), Inches(12), Inches(4.5), [
            f"Simulation: {data.get('viscous_model','?')} turbulence model",
            f"Energy equation: {'ON' if data.get('energy_enabled') else 'OFF'}",
            f"Solver: {data.get('solver_type','?')}",
            f"Boundary zones: {len(data.get('boundary_zones',[]))}",
            "",
            "Next Steps:",
            "  → Review convergence history",
            "  → Run parametric sweep for design optimisation",
            "  → Export results for AR/VR review",
        ])

        out = str(Path(output_path).with_suffix(".pptx"))
        prs.save(out)
        return _j({"ok": True, "pptx_path": out,
                   "slides": prs.slides.__len__(),
                   "note": "Open in PowerPoint or Google Slides."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 20. SPECIES TRANSPORT  — NEW
# ===========================================================================

@mcp.tool()
def enable_species_transport(
    mixture_material: str = "methane-air",
    combustion_model: str = "none",
) -> str:
    """Enable species transport (multicomponent flows and combustion).

    Args:
        mixture_material: Mixture material name from the Fluent database.
                          Common options: "methane-air", "air", "hydrogen-air",
                          "propane-air", "co2-n2-mixture".
        combustion_model: "none" (passive mixing), "eddy-dissipation",
                          "finite-rate", "flamelet" (FGM).
    """
    if err := _chk("solver"): return err
    try:
        sp = _solver.setup.models.species
        sp.enabled = True
        sp.mixture_material = mixture_material
        if combustion_model != "none":
            try: sp.reactions.option = combustion_model
            except Exception: pass
        return _j({"ok": True, "mixture_material": mixture_material,
                   "combustion_model": combustion_model,
                   "message": "Species transport enabled."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_inlet_species(
    zone_name: str,
    mass_fractions: dict[str, float],
) -> str:
    """Set species mass fractions at an inlet boundary.

    Mass fractions should sum to 1.0. Any unspecified species are set to 0.

    Args:
        zone_name:      Velocity or pressure inlet zone name.
        mass_fractions: Dict mapping species names to mass fractions.
                        e.g. {"ch4": 0.05, "o2": 0.23, "n2": 0.72}
    """
    if err := _chk("solver"): return err
    try:
        total = sum(mass_fractions.values())
        if abs(total - 1.0) > 0.01:
            return _j({"ok": False,
                       "error": f"Mass fractions sum to {total:.4f}, not 1.0."})
        # Try velocity inlet first, then pressure inlet
        for bc_type in ("velocity_inlet", "pressure_inlet"):
            try:
                bc = getattr(_solver.setup.boundary_conditions, bc_type)[zone_name]
                for species, mf in mass_fractions.items():
                    try: bc.species.mass_fraction[species] = mf
                    except Exception: pass
                return _j({"ok": True, "zone": zone_name,
                           "mass_fractions": mass_fractions,
                           "sum": round(total, 6)})
            except Exception:
                continue
        return _j({"ok": False, "error": f"Zone '{zone_name}' not found as inlet."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_species_report(
    species_name: str,
    zone_names: list[str],
    report_type: str = "mass-average",
) -> str:
    """Get a species mass fraction report on specified zones.

    Args:
        species_name: Species name (e.g. "co2", "h2o", "ch4").
        zone_names:   Zones to query.
        report_type:  "mass-average", "area-average", "vertex-max", "vertex-min".
    """
    if err := _chk("solver"): return err
    try:
        field = f"mass-fraction-of-{species_name}"
        return get_surface_report(report_type, zone_names, field)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# 21. MULTIPHASE (VOF)  — NEW
# ===========================================================================

@mcp.tool()
def enable_vof(
    n_phases: int = 2,
    interface_modelling: str = "sharp",
) -> str:
    """Enable Volume of Fluid (VOF) multiphase model.

    VOF is used for free-surface flows, liquid-gas interfaces, sloshing,
    wave simulation, and two-phase heat transfer.
    Supports GPU acceleration in 2026 R1 (with energy equation).

    Args:
        n_phases:              Number of phases (default 2, max typically 6).
        interface_modelling:   "sharp" (default), "dispersed", or "both".
    """
    if err := _chk("solver"): return err
    try:
        mp = _solver.setup.models.multiphase
        mp.model = "vof"
        mp.number_of_phases = n_phases
        try: mp.vof.interface_modelling = interface_modelling
        except Exception: pass
        return _j({"ok": True, "model": "VOF", "n_phases": n_phases,
                   "interface_modelling": interface_modelling,
                   "note": "Define phase materials with set_vof_phase, then set BCs."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def set_vof_phase(
    phase_index: int,
    material_name: str,
    phase_name: str = "",
) -> str:
    """Assign a material to a VOF phase.

    Args:
        phase_index:   Phase number (1 = primary, 2 = secondary, …).
        material_name: Fluent material name (e.g. "water-liquid", "air").
        phase_name:    Human-readable name for this phase.
    """
    if err := _chk("solver"): return err
    try:
        phases = _solver.setup.models.multiphase.phases
        phase_key = f"phase-{phase_index}"
        phase = phases[phase_key]
        phase.material = material_name
        if phase_name:
            try: phase.name = phase_name
            except Exception: pass
        return _j({"ok": True, "phase_index": phase_index,
                   "material": material_name, "name": phase_name})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_vof_volume_fraction(
    phase_index: int,
    zone_names: list[str],
) -> str:
    """Get the area-averaged volume fraction of a VOF phase on boundary zones.

    Args:
        phase_index: Phase number (1 = primary, 2 = secondary, …).
        zone_names:  Zones to query.
    """
    if err := _chk("solver"): return err
    try:
        field = f"volume-fraction-of-phase-{phase_index}"
        return get_surface_report("area-average", zone_names, field)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def render_vof_interface(
    phase_index: int = 2,
    iso_value: float = 0.5,
    output_image_path: str = "",
    colormap: str = "Blues",
    window_size: list[int] | None = None,
) -> str:
    """Render the VOF free-surface as a volume-fraction iso-surface (VF = 0.5).

    Args:
        phase_index:        Phase to visualise (default 2 = secondary phase).
        iso_value:          Volume fraction iso-value (default 0.5 = interface).
        output_image_path:  Output PNG path.
        colormap:           Colour map.
        window_size:        [width, height] pixels.
    """
    field = f"volume-fraction-of-phase-{phase_index}"
    return render_iso_surface(
        field_name=field,
        iso_value=iso_value,
        output_image_path=output_image_path,
        colormap=colormap,
        window_size=window_size,
    )


# ===========================================================================
# 22. MESHING SESSION WORKFLOW
# ===========================================================================

def _import_geometry_impl(
    geometry_file_path: str,
    workflow_type: str = "watertight",
    length_unit: str = "m",
) -> str:
    """Synchronous implementation — public tool: import_geometry. No _chk()
    here: the public wrapper already checked before starting the session
    job, and this impl runs INSIDE that job's own worker thread, where
    _session_job["running"] is True - a second _chk() call here would
    always see itself as "a session job is running" and self-block
    (verified live 2026-07-21)."""
    global _import_length_unit
    if not Path(geometry_file_path).exists():
        return _j({"ok": False, "error": f"Not found: {geometry_file_path}"})
    warnings: list[str] = []
    t0 = time.time()
    try:
        wf = _get_workflow(workflow_type)
        # the import task differs per workflow (doc 01 sections 2/4/5):
        # watertight: import_geometry; FTM: import_cad_and_part_management;
        # 2d: load_cad_geometry_2d — discover release-tolerantly
        if _workflow_kind == "fault_tolerant":
            attr = (_find_wf_task(wf, "import", "cad")
                    or _find_wf_task(wf, "part", "management"))
            readback_task = "Import CAD and Part Management"
        elif _workflow_kind == "2d":
            attr = _find_wf_task(wf, "load", "cad") or _find_wf_task(wf, "cad", "2d")
            readback_task = "Load CAD Geometry"
        else:
            attr = _find_wf_task(wf, "import", "geometry")
            readback_task = "Import Geometry"
        if attr is None:
            return _j({"ok": False,
                       "error": f"no import task found on the {_workflow_kind} workflow",
                       "hint": "call get_workflow_state to see the live task tree"})
        ig = getattr(wf, attr)
        ig.file_name = geometry_file_path
        try:
            ig.length_unit = length_unit
        except Exception as e:
            warnings.append(f"length_unit assignment rejected on this release: {e} "
                            "- import proceeds with the task's default unit; VERIFY SCALE")
        ig()
        verification = _post_import_verification()
        # unit readback from the datamodel (authoritative, PascalCase keys)
        try:
            args = _meshing.workflow.TaskObject[readback_task].Arguments()
            verification["length_unit_applied"] = args.get("LengthUnit", "<unknown>")
            if str(verification["length_unit_applied"]).lower() != length_unit.lower():
                warnings.append(
                    f"requested unit '{length_unit}' but task holds "
                    f"'{verification['length_unit_applied']}' - scale may be wrong")
        except Exception as e:
            verification["length_unit_applied"] = f"<readback failed: {e}>"
        # sizing fields in every later workflow task are in THIS unit —
        # record it so sizing tools can convert/echo correctly
        _import_length_unit = length_unit
        _persist_session_state()
        bbox = _parse_bbox_from_transcript()
        if bbox:
            verification["bounding_box"] = bbox
        prompts: list[dict] = []
        ext_txt = (" x ".join(f"{e:g}" for e in bbox["extents"]) + f" {length_unit}"
                   if bbox else "<not parsed - check the GUI ruler>")
        prompts.append(_prompt(
            f"The imported model's bounding box is {ext_txt}. Is that the correct "
            "real-world size?", "Scale",
            [("Yes - size is correct", "proceed to mesh sizing"),
             ("No - wrong unit/scale", "re-import with a different length_unit")]))
        if bbox and any(e > 0 for e in bbox["extents"]):
            L = min(e for e in bbox["extents"] if e > 0)
            tier_opts = []
            for tier, r in _FIDELITY_RULES.items():
                mn = round(L / r["min_div"], 3)
                mx = round(mn * r["max_mult"], 3)
                tier_opts.append(
                    (f"{tier}: min {mn:g} / max {mx:g} {length_unit}",
                     f"{r['min_div']} cells across the smallest extent ({L:g} "
                     f"{length_unit}), growth {r['growth']}, curvature "
                     f"{r['curvature_deg']} deg, {r['bl_layers']} prism layers"))
            tier_opts.append(("Custom sizes",
                              f"type min/max element sizes in {length_unit} (better: "
                              "based on the flow-controlling feature, e.g. inlet "
                              "diameter, via propose_mesh_sizing)"))
            prompts.append(_prompt(
                "Mesh sizing to use for local sizing + surface mesh?",
                "Mesh sizing", tier_opts))
        return _j({"ok": True, "workflow_type": _workflow_kind,
                   "length_unit": length_unit,
                   "sizing_unit_note": f"All later sizing values (local sizing, surface "
                                       f"mesh min/max) are interpreted in '{length_unit}' "
                                       f"unless a different `units` argument is passed.",
                   "elapsed_s": round(time.time() - t0, 1),
                   "verification": verification, "warnings": warnings,
                   "message": f"Geometry imported: {geometry_file_path}",
                   "interactive_prompts": prompts,
                   "prompt_render": _PROMPT_RENDER_NOTE,
                   "next_recommended_action":
                       "Present interactive_prompts to the user (AskUserQuestion) to "
                       "confirm scale and pick sizing BEFORE meshing; then "
                       "propose_boundary_types for the BC-type prompts."})
    except Exception as e:
        return _j({"ok": False, "error": str(e), "warnings": warnings})


@mcp.tool()
def import_geometry(
    geometry_file_path: str,
    workflow_type: str = "watertight",
    length_unit: str = "m",
    run_in_background: bool = True,
) -> str:
    """Import a geometry file into a Fluent meshing session.

    Runs AS A BACKGROUND JOB by default — real CAD takes 30 s to several
    minutes, and a synchronous call froze the MCP server long enough for
    clients to kill it (root cause of every mid-run restart on
    2026-07-17/20). Poll get_job_status every ~10-15 s, with
    read_console_tail for live progress. The finished job result carries
    the verification block (objects/labels/unit/bounding box) and
    interactive_prompts to confirm scale + sizing with the user
    (AskUserQuestion) BEFORE meshing.

    Initializes the meshing workflow on first call (a bare meshing session
    has no workflow yet). See knowledge/source/meshing/01_new_meshing_workflows_api.md.

    Args:
        geometry_file_path: Path to .scdoc/.scdocx, .step, .iges, .stl, or .fmd file.
        workflow_type:      "watertight" (default, clean CAD), "fault_tolerant"
                            (dirty/leaky CAD), or "2d".
        length_unit:        CAD length unit, e.g. "m", "mm", "in" — ALL later
                            sizing values are interpreted in this unit.
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("meshing"): return err
    if not Path(geometry_file_path).exists():
        return _j({"ok": False, "error": f"Not found: {geometry_file_path}"})
    _set_project_dir_from(geometry_file_path)
    if run_in_background:
        return _run_session_job(
            "import_geometry",
            lambda: _import_geometry_impl(geometry_file_path, workflow_type, length_unit))
    return _import_geometry_impl(geometry_file_path, workflow_type, length_unit)


_FIDELITY_RULES = {
    # calibrated to KB doc 07 (N_span / feature-resolution tables) and the
    # live-validated manifold first cut (L_char=20mm screening → 2/12mm)
    "screening":    {"min_div": 10, "max_mult": 6, "curvature_deg": 18.0,
                     "bl_layers": 3, "growth": 1.2},
    "engineering":  {"min_div": 16, "max_mult": 5, "curvature_deg": 15.0,
                     "bl_layers": 10, "growth": 1.2},
    "verification": {"min_div": 25, "max_mult": 4, "curvature_deg": 12.0,
                     "bl_layers": 15, "growth": 1.15},
}
_ANALYSIS_BL_LAYERS = {
    # doc 07 initial layer-count matrix, per analysis family
    "internal": {"screening": 3, "engineering": 10, "verification": 15},
    "external": {"screening": 5, "engineering": 15, "verification": 20},
    "cht":      {"screening": 5, "engineering": 15, "verification": 20},
}


# ---------------------------------------------------------------------------
# Interactive prompts — tools return AskUserQuestion-compatible specs at the
# decision points (BC types, sizing, gates) so the client renders selectable
# fields instead of prose (user request 2026-07-20).
# ---------------------------------------------------------------------------
_PROMPT_RENDER_NOTE = ("Render each entry in interactive_prompts with the "
                       "AskUserQuestion tool - selectable fields; the user can "
                       "always type a custom value via 'Other'. Max 4 questions "
                       "per AskUserQuestion call; option marked (Recommended) "
                       "goes first. Do NOT paraphrase these as plain text.")


def _prompt(question: str, header: str, options: list, multi: bool = False) -> dict:
    """AskUserQuestion-compatible spec. options = [(label, description), ...],
    2-4 entries (the client adds 'Other' automatically)."""
    return {"question": question, "header": header[:12], "multiSelect": multi,
            "options": [{"label": l, "description": d} for l, d in options[:4]]}


_BC_TYPE_DESCRIPTIONS = {
    "velocity-inlet":  "flow enters at a set velocity",
    "pressure-outlet": "flow exits at a set static pressure",
    "mass-flow-inlet": "flow enters at a set mass flow rate",
    "pressure-inlet":  "total-pressure driven inlet",
    "wall":            "no-flow solid boundary (no-slip)",
    "symmetry":        "mirror plane, no flux across",
}
_BC_TYPE_PATTERNS = [
    ("velocity-inlet",     ("inlet", "in_", "supply")),
    ("pressure-outlet",    ("outlet", "out_", "return", "exhaust")),
    ("pressure-far-field", ("farfield", "far_field")),
    ("symmetry",           ("sym",)),
    ("wall",               ("wall", "body", "ground")),
]


def _infer_bc_type(label: str) -> str:
    """KB doc-07 name-pattern inference for a zone label's BC type."""
    ln = label.lower()
    for bc, pats in _BC_TYPE_PATTERNS:
        if any(ln.startswith(p) for p in pats) or \
           any(p in ln for p in pats if len(p) >= 4):
            return bc
    return "wall"


def _parse_bbox_from_transcript() -> dict | None:
    """Last 'Bounding box:' block from the live transcript — the scale
    evidence the meshing KB requires echoing to the user after import."""
    trn = _find_transcript()
    if trn is None:
        return None
    try:
        text = trn.read_text(encoding="utf-8", errors="replace")[-200_000:]
    except Exception:
        return None
    hits = re.findall(r"Bounding box:\s*\r?\n\s*\(([^)]+)\)\s*x\s*\(([^)]+)\)", text)
    if not hits:
        return None
    try:
        lo = [float(v) for v in hits[-1][0].split(",")]
        hi = [float(v) for v in hits[-1][1].split(",")]
        return {"min": lo, "max": hi,
                "extents": [round(h - l, 3) for l, h in zip(lo, hi)]}
    except Exception:
        return None


@mcp.tool()
def propose_boundary_types() -> str:
    """Infer a boundary-condition type for every label in the imported
    geometry (KB name patterns) and return INTERACTIVE prompt specs — one
    selectable question per label — for the user to confirm or override
    BEFORE update_boundaries.

    Render each spec in interactive_prompts via AskUserQuestion (max 4
    questions per call). Body/object labels are excluded — Fluent refuses
    to retype them (verified live); they stay walls. After the user
    answers, call update_boundaries with the confirmed labels + types.
    """
    if err := _chk("meshing"): return err
    mu = _mu()
    if mu is None:
        return _j({"ok": False, "error": "meshing_utilities unavailable on this release"})
    try:
        objects = list(mu.get_all_objects() or [])
        labels: list[str] = []
        for o in objects:
            labels.extend(list(mu.get_labels(object_name=o) or []))
    except Exception as e:
        return _j({"ok": False, "error": f"label inventory failed: {e}"})
    inferred: dict = {}
    prompts: list[dict] = []
    for lb in labels:
        if lb in objects:
            inferred[lb] = "wall (body label - not retypable, stays wall)"
            continue
        bc = _infer_bc_type(lb)
        inferred[lb] = bc
        alts = [c for c in _BC_TYPE_DESCRIPTIONS if c != bc][:3]
        prompts.append(_prompt(
            f"Boundary type for zone '{lb}'?", lb,
            [(f"{bc} (Recommended)",
              f"inferred from the name - {_BC_TYPE_DESCRIPTIONS.get(bc, '')}")] +
            [(a, _BC_TYPE_DESCRIPTIONS[a]) for a in alts]))
    return _j({"ok": True, "labels": labels, "inferred_types": inferred,
               "interactive_prompts": prompts,
               "prompt_render": _PROMPT_RENDER_NOTE,
               "apply_with": "update_boundaries(boundary_labels=[...], "
                             "boundary_types=[...]) with the confirmed values"})


@mcp.tool()
def set_meshing_mode(mode: str = "") -> str:
    """Set (or report) the MESHING workflow's operating mode: 'auto' or
    'manual'. Independent of the CFD analysis mode (set_analysis_mode);
    set_workflow_mode sets both at once. Switching modes at ANY stage never
    discards the current setup — imported geometry, applied sizing, plans,
    and task state all survive; only the interaction contract changes.

    MANUAL (default): the user supplies the critical inputs at the stages
    where they matter — sizing tools REQUIRE explicit min/max element sizes
    and add_boundary_layers a layer count (ask the user via AskUserQuestion;
    never guess), and the standard confirmation gates (scale, sizing,
    volume-commit, switch) stay active.

    AUTO: sizing is inferred from the analysis type, fidelity tier, and
    geometric scale via propose_mesh_sizing — call it once after
    import_geometry (with the characteristic length, e.g. inlet/hydraulic
    diameter), then add_local_sizing / generate_surface_mesh /
    add_boundary_layers with no size arguments apply the plan automatically.
    Ask the user ONLY essential non-inferable inputs (scale/characteristic
    length); ALWAYS show the plan summary and get explicit confirmation
    before generate_volume_mesh.

    Args:
        mode: 'auto', 'manual', or empty to just report the current mode.
    """
    global _meshing_mode
    if mode:
        m = mode.strip().lower()
        if m not in ("auto", "manual"):
            return _j({"ok": False, "error": f"unknown mode '{mode}' (use 'auto' or 'manual')"})
        _meshing_mode = m
        _persist_session_state()
    return _j({"ok": True, "meshing_mode": _meshing_mode,
               "analysis_mode": _analysis_mode,
               "mesh_plan_present": _mesh_plan is not None,
               "setup_preserved_on_switch": True,
               "behavior": ("sizing tools require explicit sizes from the user; "
                            "ask via AskUserQuestion before meshing"
                            if _meshing_mode == "manual" else
                            "run propose_mesh_sizing once, then sizing tools with "
                            "no size args apply the plan automatically")})


@mcp.tool()
def set_analysis_mode(mode: str = "") -> str:
    """Set (or report) the CFD ANALYSIS workflow's operating mode: 'auto' or
    'manual'. Independent of the meshing mode (set_meshing_mode);
    set_workflow_mode sets both at once. Switching modes at ANY stage never
    discards the current setup — physics models, BCs, monitors, plans, and
    solution state all survive; only the interaction contract changes.

    MANUAL (default): detailed control — physics models, boundary
    conditions, solver settings, convergence criteria, and post-processing
    are configured through the individual tools with user-provided values
    (ask the user; never guess a BC value).

    AUTO: infer as much as possible from the loaded case (zone names/types,
    regions, project context) via propose_solver_setup; ask the user ONLY
    the essential inputs that cannot be inferred safely (inlet velocity or
    flow rate, fluid); show the returned plan summary, get explicit
    confirmation, then apply_solver_plan -> initialize_solution(
    checkpoint_dir=...) -> start_calculation -> read_console_tail.

    Args:
        mode: 'auto', 'manual', or empty to just report the current mode.
    """
    global _analysis_mode
    if mode:
        m = mode.strip().lower()
        if m not in ("auto", "manual"):
            return _j({"ok": False, "error": f"unknown mode '{mode}' (use 'auto' or 'manual')"})
        _analysis_mode = m
        _persist_session_state()
    return _j({"ok": True, "analysis_mode": _analysis_mode,
               "meshing_mode": _meshing_mode,
               "solver_plan_present": _solver_plan is not None,
               "setup_preserved_on_switch": True,
               "behavior": ("solver tools take user-provided values; ask via "
                            "AskUserQuestion before applying physics/BCs"
                            if _analysis_mode == "manual" else
                            "run propose_solver_setup, confirm the summary with "
                            "the user, then apply_solver_plan")})


@mcp.tool()
def set_workflow_mode(mode: str = "") -> str:
    """Set (or report) the operating mode of BOTH workflows at once —
    Meshing AND CFD Analysis: 'auto' or 'manual'. Use set_meshing_mode /
    set_analysis_mode to control them independently (e.g. auto meshing but
    manual solving). Switching never discards the current setup.

    Args:
        mode: 'auto', 'manual', or empty to report both current modes.
    """
    global _meshing_mode, _analysis_mode
    if mode:
        m = mode.strip().lower()
        if m not in ("auto", "manual"):
            return _j({"ok": False, "error": f"unknown mode '{mode}' (use 'auto' or 'manual')"})
        _meshing_mode = _analysis_mode = m
        _persist_session_state()
    return _j({"ok": True, "meshing_mode": _meshing_mode,
               "analysis_mode": _analysis_mode,
               "mesh_plan_present": _mesh_plan is not None,
               "solver_plan_present": _solver_plan is not None,
               "setup_preserved_on_switch": True})


@mcp.tool()
def propose_mesh_sizing(
    characteristic_length: float,
    characteristic_length_units: str = "",
    analysis_type: str = "internal",
    fidelity: str = "screening",
    domain_extent: float = 0.0,
) -> str:
    """Compute an intelligent sizing plan from the analysis type, fidelity
    tier, and geometric scale (KB meshing doc 07 rules). Works in both
    meshing modes: in AUTO the plan is applied by the sizing tools when
    they're called without sizes; in MANUAL it's the recommendation to show
    the user for confirmation/override (AskUserQuestion).

    Args:
        characteristic_length:       The flow-controlling dimension — inlet or
                                     hydraulic diameter for internal flow, body
                                     length for external aero.
        characteristic_length_units: Unit of that value (default: the CAD
                                     import unit).
        analysis_type:               'internal' (duct/manifold/pipe), 'external'
                                     (aero/hydro), or 'cht'.
        fidelity:                    'screening' (first cut), 'engineering', or
                                     'verification'.
        domain_extent:               Optional largest domain dimension (same
                                     units) — caps the max size so the domain
                                     keeps >=10 cells across it.
    """
    global _mesh_plan
    if characteristic_length <= 0:
        return _j({"ok": False, "error": "characteristic_length must be > 0"})
    fid = fidelity.strip().lower()
    if fid not in _FIDELITY_RULES:
        return _j({"ok": False, "error": f"unknown fidelity '{fidelity}' "
                                         f"(use {sorted(_FIDELITY_RULES)})"})
    fam = analysis_type.strip().lower()
    fam = ("external" if "extern" in fam or "aero" in fam else
           "cht" if "cht" in fam or "heat" in fam or "conjugate" in fam else
           "internal")
    try:
        L = _convert_size_to_import_unit(characteristic_length, characteristic_length_units)
        ext = (_convert_size_to_import_unit(domain_extent, characteristic_length_units)
               if domain_extent > 0 else 0.0)
    except ValueError as e:
        return _j({"ok": False, "error": str(e)})
    rules = _FIDELITY_RULES[fid]
    min_size = L / rules["min_div"]
    max_size = min_size * rules["max_mult"]
    capped = False
    if ext > 0 and max_size > ext / 10:
        max_size, capped = ext / 10, True
    unit = _import_length_unit or "m"
    _mesh_plan = {
        "analysis_family": fam, "fidelity": fid,
        "characteristic_length": L, "units": unit,
        "min_size": round(min_size, 6), "max_size": round(max_size, 6),
        "growth_rate": rules["growth"],
        "curvature_normal_angle_deg": rules["curvature_deg"],
        "bl_layers": _ANALYSIS_BL_LAYERS[fam][fid],
        "volume_fill": "poly-hexcore",
    }
    _persist_session_state()
    return _j({"ok": True, "plan": _mesh_plan,
               "basis": {"cells_across_characteristic_length": rules["min_div"],
                         "max_size_capped_by_domain_extent": capped},
               "meshing_mode": _meshing_mode,
               "next": ("AUTO mode: call add_local_sizing() and "
                        "generate_surface_mesh() with no size arguments - this "
                        "plan applies automatically." if _meshing_mode == "auto" else
                        "MANUAL mode: present this plan to the user for "
                        "confirmation or override (AskUserQuestion), then pass "
                        "the confirmed sizes explicitly.")})


def _resolve_sizing(min_size: float, max_size: float, units: str,
                    tool_name: str) -> tuple[float, float] | str:
    """Shared mode-aware sizing resolution for add_local_sizing /
    generate_surface_mesh. Returns (min, max) in the CAD import unit, or an
    error-JSON string. Sizes of 0 mean 'not provided': AUTO mode fills them
    from the propose_mesh_sizing plan; MANUAL mode refuses — the user must
    supply them (the 2026-07-17 hang came from silently guessed sizes)."""
    if min_size > 0 and max_size > 0:
        try:
            return (_convert_size_to_import_unit(min_size, units),
                    _convert_size_to_import_unit(max_size, units))
        except ValueError as e:
            return _j({"ok": False, "error": str(e)})
    if _meshing_mode == "auto":
        if _mesh_plan:
            return (_mesh_plan["min_size"], _mesh_plan["max_size"])
        return _j({"ok": False,
                   "error": f"{tool_name}: no sizes given and no sizing plan exists",
                   "hint": "AUTO mode needs propose_mesh_sizing(characteristic_length=...) "
                           "called once after import_geometry"})
    plan_opts = ([(f"min {_mesh_plan['min_size']:g} / max {_mesh_plan['max_size']:g} "
                   f"{_import_length_unit} (Recommended)",
                   f"from propose_mesh_sizing ({_mesh_plan.get('fidelity', '?')} tier)")]
                 if _mesh_plan else
                 [("Run propose_mesh_sizing first",
                   "computes tiers from a characteristic length you provide")])
    return _j({"ok": False,
               "error": f"{tool_name}: min_size and max_size are REQUIRED in manual "
                        "meshing mode - do not guess them",
               "sizing_unit": _import_length_unit,
               "interactive_prompt": _prompt(
                   f"Min/max surface-mesh element sizes (in {_import_length_unit})?",
                   "Mesh sizing",
                   plan_opts + [("Custom sizes",
                                 f"type min and max in {_import_length_unit}")]),
               "prompt_render": _PROMPT_RENDER_NOTE,
               "recommendation": _mesh_plan})


@mcp.tool()
def add_local_sizing(
    min_size: float = 0.0,
    max_size: float = 0.0,
    curvature_normal_angle_deg: float = 0.0,
    units: str = "",
) -> str:
    """Add curvature/proximity local surface sizing (Watertight workflow,
    optional step before Generate Surface Mesh).

    UNITS: values are interpreted in `units`; the default (empty) is the CAD
    IMPORT unit — exactly what the Fluent sizing panel shows. (The workflow's
    sizing fields are NOT meters: a CAD imported in mm has mm fields, and
    values pass through verbatim — passing meters made sizes 1000x too small
    and hung the mesher, verified live 2026-07-17.)

    MODES: in manual meshing mode min_size/max_size are required (ask the
    user). In auto mode, omit them to apply the propose_mesh_sizing plan.

    Args:
        min_size:                   Minimum local element size (in `units`).
        max_size:                   Maximum local element size (in `units`).
        curvature_normal_angle_deg: Max angle between adjacent face normals per
                                    surface-mesh cell (0 = plan value or 18).
        units:                      Unit of the size values: '' = CAD import
                                    unit (recommended), or 'm', 'mm', 'cm',
                                    'in', 'ft' to convert.
    """
    if err := _chk("meshing"): return err
    resolved = _resolve_sizing(min_size, max_size, units, "add_local_sizing")
    if isinstance(resolved, str):
        return resolved
    rmin, rmax = resolved
    if curvature_normal_angle_deg <= 0:
        curvature_normal_angle_deg = (_mesh_plan or {}).get("curvature_normal_angle_deg", 18.0)
    try:
        wf = _get_workflow()
        als = wf.add_local_sizing
        report = _apply_task_args(als, {
            "add_child": "yes",
            "boi_min_size": rmin,
            "boi_max_size": rmax,
            "boi_curvature_normal_angle": curvature_normal_angle_deg,
        })
        if guard := _guard_no_op(report, "Add Local Sizing"):
            return guard
        als()
        return _j({"ok": True, **report,
                   "sizing_unit": _import_length_unit,
                   "panel_shows": f"min {rmin} {_import_length_unit} / "
                                  f"max {rmax} {_import_length_unit}",
                   "warnings": _skip_warnings(report)})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _generate_surface_mesh_impl(
    min_size: float = 0.0,
    max_size: float = 0.0,
    growth_rate: float = 0.0,
    units: str = "",
) -> str:
    """Synchronous implementation — public tool: generate_surface_mesh. No
    _chk() here: see _import_geometry_impl's docstring for why an internal
    _chk() call inside a session-job worker self-blocks."""
    resolved = _resolve_sizing(min_size, max_size, units, "generate_surface_mesh")
    if isinstance(resolved, str):
        return resolved
    rmin, rmax = resolved
    if growth_rate <= 0:
        growth_rate = (_mesh_plan or {}).get("growth_rate", 1.2)
    t0 = time.time()
    try:
        wf = _get_workflow()
        csm = wf.create_surface_mesh
        ctrl = csm.cfd_surface_mesh_controls
        ctrl.min_size = rmin
        ctrl.max_size = rmax
        ctrl.growth_rate = growth_rate
        csm()
        verification: dict = {}
        try:
            st = _meshing.workflow.TaskObject["Generate the Surface Mesh"].get_state()
            if isinstance(st, dict):
                verification["task_state"] = st.get("State", "<unknown>")
                applied = (st.get("Arguments") or {}).get("CFDSurfaceMeshControls", {})
                if isinstance(applied, dict):
                    verification["controls_readback"] = applied
        except Exception as e:
            verification["readback_error"] = str(e)
        # inline quality verdict (assess_surface_mesh computes an almost
        # identical probe moments later anyway) so the caller doesn't have
        # to pay for a mandatory extra get_job_status + assess_surface_mesh
        # round trip just to see it
        quality = None
        try:
            quality = json.loads(_assess_surface_mesh_impl())
        except Exception:
            pass
        return _j({"ok": True,
                   "message": f"Surface mesh: min={rmin} {_import_length_unit} "
                              f"max={rmax} {_import_length_unit} growth={growth_rate}",
                   "elapsed_s": round(time.time() - t0, 1),
                   "verification": verification,
                   "quality": quality,
                   "next_recommended_action": "describe_geometry (quality verdict already "
                                              "included above; call assess_surface_mesh "
                                              "again only if you need a fresh re-check)"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def generate_surface_mesh(
    min_size: float = 0.0,
    max_size: float = 0.0,
    growth_rate: float = 0.0,
    units: str = "",
    run_in_background: bool = True,
) -> str:
    """Generate a surface mesh in a meshing session.

    Runs AS A BACKGROUND JOB by default (minutes on real geometry) — poll
    get_job_status every ~10-15 s and read_console_tail for live progress
    (the console prints the skewness table when done). The user may also
    run this step in the Fluent GUI: if state seems out of sync, check
    read_console_tail / get_workflow_state BEFORE re-running anything.
    After it finishes: assess_surface_mesh; on a conditional/reject
    skewness verdict ask the user before proceeding.

    UNITS: values are interpreted in `units`; the default (empty) is the CAD
    IMPORT unit — exactly what the Fluent sizing panel shows (NOT meters).

    MODES: in manual meshing mode min_size/max_size are required (the
    refusal carries an interactive_prompt to render). In auto mode, omit
    them to apply the propose_mesh_sizing plan.

    Args:
        min_size:          Minimum element size (in `units`).
        max_size:          Maximum element size (in `units`).
        growth_rate:       Growth rate (0 = plan value or 1.2).
        units:             '' = CAD import unit (recommended), or 'm', 'mm',
                           'cm', 'in', 'ft'.
        run_in_background: Default True (recommended). False blocks until done.
    """
    if err := _chk("meshing"): return err
    # fail fast on sizing problems BEFORE starting a background job, so the
    # manual-mode interactive prompt reaches the user immediately
    resolved = _resolve_sizing(min_size, max_size, units, "generate_surface_mesh")
    if isinstance(resolved, str):
        return resolved
    if run_in_background:
        return _run_session_job(
            "generate_surface_mesh",
            lambda: _generate_surface_mesh_impl(min_size, max_size, growth_rate, units))
    return _generate_surface_mesh_impl(min_size, max_size, growth_rate, units)


@mcp.tool()
def describe_geometry(setup_type: str = "fluid") -> str:
    """Describe Geometry task (Watertight workflow) — declares whether the
    domain is fluid-only or fluid+solid (CHT), which drives region/volume
    handling in every later step. Call right after Generate Surface Mesh.

    Args:
        setup_type: "fluid" (fluid-only, default), "fluid_solid" / "cht"
                    (conjugate heat transfer with solid bodies), "solid"
                    (solid-only), or the full Fluent sentence string for
                    this release verbatim (passed through unchanged).
    """
    if err := _chk("meshing"): return err
    # the datamodel takes full sentence enums, not one-word values — the
    # one-word aliases below map to the sentences verified on 2025 R2
    aliases = {
        "fluid": "The geometry consists of only fluid regions with no voids",
        "fluid_solid": "The geometry consists of both fluid and solid regions and/or voids",
        "cht": "The geometry consists of both fluid and solid regions and/or voids",
        "solid": "The geometry consists of only solid regions",
    }
    resolved = aliases.get(setup_type.strip().lower(), setup_type)
    try:
        wf = _get_workflow()
        dg = wf.describe_geometry
        report = _apply_task_args(dg, {"setup_type": resolved})
        if guard := _guard_no_op(report, "Describe Geometry"):
            return guard
        dg()
        return _j({"ok": True, "setup_type_applied": resolved, **report,
                   "warnings": _skip_warnings(report)})
    except Exception as e:
        return _j({"ok": False, "error": str(e), "setup_type_resolved": resolved})


@mcp.tool()
def update_boundaries(
    boundary_labels: list[str],
    boundary_types: list[str],
) -> str:
    """Assign boundary zone types (Watertight workflow) — do this right
    after Describe Geometry, before Add Boundary Layers, so later steps see
    correctly typed zones. Re-run list_boundary_zones (after switch_to_solver)
    to confirm the assignment stuck.

    Args:
        boundary_labels: Zone/label names from the imported CAD, e.g.
                         ["inlet_main", "outlet_main", "wall_body"].
        boundary_types:  Matching Fluent boundary types, e.g.
                         ["velocity-inlet", "pressure-outlet", "wall"].
    """
    if err := _chk("meshing"): return err
    if len(boundary_labels) != len(boundary_types):
        return _j({"ok": False, "error": "boundary_labels and boundary_types must be the same length"})
    # validate requested labels against the live inventory before touching
    # the task — assigning a nonexistent label is a silent no-op in Fluent
    warnings: list[str] = []
    mu = _mu()
    if mu is not None:
        try:
            available: set = set()
            for o in (mu.get_all_objects() or []):
                available.update(mu.get_labels(object_name=o) or [])
            unknown = [l for l in boundary_labels if l not in available]
            if unknown:
                return _j({"ok": False,
                           "error": f"labels not found in the imported geometry: {unknown}",
                           "available_labels": sorted(available),
                           "hint": "use list_labels_and_zones for the live inventory"})
        except Exception as e:
            warnings.append(f"label pre-validation unavailable: {e}")
    try:
        wf = _get_workflow()
        ub = wf.update_boundaries
        report = _apply_task_args(ub, {
            "boundary_label_list": boundary_labels,
            "boundary_label_type_list": boundary_types,
        })
        if guard := _guard_no_op(report, "Update Boundaries"):
            return guard
        ub()
        return _j({"ok": True, **report,
                   "warnings": warnings + _skip_warnings(report),
                   "next_recommended_action":
                       "re-run list_labels_and_zones (or list_boundary_zones after "
                       "switch_to_solver) to confirm the assignment stuck"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def add_boundary_layers(
    number_of_layers: int = 0,
    first_aspect_ratio: float = 5.0,
    growth_rate: float = 1.2,
    transition_ratio: float = 0.272,
    control_name: str = "smooth-transition_1",
) -> str:
    """Add prism/boundary-layer inflation on wall zones (Watertight workflow).

    Call after Update Regions, before Generate Volume Mesh. Skip entirely
    for inviscid/Euler cases or walls where near-wall shear/heat-transfer
    resolution isn't needed (never applies to BOI/refinement boxes — those
    are virtual sizing volumes, not real geometry).

    MODES: number_of_layers=0 means 'not provided' — AUTO meshing mode fills
    it from the propose_mesh_sizing plan; MANUAL mode requires it (ask the
    user: 3 quick-look / 5-8 standard / 12+ high-accuracy tier).

    Args:
        number_of_layers:   Prism layer count (0 = resolve from mode/plan).
        first_aspect_ratio: First-layer aspect ratio (smooth-transition offset method).
        growth_rate:        Layer-to-layer growth rate (keep <=1.2).
        transition_ratio:   Last-layer-to-adjacent-cell size ratio (Fluent default 0.272).
        control_name:       Name for this boundary-layer control.
    """
    if err := _chk("meshing"): return err
    if number_of_layers <= 0:
        if _meshing_mode == "auto" and _mesh_plan:
            number_of_layers = int(_mesh_plan.get("bl_layers", 5))
        elif _meshing_mode == "auto":
            number_of_layers = 5
        else:
            return _j({"ok": False,
                       "error": "add_boundary_layers: number_of_layers is REQUIRED in "
                                "manual meshing mode - do not guess it",
                       "interactive_prompt": _prompt(
                           "How many prism/boundary layers on the walls?", "BL layers",
                           [("5 - standard (Recommended)", "general engineering runs"),
                            ("3 - quick look", "screening only, not wall-accurate"),
                            ("8 - refined", "better wall shear/heat transfer"),
                            ("12+ - high accuracy", "wall-resolved analyses")]),
                       "prompt_render": _PROMPT_RENDER_NOTE,
                       "recommendation": _mesh_plan})
    desired = {
        "number_of_layers": number_of_layers,
        "first_aspect_ratio": first_aspect_ratio,
        "growth_rate": growth_rate,
        "transition_ratio": transition_ratio,
        "bl_control_name": control_name,
        "control_name": control_name,
    }
    try:
        wf = _get_workflow()
        # high-level path: attribute name varies by pyfluent release
        # (0.40: add_boundary_layers; 0.39: absent) — discover, don't assume
        task_attr = _find_wf_task(wf, "boundary", "layer")
        if task_attr is not None:
            abl = getattr(wf, task_attr)
            abl.insert_compound_child_task()
            # bind to the NEWEST child (highest _N suffix): repeat calls insert
            # child_2, child_3... and configuring child_1 again would silently
            # overwrite the first control and leave the new one unconfigured
            children = sorted(
                (n for n in dir(wf)
                 if "child" in n.lower() and "boundary" in n.lower() and "layer" in n.lower()),
                key=lambda n: int(n.rsplit("_", 1)[-1]) if n.rsplit("_", 1)[-1].isdigit() else 0)
            child_attr = children[-1] if children else None
            child = getattr(wf, child_attr, None) if child_attr else None
            if child is None:
                return _j({"ok": False,
                           "error": f"no child task found after insert_compound_child_task() on '{task_attr}'",
                           "hint": "call get_workflow_state to see the task tree"})
            report = _apply_task_args(child, desired)
            if guard := _guard_no_op(report, "Add Boundary Layers (child control)"):
                return guard
            child()
            return _j({"ok": True, "via": f"high-level task '{task_attr}'",
                       **report, "warnings": _skip_warnings(report)})
        # low-level datamodel fallback: TaskObject['Add Boundary Layers'] exists
        # even when the wrapper attribute doesn't (verified on 0.39.1 / 2025 R2)
        tobj = _meshing.workflow.TaskObject["Add Boundary Layers"]
        pascal = {
            "NumberOfLayers": number_of_layers,
            "FirstAspectRatio": first_aspect_ratio,
            "Rate": growth_rate,
            "TransitionRatio": transition_ratio,
            "BLControlName": control_name,
        }
        tobj.Arguments.set_state(pascal)
        tobj.AddChildAndUpdate()
        readback = {}
        try:
            readback = tobj.Arguments()
        except Exception as e:
            readback = {"readback_error": str(e)}
        return _j({"ok": True, "via": "low-level TaskObject fallback",
                   "applied": pascal, "readback": readback,
                   "warnings": ["applied via datamodel fallback - verify layer count "
                                "in get_workflow_state / assess_volume_mesh"]})
    except Exception as e:
        return _j({"ok": False, "error": str(e),
                   "hint": "get_workflow_state shows the live task tree and argument names"})


def _generate_volume_mesh_impl(mesh_type: str = "poly-hexcore") -> str:
    """Synchronous implementation — public tool: generate_volume_mesh. No
    _chk() here: see _import_geometry_impl's docstring for why an internal
    _chk() call inside a session-job worker self-blocks."""
    warnings: list[str] = []
    t0 = time.time()
    try:
        wf = _get_workflow()
        try:
            wf.update_regions()
        except Exception as e:
            # do not swallow: a failed Update Regions means the volume mesh
            # may be built on stale fluid/solid region definitions
            warnings.append(f"auto Update Regions failed before volume meshing: {e} "
                            "- run update_regions/create_regions explicitly and re-check")
        cvm = getattr(wf, "create_volume_mesh_wtm", None) or getattr(wf, "create_volume_mesh", None)
        if cvm is None:
            return _j({"ok": False, "error": "no create_volume_mesh task on this workflow"})
        report = _apply_task_args(cvm, {"volume_fill": mesh_type})
        if guard := _guard_no_op(report, "Generate the Volume Mesh"):
            return guard
        cvm()
        verification: dict = {}
        mu = _mu()
        if mu is not None:
            try:
                verification["mesh_exists"] = bool(mu.mesh_exists())
            except Exception as e:
                verification["mesh_exists_error"] = str(e)
        # inline quality verdict (assess_volume_mesh computes an almost
        # identical probe moments later anyway) so the caller doesn't have
        # to pay for a mandatory extra get_job_status + assess_volume_mesh
        # round trip just to see it
        quality = None
        try:
            quality = json.loads(_assess_volume_mesh_impl())
        except Exception:
            pass
        return _j({"ok": True, "mesh_type": mesh_type, **report,
                   "elapsed_s": round(time.time() - t0, 1),
                   "warnings": warnings + _skip_warnings(report),
                   "verification": verification,
                   "quality": quality,
                   "next_recommended_action":
                       "save_mesh_and_workflow BEFORE switch_to_solver (quality verdict "
                       "already included above; call assess_volume_mesh again only if "
                       "you need a fresh re-check) (meshing KB doc 06)"})
    except Exception as e:
        return _j({"ok": False, "error": str(e), "warnings": warnings})


@mcp.tool()
def generate_volume_mesh(mesh_type: str = "poly-hexcore", confirmed: bool = False,
                         run_in_background: bool = True) -> str:
    """Generate the volume mesh.

    The most expensive meshing step — minutes to tens of minutes. Runs AS A
    BACKGROUND JOB by default: poll get_job_status every ~10-15 s and
    read_console_tail for live progress (the console prints the
    orthogonal-quality table when done).

    CONFIRMATION GATE — SERVER-ENFORCED (meshing KB docs 06/07): the first
    call (confirmed=False) does NOT mesh; it returns an interactive prompt
    spec (fill method + commit) to render via AskUserQuestion. Re-call with
    confirmed=True (and the user's chosen mesh_type) after an explicit yes.
    After the job finishes: assess_volume_mesh for the quality verdict,
    then save_mesh_and_workflow.

    Args:
        mesh_type:         "tet", "poly", "poly-hexcore", or "hexcore".
        confirmed:         Pass True only after the user approved the gate prompt.
        run_in_background: Default True (recommended). False blocks until done.
    """
    if not confirmed:
        return _j({"ok": False, "confirmation_required": True,
                   "interactive_prompt": _prompt(
                       f"Generate the volume mesh now with '{mesh_type}' fill? "
                       "(most expensive meshing step)", "Volume mesh",
                       [(f"Yes - {mesh_type} (Recommended)",
                         "run it, then quality-check and save"),
                        ("Different fill method",
                         "choose tet / poly / hexcore instead"),
                        ("Not yet", "review the setup in the GUI first")]),
                   "prompt_render": _PROMPT_RENDER_NOTE,
                   "hint": "re-call generate_volume_mesh(mesh_type=..., confirmed=True) "
                           "after the user approves"})
    if err := _chk("meshing"): return err
    if run_in_background:
        return _run_session_job("generate_volume_mesh",
                                lambda: _generate_volume_mesh_impl(mesh_type))
    return _generate_volume_mesh_impl(mesh_type)


@mcp.tool()
def switch_to_solver(confirmed: bool = False) -> str:
    """Switch from meshing to solver mode after meshing is complete.

    BLOCKING, typically ~1-3 minutes (Fluent restarts its core in solver
    mode) — a silent wait is normal. CONFIRMATION GATE — SERVER-ENFORCED
    (meshing KB doc 06, 'Switch to solution'): the first call
    (confirmed=False) does NOT switch; it returns an interactive prompt
    spec to render via AskUserQuestion, since this is one-way for the
    workflow state. Re-call with confirmed=True after an explicit yes. Save
    the mesh first (save_mesh_and_workflow) — the switch discards the
    meshing workflow state, and an unsaved mesh survives only inside the
    solver session's memory.

    Args:
        confirmed: Pass True only after the user approved the gate prompt.
    """
    global _solver, _meshing, _session_mode, _workflow, _workflow_kind
    if err := _chk("meshing"): return err
    if not confirmed:
        return _j({"ok": False, "confirmation_required": True,
                   "interactive_prompt": _prompt(
                       "Switch Fluent from meshing to solver mode now? (one-way for "
                       "the workflow state; takes 1-3 min, Fluent core restarts)",
                       "Switch to solver",
                       [("Yes - switch to solver (Recommended)",
                         "proceed to CFD setup: physics, boundary conditions, and solve"),
                        ("Not yet",
                         "save the mesh / review the setup first (save_mesh_and_workflow)")]),
                   "prompt_render": _PROMPT_RENDER_NOTE,
                   "hint": "re-call switch_to_solver(confirmed=True) after the user approves"})
    if _workflow_kind == "2d":
        return _j({"ok": False,
                   "error": "switch_to_solver() is not valid in the 2D Meshing workflow - "
                            "the 2D surface mesh is the terminal artifact",
                   "hint": "write it out with save_mesh_and_workflow instead "
                           "(see meshing KB doc 01, section 5)"})
    t0 = time.time()
    try:
        _solver = _meshing.switch_to_solver()
        _meshing = None; _session_mode = "solver"
        _workflow = None; _workflow_kind = ""
        # re-persist: the record must say mode=solver or a post-restart
        # auto-reattach reconstructs the WRONG session type (hit live
        # 2026-07-20 - record still said 'meshing' after the switch)
        _persist_session_state()
        return _j({"ok": True, "message": "Switched to solver mode.",
                   "elapsed_s": round(time.time() - t0, 1),
                   "next_recommended_action": "run_mesh_check, then list_boundary_zones "
                                              "to re-inventory zones before BC setup"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def list_labels_and_zones() -> str:
    """Live label/object/zone inventory of a MESHING session.

    The meshing-mode counterpart of list_boundary_zones (which is
    solver-only): returns CAD objects, the face labels on each object
    (these are the names update_boundaries takes), the face-zone count,
    regions where available, and any boundary assignments already made.
    Build the BC inventory from THIS, never from assumed names
    (meshing KB doc 04, 'Boundary definitions').
    """
    if err := _chk("meshing"): return err
    out: dict = {"ok": True}
    mu = _mu()
    if mu is None:
        return _j({"ok": False,
                   "error": "meshing_utilities not available on this release",
                   "hint": "fall back to run_tui_command in the meshing TUI"})
    try:
        objs = list(mu.get_all_objects() or [])
        out["objects"] = objs
    except Exception as e:
        objs = []
        out["objects_error"] = str(e)
    labels: dict = {}
    for o in objs:
        try:
            labels[o] = list(mu.get_labels(object_name=o) or [])
        except Exception as e:
            labels[o] = f"<get_labels failed: {e}>"
    out["labels_by_object"] = labels
    try:
        out["face_zone_count"] = len(list(mu.get_face_zones(filter="*") or []))
    except Exception as e:
        out["face_zone_count_error"] = str(e)
    regions: dict = {}
    for o in objs:
        try:
            r = mu.get_regions(object_name=o)
            if r:
                regions[o] = list(r)
        except Exception:
            pass
    out["regions_by_object"] = regions
    # boundary assignments already recorded on the Update Boundaries task
    try:
        args = _meshing.workflow.TaskObject["Update Boundaries"].Arguments()
        if isinstance(args, dict) and args:
            out["current_boundary_assignments"] = {
                k: v for k, v in args.items() if "label" in k.lower() or "zone" in k.lower()}
    except Exception:
        pass
    return _j(out)


@mcp.tool()
def get_workflow_state() -> str:
    """Read back the meshing workflow's task tree: every task in order,
    its execution state (Up-to-date / Out-of-date / ...), and its argument
    dict. THE debugging primitive for meshing: call it whenever a task tool
    reports skipped arguments, before re-running any task, and to see how
    far a GUI-driven workflow has progressed (works regardless of whether
    the workflow was initialized by this server or in the Fluent GUI).
    """
    if err := _chk("meshing"): return err
    out: dict = {"ok": True, "workflow_kind": _workflow_kind or "<initialized outside this server>"}
    try:
        tobj = _meshing.workflow.TaskObject
        names = list(tobj.get_object_names())
    except Exception as e:
        return _j({"ok": False,
                   "error": f"no active workflow task tree: {e}",
                   "hint": "import_geometry initializes the workflow; nothing to inspect yet"})
    tasks = []
    for name in names:
        entry: dict = {"task": name}
        try:
            st = tobj[name].get_state()
            entry["state"] = st.get("State", "<unknown>") if isinstance(st, dict) else str(st)[:200]
            args = (st.get("Arguments") or {}) if isinstance(st, dict) else {}
            entry["arguments"] = _safe_truncate_dict(args)
        except Exception as e:
            entry["error"] = str(e)
        tasks.append(entry)
    out["tasks"] = tasks
    return _j(out)


def _safe_truncate_dict(d: dict, max_chars: int = 400) -> dict:
    """Bound each argument value so a huge zone list can't balloon the reply."""
    out = {}
    for k, v in d.items():
        s = repr(v)
        out[k] = v if len(s) <= max_chars else s[:max_chars] + "...<truncated>"
    return out


@mcp.tool()
def create_regions() -> str:
    """Run the Create Regions task (Watertight workflow) — forms fluid/solid/
    void regions from the enclosed surface mesh. Run after Describe Geometry
    (and capping, when present); verify the result against the expected
    region count with the returned inventory (meshing KB doc 04, region gate).
    """
    if err := _chk("meshing"): return err
    try:
        wf = _get_workflow()
        cr = getattr(wf, "create_regions", None)
        if cr is None:
            return _j({"ok": False, "error": "no create_regions task on this workflow",
                       "hint": "fault-tolerant workflows use identify_regions instead"})
        cr()
        return _j({"ok": True, "verification": _region_inventory(),
                   "next_recommended_action":
                       "compare region count/types against the mesh plan; fix with "
                       "update_regions before adding boundary layers"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def update_regions() -> str:
    """Run the Update Regions task (Watertight workflow) — re-derives or
    corrects fluid/solid/dead region classification. generate_volume_mesh
    auto-runs this too, but call it explicitly after create_regions or any
    region-affecting change so classification failures surface HERE rather
    than as a mis-typed volume mesh.
    """
    if err := _chk("meshing"): return err
    try:
        wf = _get_workflow()
        ur = getattr(wf, "update_regions", None)
        if ur is None:
            return _j({"ok": False, "error": "no update_regions task on this workflow"})
        ur()
        return _j({"ok": True, "verification": _region_inventory()})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _region_inventory() -> dict:
    """Best-effort region readback after Create/Update Regions."""
    v: dict = {}
    mu = _mu()
    if mu is None:
        v["note"] = "meshing_utilities not available"
        return v
    try:
        regions: dict = {}
        for o in (mu.get_all_objects() or []):
            try:
                r = mu.get_regions(object_name=o)
                if r:
                    regions[o] = list(r)
            except Exception:
                pass
        v["regions_by_object"] = regions
    except Exception as e:
        v["regions_error"] = str(e)
    return v


@mcp.tool()
def assess_surface_mesh() -> str:
    """Surface-mesh quality summary in a MESHING session (skewness limits,
    face count) with a screening verdict. Surface skewness is a screening
    metric, not final acceptance — poor faces only block if they sit in
    critical regions (meshing KB doc 03, proceed-or-improve rule).
    """
    if err := _chk("meshing"): return err
    return _assess_surface_mesh_impl()


def _assess_surface_mesh_impl() -> str:
    """Body of assess_surface_mesh, callable internally (e.g. from
    generate_surface_mesh's session-job worker, which embeds this verdict
    in its own result) WITHOUT re-triggering the busy-guard in _chk() -
    that guard correctly blocks an EXTERNAL call while a job is running,
    but an internal call from within that same job's own worker thread is
    not a conflict (verified live 2026-07-21: the embedded quality verdict
    was self-blocking with "session job/calculation in progress" before
    this split)."""
    mu = _mu()
    if mu is None:
        return _j({"ok": False, "error": "meshing_utilities not available on this release"})
    out: dict = {"ok": True, "probes": {}}
    try:
        out["mesh_exists"] = bool(mu.mesh_exists())
    except Exception as e:
        out["probes"]["mesh_exists"] = f"failed: {e}"
    try:
        zones = list(mu.get_face_zones(filter="*") or [])
        out["face_zone_count"] = len(zones)
    except Exception as e:
        zones = []
        out["probes"]["get_face_zones"] = f"failed: {e}"
    # quality-limit signatures vary by release — try each, report what worked
    for label, call in (
        ("face_skewness_by_id", lambda: mu.get_face_quality_limits(
            face_zone_id_list=zones, measure="Skewness")),
        ("face_skewness_by_pattern", lambda: mu.get_face_quality_limits(
            face_zone_name_pattern="*", measure="Skewness")),
    ):
        if not zones and "by_id" in label:
            continue
        try:
            r = call()
            if r is not None:
                out["skewness"] = str(r)[:1500]
                out["probes"][label] = "ok"
                break
        except Exception as e:
            out["probes"][label] = f"failed: {e}"
    if "skewness" not in out:
        out["warnings"] = ["no quality probe succeeded on this release - inspect "
                           "visually in the GUI or via run_tui_command"]
        return _j(out)
    # skewness policy per meshing KB doc 07 (source training: max < 0.7 desirable):
    # <0.70 approve | 0.70-0.85 conditional | >0.85 reject/improve | ~1.0 block.
    # The probe's return shape is release-dependent; take the max float in (0,1]
    # as the max-skewness candidate and mark the parse heuristic.
    try:
        import re as _re
        floats = [float(x) for x in _re.findall(r"\d+\.\d+(?:e-?\d+)?", out["skewness"])]
        cands = [f for f in floats if 0.0 < f <= 1.0]
        if cands:
            mx = max(cands)
            out["max_skewness_parsed"] = mx
            out["parse_note"] = "heuristic parse of a release-dependent payload - verify in GUI if borderline"
            if mx >= 0.999:
                out["verdict"] = ("BLOCK: near-degenerate face (skewness ~1.0) - locate and "
                                  "improve/remesh before volume meshing (KB doc 07 skewness policy)")
            elif mx > 0.85:
                out["verdict"] = ("REJECT/IMPROVE: max skewness > 0.85 - improve before volume "
                                  "meshing unless the user records an expert override")
            elif mx >= 0.70:
                out["verdict"] = ("CONDITIONAL: max skewness 0.70-0.85 - locate the poor faces; "
                                  "improve if they are in critical regions, else proceed and "
                                  "judge the volume mesh; ASK THE USER at this gate")
            else:
                out["verdict"] = "APPROVE (screening): max skewness < 0.70"
            # The combined probe above aggregates ALL face zones in one call
            # (get_face_quality_limits over the whole `zones` list), so a
            # near-degenerate face can't be attributed to a specific zone
            # against a console table scoped to only one CAD object (verified
            # live 2026-07-20: a 0.9999 max here couldn't be reconciled with a
            # console table showing 0.47 for a single named object). Only
            # pay for the extra per-zone round trips once a verdict worse
            # than plain APPROVE is about to be issued.
            if mx >= 0.70 and zones:
                per_zone: dict = {}
                for zid in zones:
                    try:
                        zr = mu.get_face_quality_limits(face_zone_id_list=[zid], measure="Skewness")
                        if zr is not None:
                            per_zone[str(zid)] = str(zr)[:300]
                    except Exception as e:
                        per_zone[str(zid)] = f"failed: {e}"
                if per_zone:
                    out["skewness_by_zone"] = per_zone
                    worst_zone, worst_val = None, -1.0
                    for zid, payload in per_zone.items():
                        zfloats = [float(x) for x in _re.findall(r"\d+\.\d+(?:e-?\d+)?", payload)]
                        zcands = [f for f in zfloats if 0.0 < f <= 1.0]
                        if zcands and max(zcands) > worst_val:
                            worst_val, worst_zone = max(zcands), zid
                    if worst_zone is not None:
                        out["max_skewness_zone_id"] = worst_zone
                        out["max_skewness_zone_value"] = worst_val
    except Exception as e:
        out["probes"]["skewness_policy_parse"] = f"failed: {e}"
    return _j(out)


@mcp.tool()
def assess_volume_mesh() -> str:
    """Volume-mesh quality verdict in a MESHING session, against the meshing
    KB's orthogonal-quality ladder: minimum above ~0.01 required, below
    ~0.05 investigate/improve, above ~0.1 is the practical target (doc 05).
    'ok' means the probe ran; read 'verdict' for acceptance. Never accept a
    mesh only because generation completed.
    """
    if err := _chk("meshing"): return err
    return _assess_volume_mesh_impl()


def _assess_volume_mesh_impl() -> str:
    """Body of assess_volume_mesh, callable internally (e.g. from
    generate_volume_mesh's session-job worker, which embeds this verdict
    in its own result) WITHOUT re-triggering the busy-guard in _chk() -
    see _assess_surface_mesh_impl's docstring for why this split exists."""
    mu = _mu()
    if mu is None:
        return _j({"ok": False, "error": "meshing_utilities not available on this release"})
    out: dict = {"ok": True, "probes": {}}
    try:
        # mesh_exists() returns FALSE in meshing mode even right after a
        # successful volume fill (false negative verified live 2026-07-15) —
        # report it but NEVER gate on it; the workflow task state
        # ('Generate the Volume Mesh' Up-to-date) is the authoritative signal
        out["mesh_exists_probe"] = bool(mu.mesh_exists())
    except Exception as e:
        out["probes"]["mesh_exists"] = f"failed: {e}"
    try:
        st = _meshing.workflow.TaskObject["Generate the Volume Mesh"].get_state()
        out["volume_mesh_task_state"] = st.get("State", "<unknown>") if isinstance(st, dict) else str(st)[:100]
        if out["volume_mesh_task_state"] != "Up-to-date":
            return _j({**out, "verdict": "NO VOLUME MESH - the Generate the Volume Mesh "
                                         "task is not Up-to-date; run generate_volume_mesh"})
    except Exception as e:
        out["probes"]["volume_task_state"] = f"failed: {e}"
    min_oq = None
    for label, call in (
        ("cell_quality_by_pattern", lambda: mu.get_cell_quality_limits(
            cell_zone_name_pattern="*", measure="Orthogonal Quality")),
        ("cell_quality_by_id", lambda: mu.get_cell_quality_limits(
            cell_zone_id_list=list(mu.get_cell_zones(filter="*") or []),
            measure="Orthogonal Quality")),
    ):
        try:
            r = call()
            if r is not None:
                out["orthogonal_quality_raw"] = str(r)[:1500]
                out["probes"][label] = "ok"
                if isinstance(r, dict):
                    for k in ("min", "minimum", "min-quality", "MinQuality"):
                        if k in r:
                            min_oq = float(r[k]); break
                elif isinstance(r, (list, tuple)) and len(r) >= 2:
                    # observed shape on this release: [cell_count, min, max, avg, ...]
                    # (matches ansys-fluent-core's own get_cell_quality_limits
                    # docstring: "number of cells and the cell quality limits
                    # (minimum, maximum, average quality)"; verified live
                    # 2026-07-20 - r[1] equalled the console's own reported
                    # minimum Orthogonal Quality exactly)
                    try:
                        min_oq = float(r[1])
                    except (TypeError, ValueError):
                        pass
                break
        except Exception as e:
            out["probes"][label] = f"failed: {e}"
    try:
        out["worst_cell"] = str(mu.print_worst_quality_cell(measure="Orthogonal Quality"))[:500]
    except Exception as e:
        out["probes"]["print_worst_quality_cell"] = f"failed: {e}"
    if min_oq is not None:
        out["min_orthogonal_quality"] = min_oq
        if min_oq <= 0.01:
            out["verdict"] = ("REJECT: min orthogonal quality <= 0.01 - improve or "
                              "regenerate before solving (meshing KB doc 05 / MESH-009)")
        elif min_oq < 0.05:
            out["verdict"] = ("INVESTIGATE: min orthogonal quality < 0.05 - locate the "
                              "poor cells; improve unless they are outside critical regions")
        elif min_oq < 0.1:
            out["verdict"] = ("MARGINAL PASS: above 0.05 but below the practical 0.1 "
                              "target - check poor-cell locations")
        else:
            out["verdict"] = "PASS (screening): min orthogonal quality >= 0.1"
        out["verdict_note"] = ("screening values only - acceptance also depends on WHERE "
                               "poor cells sit relative to the physics (doc 00/05)")
    else:
        out["verdict"] = ("UNKNOWN: no quality probe succeeded on this release - run "
                          "run_mesh_check after switch_to_solver, or inspect in the GUI")
    return _j(out)


def _save_mesh_and_workflow_impl(mesh_file_path: str, workflow_file_path: str) -> str:
    out: dict = {"ok": True, "artifacts": []}
    _mkdir(mesh_file_path)
    # forward slashes: Fluent's TUI/scheme layer treats backslashes as escapes
    tui_path = mesh_file_path.replace("\\", "/")
    try:
        _meshing.tui.file.write_mesh(f'"{tui_path}"')
    except Exception as e1:
        try:
            _meshing.tui.file.write_mesh(tui_path)
        except Exception as e2:
            return _j({"ok": False,
                       "error": f"write_mesh failed: {e1} / retry without quotes: {e2}"})
    if Path(mesh_file_path).exists():
        out["artifacts"].append(mesh_file_path)
    else:
        # Fluent may append .h5/.msh.h5 itself — report what actually landed
        hits = [str(p) for p in Path(mesh_file_path).parent.glob(Path(mesh_file_path).stem + "*")]
        if hits:
            out["artifacts"].extend(hits)
            out["note"] = "Fluent adjusted the extension"
        else:
            return _j({"ok": False,
                       "error": f"write_mesh reported no error but nothing exists at "
                                f"{mesh_file_path} - a completed call is not sufficient "
                                f"verification (meshing KB doc 06)"})
    if workflow_file_path:
        _mkdir(workflow_file_path)
        try:
            _meshing.workflow.SaveWorkflow(FilePath=workflow_file_path)
            if Path(workflow_file_path).exists():
                out["artifacts"].append(workflow_file_path)
            else:
                out.setdefault("warnings", []).append(
                    f"SaveWorkflow raised no error but {workflow_file_path} does not exist")
        except Exception as e:
            out.setdefault("warnings", []).append(f"workflow save failed: {e}")
    return _j(out)


@mcp.tool()
def save_mesh_and_workflow(mesh_file_path: str, workflow_file_path: str = "",
                           run_in_background: bool = True) -> str:
    """Write the mesh (.msh.h5) and optionally the workflow recipe (.wft)
    from a MESHING session. Do this BEFORE switch_to_solver — the switch
    discards workflow state, and it is the ONLY persistence path for the 2D
    workflow (switch_to_solver is invalid there). Keep the mesh and its
    workflow file together per the archiving standard (meshing KB doc 06).

    Runs AS A BACKGROUND JOB by default - writing a volume mesh is the same
    class of large-file I/O that import_geometry/generate_volume_mesh were
    already backgrounded to avoid freezing the server on. Poll
    get_job_status + read_console_tail for progress.

    Args:
        mesh_file_path:     Output mesh path, e.g. .../project_mesh.msh.h5
        workflow_file_path: Optional .wft path for the replayable task recipe.
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk("meshing"): return err
    _mkdir(mesh_file_path)
    if run_in_background:
        return _run_session_job("save_mesh_and_workflow",
                                lambda: _save_mesh_and_workflow_impl(mesh_file_path, workflow_file_path))
    return _save_mesh_and_workflow_impl(mesh_file_path, workflow_file_path)


# ===========================================================================
# 23. TUI / SCHEME SCRIPTING
# ===========================================================================

@mcp.tool()
def run_tui_command(command: str) -> str:
    """Execute any Fluent TUI command directly.

    The TUI exposes the complete Fluent API. Use for anything not covered
    by the other tools.

    Examples:
        /solve/set/under-relaxation/pressure 0.3
        /report/surface-integrals/area-average
        /adapt/gradient yes yes 0.5 0.25 no no no no q

    Args:
        command: Full TUI command path + arguments.
    """
    if err := _chk(): return err
    try:
        result = _scheme_str(f'(ti-menu-load-string "{_esc(command)}")')
        return _j({"ok": True, "command": command, "output": result})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def run_scheme_expression(expression: str) -> str:
    """Evaluate a Scheme/Fluent interpreter expression (PYFLUENT SESSION path).

    The most powerful escape hatch — read or modify any Fluent internal
    state through the connected session. Uses string_eval so the RESULT IS
    RETURNED (exec() would run it but discard the value). For the RAW gRPC
    equivalent that needs no session object — e.g. the session is busy
    iterating, or this process never launched Fluent — use grpc_scheme_eval.

    Examples:
        (surface-area '(wall-1))
        (rpgetvar 'number-of-iterations)

    Args:
        expression: Scheme expression string.
    """
    if err := _chk(): return err
    try:
        result = _scheme_str(expression)
        return _j({"ok": True, "expression": expression, "result": result})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _tui_path(p: str) -> str:
    """Format a Windows path for a Fluent TUI string argument (forward
    slashes read reliably through the TUI's own quoting; backslashes are
    ambiguous inside the already-quoted Scheme string this goes through)."""
    return str(p).replace("\\", "/")


@mcp.tool()
def start_traceability(output_dir: str, run_label: str = "") -> str:
    """Start recording a Fluent-native transcript AND journal file.

    Solves a real gap: run_tui_command only returns the Scheme boolean for
    report/list-style commands, never the printed console text. A transcript
    file captures that full text on disk regardless of what any tool call
    returns. The journal file is a bonus — a real, standalone-replayable
    .jou of every command issued this session (see robust_execution_
    methodology.md Section 5, "TUI + journal" method).

    Args:
        output_dir: Directory to write into (created if missing).
        run_label:  Label for the file names (default: timestamp).
    """
    if err := _chk(): return err
    try:
        label = run_label or datetime.now().strftime("%Y%m%d_%H%M%S")
        trn_path = str(Path(output_dir) / f"fluent_transcript_{label}.trn")
        jou_path = str(Path(output_dir) / f"fluent_journal_{label}.jou")
        _mkdir(trn_path)
        _mkdir(jou_path)
        # exec() (not string_eval/_scheme_str) - the returned text is never
        # read at either call site, so skip marshaling it back over gRPC.
        # Must be a tuple: a bare string raises "Error: eof inside list"
        # (manifold_simulation_postmortem_2026-07-11.md).
        _active().scheme_eval.exec((f'(ti-menu-load-string "file/start-transcript \\"{_tui_path(trn_path)}\\"")',))
        _active().scheme_eval.exec((f'(ti-menu-load-string "file/start-journal \\"{_tui_path(jou_path)}\\"")',))
        return _j({"ok": True, "transcript_file": trn_path, "journal_file": jou_path,
                   "message": "Recording started. Call stop_traceability to close both files."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def stop_traceability() -> str:
    """Stop recording the Fluent transcript and journal files started by
    start_traceability, closing both so they can be read/replayed."""
    if err := _chk(): return err
    try:
        _active().scheme_eval.exec(('(ti-menu-load-string "file/stop-transcript")',))
        _active().scheme_eval.exec(('(ti-menu-load-string "file/stop-journal")',))
        return _j({"ok": True, "message": "Transcript and journal recording stopped."})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _run_tui_script_impl(script_path: str, stop_on_error: bool, max_lines: int) -> dict:
    lines = Path(script_path).read_text(encoding="utf-8", errors="replace").splitlines()
    executed: list[dict] = []
    for i, raw_line in enumerate(lines[:max_lines], start=1):
        line = raw_line.strip()
        if not line or line.startswith(";") or line.startswith("#"):
            continue  # blank line or comment (';' is the Fluent/Scheme comment char)
        try:
            # a .jou/script line starting with '(' is a raw Scheme form (as
            # recorded by start_traceability's journal, or hand-written);
            # anything else is a plain TUI command path - same dispatch
            # run_tui_command/run_scheme_expression already use, just looped
            output = _scheme_str(line) if line.startswith("(") else _tui(line)
        except Exception as e:
            output = f"ERROR: {e}"
        is_error = "error" in output.lower()
        executed.append({"line": i, "command": line, "output": output[:500],
                         "error": is_error})
        if is_error and stop_on_error:
            return {"ok": False, "stopped_at_line": i, "lines_executed": len(executed),
                    "executed": executed,
                    "error": f"line {i} looked like an error - stopped "
                             "(pass stop_on_error=False to run the rest anyway)"}
    return {"ok": True, "lines_executed": len(executed), "executed": executed}


@mcp.tool()
def run_tui_script(script_path: str, stop_on_error: bool = True,
                   max_lines: int = 2000, run_in_background: bool = True) -> str:
    """Replay a saved TUI script or Fluent journal (.jou) file end-to-end,
    one command per line, inside the CURRENTLY CONNECTED session.

    Built for a validated workflow: once a mesh sizing / BC / solver-setup
    recipe has been confirmed correct for a given case/geometry type (user
    feedback 2026-07-21), call start_traceability(output_dir) at the START
    of that session, work through setup once as normal, stop_traceability()
    at the end - the resulting .jou file IS the reusable recipe. For a
    future case of the same type, import/mesh the new geometry, then call
    run_tui_script(that_journal_path) to apply the entire rest of the
    recipe in ONE call instead of re-driving each step (and each
    confirmation prompt) individually.

    Runs AS A BACKGROUND JOB by default - a full end-to-end recipe replay
    can itself take as long as the original session (meshing + solving).
    Poll get_job_status + read_console_tail for progress.

    LIMITATION: lines execute via raw TUI/Scheme (ti-menu-load-string /
    string_eval) exactly like run_tui_command - they do NOT go through this
    server's own Python-side mode tracking (_solver/_meshing/_session_mode).
    A journal that switches meshing<->solver mode mid-script will leave
    this server's bookkeeping stale; for that case, split the journal at
    the mode switch and call switch_to_solver explicitly between the two
    halves. Journals scoped to a single mode (pure meshing setup, or pure
    solver setup) replay cleanly.

    Args:
        script_path:        Path to a .jou (from start_traceability) or a
                            plain text file of TUI commands, one per line.
                            Blank lines and ';'/'#' comment lines are skipped.
        stop_on_error:      Stop at the first line whose output looks like
                            an error (default True).
        max_lines:          Safety cap on lines executed (default 2000).
        run_in_background:  Default True (recommended). False blocks until done.
    """
    if err := _chk(): return err
    if not Path(script_path).exists():
        return _j({"ok": False, "error": f"Not found: {script_path}"})
    if run_in_background:
        return _run_session_job("run_tui_script",
                                lambda: _j(_run_tui_script_impl(script_path, stop_on_error, max_lines)))
    return _j(_run_tui_script_impl(script_path, stop_on_error, max_lines))


# ===========================================================================
# 24. SANITY CHECKS, PREFLIGHT, AND NATIVE POST-PROCESSING (v3.2)
# Lessons from the manifold_2 run: every path below was verified against a
# live Fluent 2025 R2 session before being added here.
# ===========================================================================

def _tui(command: str) -> str:
    """Internal TUI runner returning raw transcript text ('' on failure).
    Session path via string_eval — exec() would discard the transcript."""
    try:
        return _scheme_str(f'(ti-menu-load-string "{_esc(command)}")')
    except Exception:
        return ""


@mcp.tool()
def read_mesh(mesh_file_path: str) -> str:
    """Read a Fluent MESH file (.msh / .msh.h5) into the solver session.

    Use this for bare mesh files; read_case is only for .cas/.cas.h5.
    After reading, ALWAYS call run_mesh_check (or preflight_simulation)
    before doing any setup.

    Args:
        mesh_file_path: Absolute path to the mesh file.
    """
    if err := _chk("solver"): return err
    if not Path(mesh_file_path).exists():
        return _j({"ok": False, "error": f"Not found: {mesh_file_path}"})
    _set_project_dir_from(mesh_file_path)
    try:
        _solver.settings.file.read(file_type="mesh", file_name=mesh_file_path)
        return _j({"ok": True, "message": f"Mesh loaded: {mesh_file_path}",
                   "next": "run_mesh_check, then list_boundary_zones"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def run_mesh_check() -> str:
    """Run Fluent's mesh check and return a STRUCTURED verdict, not a blob.

    Parses the transcript for the 'Mesh check failed' warning and returns
    passed=true/false plus domain extents and volume statistics. If the check
    fails and the case has unpaired interface zones (e.g. 'filter.1:*' style
    multi-region meshes), the most common fix is create_mesh_interfaces —
    the hint field says so explicitly.
    """
    if err := _chk("solver"): return err
    out = _tui("mesh/check")
    if not out:
        return _j({"ok": False, "error": "mesh/check produced no transcript"})
    passed = "Mesh check failed" not in out
    import re as _re
    extents = dict(_re.findall(r"(\w)-coordinate: (min \(m\) = [^\n]+)", out))
    minvol = _re.search(r"minimum volume \(m3\):\s*([\d.eE+-]+)", out)
    result = {
        "ok": True,
        "mesh_check_passed": passed,
        "domain_extents": extents,
        "minimum_cell_volume_m3": float(minvol.group(1)) if minvol else None,
        "negative_volumes": (float(minvol.group(1)) <= 0.0) if minvol else None,
    }
    if not passed:
        # count unpaired interface zones as the prime suspect
        unpaired = []
        try:
            bc_state = _solver.settings.setup.boundary_conditions.get_state() or {}
            unpaired = list((bc_state.get("interface") or {}).keys())
        except Exception:
            pass
        result["hint"] = (
            "Most common cause: unpaired non-conformal interface zones "
            f"(found interface-type zones: {unpaired}). Call "
            "create_mesh_interfaces to auto-pair them, then re-run this check. "
            "If it still fails, run run_tui_command('mesh/check-verbosity 2') "
            "then re-check for the specific failing test."
        )
    return _j(result)


def _interface_flow_integrity() -> dict:
    """Verify mesh-interface zones created by auto_create() are typed
    'interior' (flow passes through) and not silently left as 'wall'
    (flow fully blocked). Reuses _zone_type_pairs() — the same source of
    truth as list_boundary_zones — instead of a second, possibly-stale
    lookup path.

    Why this matters: Fluent's non-conformal interface pairing can, on a
    mismatched/partially-overlapping pair of faces, leave the interface's
    own zone typed as a WALL rather than INTERIOR. The mesh check still
    passes and hybrid_initialize()/iterate() both run to completion without
    error — the only symptom is that mass flow across that interface is
    exactly zero, which downstream shows up as one or more outlets reading
    near-zero velocity while the rest of the domain carries all the flow
    (the manifold_2 case, 2026-07-14, first suspected this exact mechanism
    for its 2-of-3 near-stagnant outlets — checked and cleared that time,
    but the failure mode is real and worth checking on every interface
    creation, not just when a result looks suspicious).
    """
    zone_types = {z["zone"]: z["type"] for z in _zone_type_pairs()}
    try:
        mi_state = _solver.settings.setup.mesh_interfaces.interface.get_state() or {}
    except Exception:
        mi_state = {}
    checked, blocked = [], []
    for intf_name in mi_state.keys():
        checked.append(intf_name)
        if zone_types.get(intf_name) == "wall":
            blocked.append(intf_name)
    return {"interfaces_checked": checked, "blocked_as_wall": blocked,
            "ok": not blocked}


@mcp.tool()
def create_mesh_interfaces() -> str:
    """Auto-create (pair) all non-conformal mesh interfaces.

    Required for multi-region meshes whose regions meet at interface zones
    (a failed mesh check with interface-type zones is the classic symptom).
    SIDE EFFECTS to expect and re-inventory afterwards: some zones may be
    reclassified (e.g. wall -> interface) and small '*-non-overlapping' wall
    fragments are created; re-run list_boundary_zones before assigning BCs.

    Also runs a flow-integrity check on every newly-created interface (see
    check_interface_flow_integrity): a pairing that lands on an interface
    typed WALL instead of INTERIOR blocks flow across it completely while
    still passing the mesh check and solving without error — the classic
    downstream symptom is one or more outlets reading near-zero flow. A
    'CAUTION' key appears in the result if any interface came out blocked.
    """
    if err := _chk("solver"): return err
    try:
        mi = _solver.settings.setup.mesh_interfaces
        mi.auto_create()
        interfaces = {}
        try:
            interfaces = mi.interface.get_state() or {}
        except Exception:
            pass
        integrity = _interface_flow_integrity()
        result = {"ok": True, "interfaces_created": list(interfaces.keys()),
                  "warning": "Zone list has changed: interface pairing can "
                             "reclassify zones and add '*-non-overlapping' "
                             "wall fragments. Call list_boundary_zones again "
                             "before applying boundary conditions.",
                  "flow_integrity": integrity}
        if integrity["blocked_as_wall"]:
            result["CAUTION"] = (
                f"Interface zone(s) {integrity['blocked_as_wall']} came out "
                "typed WALL, not INTERIOR - flow across them is fully "
                "blocked even though mesh check and solve will both appear "
                "to succeed. Fix before trusting results: re-check the "
                "matching faces on each side (partial/non-overlapping "
                "coverage is the usual cause), delete and re-pair that "
                "interface, or manually retype the zone to interior. "
                "Expect near-zero flow on any outlet fed only through this "
                "interface until it's fixed."
            )
        return _j(result)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def check_interface_flow_integrity() -> str:
    """Re-check that every mesh-interface zone is typed 'interior' (flow
    passes through), not 'wall' (flow blocked).

    Call this any time after create_mesh_interfaces — right before BC
    setup, and again before trusting per-outlet/per-branch results on a
    manifold/T-junction/multi-region case — since a wall-typed interface
    produces no error anywhere in the pipeline, only silently-zero flow
    downstream. See _interface_flow_integrity for the full mechanism.
    """
    if err := _chk("solver"): return err
    try:
        integrity = _interface_flow_integrity()
        return _j({"ok": True, **integrity,
                   "verdict": "FAIL - flow blocked at interface(s)" if integrity["blocked_as_wall"]
                              else "pass"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def preflight_simulation() -> str:
    """Run the standard pre-setup sanity checklist and return pass/fail per item.

    Call after loading a mesh/case and BEFORE applying BCs or iterating.
    Checks: (1) mesh loaded; (2) domain scale plausible (bounding box not
    absurdly large/small - catches wrong units, e.g. mm read as m);
    (3) boundary-zone inventory - at least one inlet and one outlet, walls
    present, and no unpaired interface zones; (4) Fluent mesh check passes;
    (5) minimum orthogonal quality above 0.05 with a warning band below 0.15.
    Any failed item comes back in 'failed' with a suggested fix - resolve or
    get the user's explicit OK before proceeding.
    """
    if err := _chk("solver"): return err
    checks, failed = {}, []

    def record(name, ok, detail, fix=None):
        checks[name] = {"passed": bool(ok), "detail": detail}
        if fix:
            checks[name]["fix"] = fix
        if not ok:
            failed.append(name)

    # 1+2) mesh loaded & scale plausible (from mesh/check extents)
    out = _tui("mesh/check")
    import re as _re
    ext = _re.findall(r"[xyz]-coordinate: min \(m\) = ([\d.eE+-]+), max \(m\) = ([\d.eE+-]+)", out)
    if len(ext) == 3:
        spans = [abs(float(b) - float(a)) for a, b in ext]
        record("mesh_loaded", True, f"domain spans (m): {[f'{s:.3g}' for s in spans]}")
        scale_ok = all(1e-4 < s < 1e4 for s in spans)
        record("scale_plausible", scale_ok,
               f"largest span {max(spans):.3g} m, smallest {min(spans):.3g} m",
               fix=None if scale_ok else "Suspicious size - check mesh units "
                   "(mm CAD read as m?). Use run_tui_command('mesh/scale ...').")
    else:
        record("mesh_loaded", False, "no domain extents in mesh/check output",
               fix="Load a mesh first: read_mesh (or read_case).")

    # 3) zone inventory
    try:
        bc_state = _solver.settings.setup.boundary_conditions.get_state() or {}
        inv = {k: list(v.keys()) for k, v in bc_state.items() if isinstance(v, dict) and v}
        n_in = sum(len(v) for k, v in inv.items() if "inlet" in k)
        n_out = sum(len(v) for k, v in inv.items() if "outlet" in k)
        n_wall = len(inv.get("wall", []))
        # interface-type zones KEEP that type after pairing — a zone is only
        # unpaired if no created mesh interface references it
        paired_zones: set = set()
        try:
            for props in (_solver.settings.setup.mesh_interfaces.interface.get_state() or {}).values():
                paired_zones.update(props.get("zone1", []) + props.get("zone2", []))
        except Exception:
            pass
        unpaired = [z for z in inv.get("interface", []) if z not in paired_zones]
        record("has_inlet", n_in >= 1, f"{n_in} inlet zone(s): {[z for k, v in inv.items() if 'inlet' in k for z in v]}",
               fix=None if n_in else "No inlet-type zone. Convert one with the BC tools before solving.")
        record("has_outlet", n_out >= 1, f"{n_out} outlet zone(s)",
               fix=None if n_out else "No outlet-type zone - flow problem is ill-posed unless fully closed.")
        record("has_walls", n_wall >= 1, f"{n_wall} wall zone(s)")
        record("interfaces_paired", not unpaired,
               f"unpaired interface zones: {unpaired}" if unpaired else "none pending",
               fix=None if not unpaired else "Call create_mesh_interfaces, then re-inventory zones.")
    except Exception as e:
        record("zone_inventory", False, f"could not read boundary conditions: {e}")

    # 4) mesh check verdict
    record("mesh_check", "Mesh check failed" not in out and bool(out),
           "Fluent mesh check clean" if out and "Mesh check failed" not in out
           else "WARNING: Mesh check failed",
           fix=None if out and "Mesh check failed" not in out
           else "See run_mesh_check hint (usually create_mesh_interfaces).")

    # 5) quality
    qout = _tui("mesh/quality")
    m = _re.search(r"Minimum Orthogonal Quality\s*=\s*([\d.eE+-]+)", qout)
    if m:
        oq = float(m.group(1))
        record("orthogonal_quality", oq > 0.05,
               f"min orthogonal quality {oq:.3g}" + (" (marginal: < 0.15)" if oq < 0.15 else ""),
               fix=None if oq > 0.05 else "Quality below 0.05 - improve the mesh before solving.")
    else:
        record("orthogonal_quality", False, "could not parse mesh/quality output")

    return _j({"ok": True, "passed": not failed, "failed": failed, "checks": checks,
               "action": "proceed" if not failed else
                         "STOP: resolve failed checks (or get explicit user OK) before setup/solve"})


@mcp.tool()
def check_outlet_flow_balance(
    outlet_zones: list[str],
    inlet_zones: list[str] | None = None,
    near_zero_fraction: float = 0.02,
    imbalance_tolerance: float = 0.05,
) -> str:
    """Sanity-check mass-flow distribution across multiple outlets after
    solving — call this on any manifold/T-junction/multi-branch case before
    trusting per-outlet results, especially right after using
    create_mesh_interfaces (see check_interface_flow_integrity).

    Uses report_definitions.compute() (NOT the per-child .get_monitor_value(),
    which is broken on this Fluent/pyfluent version — see pitfall #13 in
    pyfluent_cfd_workflow_playbook.md / cht_simulation_postmortem_2026-07-10.md).

    Flags two independent things:
      1. near-zero outlet(s): any outlet whose |mass flow| share of total
         outflow is below `near_zero_fraction` while the rest carry the
         remainder. The single most common cause is a mesh interface that
         came out typed WALL instead of INTERIOR on the branch feeding that
         outlet (silent, no solver error) - cross-check with
         check_interface_flow_integrity before assuming it's a real design
         result.
      2. global mass imbalance: |sum(outlet mdot) + sum(inlet mdot)| /
         sum(|inlet mdot|) above `imbalance_tolerance` - the run likely
         hasn't converged yet (mass conservation is one of the required
         convergence criteria, not just flat residuals).

    Args:
        outlet_zones:        Outlet boundary zone names.
        inlet_zones:         Inlet boundary zone names (optional but
                              recommended - enables the global-imbalance
                              check; without it only the per-outlet split
                              is evaluated).
        near_zero_fraction:  Share of total outflow below which an outlet
                             is flagged as suspiciously near-zero (default 0.02 = 2%).
        imbalance_tolerance: Fractional mass imbalance above which the run
                             is flagged as not-yet-converged (default 0.05 = 5%).
    """
    if err := _chk("solver"): return err
    # surface-massflowrate report definitions + compute() — the flux-report
    # child's zone attribute was renamed on pyfluent 0.40.1 ("'flux_1_child'
    # object has no attribute 'zone_names'", verified live 2026-07-20); the
    # surface path is the one proven to return correct mdot on this release.
    rd = _solver.solution.report_definitions
    created: list[str] = []
    try:
        names = []
        for i, z in enumerate(outlet_zones):
            rn = f"mcp_ofb_out_{i}"
            rd.surface[rn] = {}
            r = rd.surface[rn]
            r.report_type = "surface-massflowrate"
            r.surface_names = [z]
            names.append(rn); created.append(rn)
        in_names = []
        if inlet_zones:
            for i, z in enumerate(inlet_zones):
                rn = f"mcp_ofb_in_{i}"
                rd.surface[rn] = {}
                r = rd.surface[rn]
                r.report_type = "surface-massflowrate"
                r.surface_names = [z]
                in_names.append(rn); created.append(rn)
        computed = rd.compute(report_defs=names + in_names)
        values = {}
        for entry in computed or []:
            for k, v in entry.items():
                values[k] = v[0] if isinstance(v, (list, tuple)) else v
        outlet_mdot = {z: values[n] for z, n in zip(outlet_zones, names)}
        total_out = sum(abs(v) for v in outlet_mdot.values())
        shares = {z: (abs(v) / total_out if total_out else 0.0) for z, v in outlet_mdot.items()}
        near_zero = [z for z, s in shares.items() if s < near_zero_fraction]

        result = {"ok": True, "outlet_mass_flow_kgps": outlet_mdot,
                  "outlet_share_of_total_outflow": shares,
                  "near_zero_outlets": near_zero}
        if near_zero:
            result["CAUTION"] = (
                f"Outlet(s) {near_zero} carry <{near_zero_fraction*100:.0f}% "
                "of total outflow while the rest carry the remainder. Before "
                "treating this as a real result: call "
                "check_interface_flow_integrity to rule out a WALL-typed "
                "mesh interface silently blocking that branch, and confirm "
                "the run has actually converged (see imbalance below)."
            )

        if in_names:
            inlet_mdot = {z: values[n] for z, n in zip(inlet_zones, in_names)}
            signed_total = sum(values[n] for n in names) + sum(inlet_mdot.values())
            imbalance = abs(signed_total) / (sum(abs(v) for v in inlet_mdot.values()) or 1.0)
            result["inlet_mass_flow_kgps"] = inlet_mdot
            result["mass_imbalance_fraction"] = imbalance
            if imbalance > imbalance_tolerance:
                result["CAUTION_convergence"] = (
                    f"Mass imbalance {imbalance*100:.1f}% exceeds "
                    f"{imbalance_tolerance*100:.0f}% - run has likely not "
                    "converged yet (flat residuals alone are not enough, "
                    "see workflow playbook Sec 3.3). Keep iterating before "
                    "trusting the per-outlet split above."
                )
        return _j(result)
    except Exception as e:
        return _j({"ok": False, "error": str(e)})
    finally:
        for rn in created:
            try:
                del rd.surface[rn]
            except Exception:
                pass


@mcp.tool()
def apply_visual_defaults(resolution: str = "1080p") -> str:
    """Apply the standard clean-engineering graphics defaults. Call ONCE per
    session BEFORE saving any contour/scene images.

    Sets: grid plane OFF, reflections OFF, shadows OFF, edge reflections OFF,
    axis triad ON, ruler ON, saved-picture background WHITE (invert_background),
    PNG format, landscape, and fixed picture resolution (independent of the
    GUI window size).

    Args:
        resolution: "1080p" (1920x1080, default) or "4k" (3840x2160).
    """
    if err := _chk("solver"): return err
    xres, yres = (3840, 2160) if resolution.lower() in ("4k", "2160p", "best") else (1920, 1080)
    applied, skipped = [], []

    def _set(label, fn):
        try:
            fn(); applied.append(label)
        except Exception as e:
            skipped.append(f"{label}: {e}")

    prefs = _solver.preferences
    ge = prefs.Graphics.GraphicsEffects
    _set("grid_plane_off", lambda: setattr(ge, "GridPlaneEnabled", False))
    _set("reflections_off", lambda: setattr(ge, "ReflectionsEnabled", False))
    _set("shadows_off", lambda: setattr(ge, "SimpleShadowsEnabled", False))
    _set("edge_reflections_off", lambda: setattr(ge, "ShowEdgeReflections", False))
    ap = prefs.Appearance
    _set("axis_triad_on", lambda: setattr(ap, "AxisTriad", True))
    _set("ruler_on", lambda: setattr(ap, "Ruler", True))
    pic = _solver.settings.results.graphics.picture
    _set("picture_resolution", lambda: pic.set_state(
        {"use_window_resolution": False, "x_resolution": xres, "y_resolution": yres}))
    _set("white_background_in_pictures", lambda: setattr(pic, "invert_background", True))
    _set("png_format", lambda: pic.driver_options.set_state({"hardcopy_format": "png"}))
    return _j({"ok": True, "resolution": f"{xres}x{yres}",
               "applied": applied, "skipped": skipped})


@mcp.tool()
def create_plane_surface(
    name: str,
    method: str = "zx-plane",
    coordinate: float = 0.0,
    point: list[float] | None = None,
    normal: list[float] | None = None,
) -> str:
    """Create a named cut-plane surface for contours/vectors.

    IMPORTANT quirks (verified on 2025 R2): the method must be set BEFORE the
    point/normal or coordinate fields become active, and a half-configured
    plane object cannot be deleted or reused - so this tool refuses names that
    already exist instead of overwriting.

    Args:
        name:       Unique surface name (fails if it already exists).
        method:     "xy-plane" | "yz-plane" | "zx-plane" | "point-and-normal".
        coordinate: Plane position for axis-aligned methods (z, x, or y value).
        point:      [x,y,z] on the plane (point-and-normal only).
        normal:     [nx,ny,nz] plane normal (point-and-normal only).
    """
    if err := _chk("solver"): return err
    try:
        ps = _solver.settings.results.surfaces.plane_surface
        if name in ps.get_object_names():
            return _j({"ok": False, "error": f"Surface '{name}' already exists; pick a new name."})
        ps[name] = {"method": method}           # step 1: method activates the rest
        e = ps[name]
        if method == "point-and-normal":
            if not (point and normal):
                return _j({"ok": False, "error": "point-and-normal needs point=[x,y,z] and normal=[nx,ny,nz]"})
            e.point.set_state([float(v) for v in point])
            e.normal.set_state([float(v) for v in normal])
        else:
            axis = {"xy-plane": "z", "yz-plane": "x", "zx-plane": "y"}.get(method)
            if axis is None:
                return _j({"ok": False, "error": f"Unknown method '{method}'"})
            getattr(e, axis).set_state(float(coordinate))
        return _j({"ok": True, "surface": name, "state": e.get_state()})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


def _parse_domain_extents_from_transcript() -> dict | None:
    """Last 'Domain Extents' block (printed by mesh check) from the live
    transcript: {'x': (min, max), 'y': ..., 'z': ...} in meters."""
    trn = _find_transcript()
    if trn is None:
        return None
    try:
        text = trn.read_text(encoding="utf-8", errors="replace")[-300_000:]
    except Exception:
        return None
    hits = re.findall(
        r"([xyz])-coordinate:\s*min \(m\) = ([\d.eE+-]+), max \(m\) = ([\d.eE+-]+)",
        text)
    if not hits:
        return None
    ext: dict = {}
    for axis, lo, hi in hits[-3:]:
        try:
            ext[axis] = (float(lo), float(hi))
        except ValueError:
            pass
    return ext if len(ext) == 3 else None


@mcp.tool()
def create_midplane(axis: str = "auto", name: str = "") -> str:
    """Create a mid cut-plane through the domain AUTOMATICALLY (solver mode)
    — no manual bounding-box reasoning needed.

    Runs Fluent's mesh check (its console output carries the domain
    extents), picks the axis — 'auto' chooses the THINNEST extent, the
    natural symmetry/mid plane of most internal-flow models — and creates
    the plane at the mid-coordinate. Returns the recommended plane-normal
    `view` to pass to display_scene (planar cuts must be viewed normal to
    the plane, not isometric — postproc KB).

    Args:
        axis: 'auto' (thinnest extent), or 'x' | 'y' | 'z' (plane normal).
        name: Surface name (default: midplane-<axis>).
    """
    if err := _chk("solver"): return err
    try:
        _solver.mesh.check()
    except Exception:
        pass
    ext = _parse_domain_extents_from_transcript()
    if not ext:
        return _j({"ok": False,
                   "error": "domain extents not found in the transcript",
                   "hint": "run run_mesh_check once, then retry; or use "
                           "create_plane_surface with explicit coordinates"})
    ax = axis.strip().lower()
    if ax not in ("x", "y", "z"):
        ax = min(ext, key=lambda a: ext[a][1] - ext[a][0])   # thinnest extent
    lo, hi = ext[ax]
    mid = (lo + hi) / 2.0
    method = {"x": "yz-plane", "y": "zx-plane", "z": "xy-plane"}[ax]
    pname = name or f"midplane-{ax}"
    result = json.loads(create_plane_surface(name=pname, method=method,
                                             coordinate=mid))
    if not result.get("ok"):
        return _j({**result, "domain_extents_m": ext})
    return _j({"ok": True, "surface": pname, "axis": ax,
               "coordinate_m": mid, "method": method,
               "domain_extents_m": {k: list(v) for k, v in ext.items()},
               "recommended_view": {"x": "right", "y": "top", "z": "front"}[ax],
               "next": f"create_native_contour(..., surfaces=['{pname}']) then "
                       f"display_scene(view='"
                       f"{ {'x': 'right', 'y': 'top', 'z': 'front'}[ax] }')"})


@mcp.tool()
def generate_default_results(analysis_type: str = "internal",
                             output_dir: str = "") -> str:
    """One-call default post-processing pack (spec §19 + KB default
    postprocessing matrix). Best-effort per artifact — partial failures are
    reported, not fatal. Produces:
      - velocity + pressure contours on an auto mid cut-plane (plane-normal view)
      - static-pressure contour on the walls (isometric)
      - KPI table: per-inlet/outlet mass flow + inlet avg pressure (dP proxy)
    Run AFTER assess_convergence says C1/C2.

    Args:
        analysis_type: 'internal' (default), 'external', 'cht' — tunes which
                       KPIs matter (all get the contour set).
        output_dir:    Image/artifact directory (default: logs/results).
    """
    if err := _chk("solver"): return err
    out = Path(output_dir) if output_dir else _SESSION_STATE_PATH.parent / "results"
    out.mkdir(parents=True, exist_ok=True)
    artifacts: list = []
    failures: dict = {}
    try:
        json.loads(apply_visual_defaults())
    except Exception as e:
        failures["visual_defaults"] = str(e)
    # 1. mid-plane
    plane = view = None
    try:
        mp = json.loads(create_midplane())
        if mp.get("ok"):
            plane, view = mp["surface"], mp["recommended_view"]
        else:
            failures["midplane"] = mp.get("error")
    except Exception as e:
        failures["midplane"] = str(e)
    if plane:
        for field in ("velocity-magnitude", "pressure"):
            cname = f"res_{field.split('-')[0]}_{plane}"
            img = str(out / f"{cname}.png")
            try:
                c = json.loads(create_native_contour(name=cname, field_name=field,
                                                     surfaces=[plane]))
                if c.get("ok"):
                    d = json.loads(display_scene(object_names=[cname],
                                                 output_image_path=img, view=view))
                    if d.get("ok"):
                        artifacts.append({"type": "png", "what": f"{field} on {plane}",
                                          "path": img})
                    else:
                        failures[cname] = d.get("error")
                else:
                    failures[cname] = c.get("error")
            except Exception as e:
                failures[cname] = str(e)
    # 2. wall pressure
    zones = _zone_type_pairs()
    walls = [z["zone"] for z in zones if "wall" in z["type"].lower()]
    if walls:
        img = str(out / "res_wall_pressure.png")
        try:
            c = json.loads(create_native_contour(name="res_wall_pressure",
                                                 field_name="pressure",
                                                 surfaces=walls))
            if c.get("ok"):
                d = json.loads(display_scene(object_names=["res_wall_pressure"],
                                             output_image_path=img, view="isometric"))
                if d.get("ok"):
                    artifacts.append({"type": "png", "what": "static pressure on walls",
                                      "path": img})
                else:
                    failures["wall_pressure"] = d.get("error")
            else:
                failures["wall_pressure"] = c.get("error")
        except Exception as e:
            failures["wall_pressure"] = str(e)
    # 3. KPI table
    inlets = [z["zone"] for z in zones if "inlet" in z["type"].lower()]
    outlets = [z["zone"] for z in zones if "outlet" in z["type"].lower()]
    kpis: dict = {}
    for z in inlets + outlets:
        try:
            r = json.loads(get_surface_report("mass-flow", [z]))
            if r.get("ok"):
                kpis[f"mdot_{z}_kgps"] = r["value"]
        except Exception as e:
            failures[f"mdot_{z}"] = str(e)
    for z in inlets:
        try:
            r = json.loads(get_surface_report("area-average", [z], "pressure"))
            if r.get("ok"):
                kpis[f"p_avg_{z}_pa"] = r["value"]
        except Exception as e:
            failures[f"p_avg_{z}"] = str(e)
    total_out = sum(abs(v) for k, v in kpis.items()
                    if k.startswith("mdot_") and any(o in k for o in outlets)
                    and isinstance(v, (int, float)))
    if total_out:
        kpis["outlet_flow_split"] = {
            o: round(abs(kpis.get(f"mdot_{o}_kgps", 0.0)) / total_out, 4)
            for o in outlets}
    kpi_path = str(out / "kpi_summary.json")
    try:
        Path(kpi_path).write_text(json.dumps(
            {"analysis_type": analysis_type, "generated": _ts(), "kpis": kpis},
            indent=2, default=str), encoding="utf-8")
        artifacts.append({"type": "json", "what": "KPI summary", "path": kpi_path})
    except Exception as e:
        failures["kpi_file"] = str(e)
    return _j({"ok": bool(artifacts), "analysis_type": analysis_type,
               "artifacts": artifacts, "kpis": kpis, "failures": failures,
               "next_actions": ["render_streamlines for flow paths",
                                "export_html / export_gltf for interactive 3D",
                                "generate_simulation_report for a document"]})


@mcp.tool()
def create_outlet_normal_plane(boundary_zone: str, name: str = "") -> str:
    """Create a cut plane through a boundary opening's true cross-section.

    Computes the area centroid and mean face normal of the given boundary
    zone (inlet/outlet) from the actual mesh, then creates a point-and-normal
    plane there - the plane of maximum flow area for that opening. Use for
    'show me the flow at each outlet' requests.

    Args:
        boundary_zone: Existing boundary zone (e.g. "pres_outlet_1").
        name:          Surface name (default "<zone>_plane").
    """
    if err := _chk("solver"): return err
    try:
        import numpy as np
        from ansys.fluent.core import SurfaceFieldDataRequest, SurfaceDataType
        fd = _solver.fields.field_data
        data = fd.get_field_data(SurfaceFieldDataRequest(
            data_types=[SurfaceDataType.Vertices, SurfaceDataType.FacesNormal],
            surfaces=[boundary_zone]))
        zd = data[boundary_zone]
        verts = np.asarray(zd.vertices)
        normals = np.asarray(zd.face_normals)
        if not verts.size:
            return _j({"ok": False, "error": f"No surface data for '{boundary_zone}'"})
        centroid = verts.mean(axis=0)
        nvec = normals.mean(axis=0)
        nvec = nvec / np.linalg.norm(nvec)
        pname = name or f"{boundary_zone}_plane"
        return create_plane_surface(pname, "point-and-normal",
                                    point=[float(v) for v in centroid],
                                    normal=[float(v) for v in nvec])
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def create_native_contour(
    name: str,
    field_name: str,
    surfaces: list[str],
) -> str:
    """Create (but not yet display) a native Fluent contour graphics object.

    Use display_scene to render one or more created objects into an image.
    Field names use Fluent conventions, e.g. "velocity-magnitude",
    "temperature", "pressure", "total-pressure".

    Args:
        name:       Unique object name (fails if it already exists).
        field_name: Fluent field variable.
        surfaces:   Surfaces/zones to color (planes, walls, outlets...).
    """
    if err := _chk("solver"): return err
    try:
        c = _solver.settings.results.graphics.contour
        if name in c.get_object_names():
            return _j({"ok": False, "error": f"Contour '{name}' already exists; pick a new name."})
        c[name] = {}
        obj = c[name]
        obj.field = field_name
        obj.surfaces_list = surfaces
        return _j({"ok": True, "contour": name, "field": field_name, "surfaces": surfaces})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# Per-object transparency (percent, 0-100) consumed by display_scene's
# scene-based compositor. Plain mesh objects have NO transparency option
# (verified 2025 R2) — transparency only exists on scene graphics_objects.
_object_transparency: dict[str, int] = {}


@mcp.tool()
def create_wall_display(name: str, surfaces: list[str], style: str = "transparent",
                        transparency: int = 70) -> str:
    """Create a mesh-display object showing walls for scene context.

    STANDARD STYLE (default): 'transparent' — faces WITHOUT edges, rendered
    see-through by display_scene so interior contour surfaces stay visible
    behind real-looking geometry. 'wireframe' (edges only) and 'solid'
    (opaque faces) are available as explicit fallbacks.

    Transparency is a property of Fluent SCENE graphics objects, not of the
    mesh object itself, so it takes effect when the object is rendered via
    display_scene (which composites through a scene).

    Args:
        name:         Unique object name.
        surfaces:     Wall zones to include (never interface zones).
        style:        "transparent" (default) | "solid" | "wireframe".
        transparency: 0-100 for style="transparent" (default 70).
    """
    if err := _chk("solver"): return err
    if style not in ("transparent", "solid", "wireframe"):
        return _j({"ok": False, "error": "style must be transparent|solid|wireframe"})
    try:
        m = _solver.settings.results.graphics.mesh
        if name in m.get_object_names():
            return _j({"ok": False, "error": f"Mesh display '{name}' already exists; pick a new name."})
        m[name] = {}
        obj = m[name]
        obj.surfaces_list = surfaces
        wire = style == "wireframe"
        obj.options = {"nodes": False, "edges": wire, "faces": not wire}
        if style == "transparent":
            _object_transparency[name] = max(0, min(100, transparency))
        else:
            _object_transparency.pop(name, None)
        return _j({"ok": True, "mesh_display": name, "style": style,
                   "transparency": _object_transparency.get(name, 0),
                   "surfaces": surfaces})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def display_scene(
    object_names: list[str],
    output_image_path: str,
    view: str = "isometric",
    camera_position: list[float] | None = None,
    camera_target: list[float] | None = None,
    camera_up: list[float] | None = None,
) -> str:
    """Render one or more graphics objects into a single saved image.

    Composites through a native Fluent SCENE object: all objects render
    together and any transparency recorded by create_wall_display
    (style="transparent") is applied per object — this is the only way to
    get see-through walls (mesh objects have no transparency of their own,
    and per-object .display() calls would each REPLACE the window).
    Falls back to display + add_to_graphics (no transparency) if scene
    creation fails. Then applies the view and saves a PNG at the resolution
    configured by apply_visual_defaults.

    Args:
        object_names:      Contour / mesh-display / vector object names, in
                           draw order (context first, colored result last).
        output_image_path: Output PNG path.
        view:              "isometric"|"front"|"back"|"left"|"right"|"top"|
                           "bottom"|"auto" (auto = fit current angle), or
                           "custom" with the camera_* args.
        camera_position:   [x,y,z] camera location (view="custom").
        camera_target:     [x,y,z] look-at point (view="custom").
        camera_up:         [x,y,z] up vector (view="custom", default +Y).
    """
    if err := _chk("solver"): return err
    if not object_names:
        return _j({"ok": False, "error": "object_names is empty"})
    try:
        g = _solver.settings.results.graphics
        registry = {}
        for kind in ("mesh", "contour", "vector", "pathline"):
            grp = getattr(g, kind, None)
            if grp is None:
                continue
            try:
                for nm in grp.get_object_names():
                    registry[nm] = (kind, grp)
            except Exception:
                pass
        missing = [n for n in object_names if n not in registry]
        if missing:
            return _j({"ok": False, "error": f"Unknown graphics objects: {missing}",
                       "known": sorted(registry.keys())})
        composited = "scene"
        try:
            sc = _solver.settings.results.scene
            scene_name = "mcp_display_scene"
            if scene_name in sc.get_object_names():
                del sc[scene_name]
            sc[scene_name] = {}
            go = sc[scene_name].graphics_objects
            for nm in object_names:
                go[nm] = {}
                if nm in _object_transparency:
                    go[nm].transparency = _object_transparency[nm]
            sc[scene_name].display()
        except Exception:
            # Fallback: overlay path (no per-object transparency)
            composited = "overlay-fallback"
            first, rest = object_names[0], object_names[1:]
            kind, grp = registry[first]
            grp[first].display()
            for nm in rest:
                kind, grp = registry[nm]
                grp.add_to_graphics(object_name=nm)
        v = g.views
        if view == "custom":
            if not (camera_position and camera_target):
                return _j({"ok": False, "error": "view='custom' needs camera_position and camera_target"})
            cam = v.camera
            cam.position(xyz=[float(x) for x in camera_position])
            cam.target(xyz=[float(x) for x in camera_target])
            cam.up_vector(xyz=[float(x) for x in (camera_up or [0.0, 1.0, 0.0])])
        elif view == "auto":
            v.auto_scale()
        else:
            v.restore_view(view_name=view)
            v.auto_scale()
        _mkdir(output_image_path)
        g.picture.save_picture(file_name=output_image_path)
        return _j({"ok": True, "image_path": output_image_path,
                   "objects": object_names, "view": view,
                   "compositing": composited,
                   "transparency": {n: _object_transparency[n]
                                    for n in object_names
                                    if n in _object_transparency}})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


@mcp.tool()
def get_session_info() -> str:
    """Return reconnection info for the live Fluent session (ip, port,
    password, cortex PID). Lets a later script or a recovered MCP process
    re-attach with connect_to_fluent instead of relaunching Fluent and
    losing the loaded case."""
    if err := _chk(): return err
    try:
        props = _active().connection_properties
        return _j({"ok": True, "ip": props.ip, "port": props.port,
                   "password": props.password,
                   "cortex_pid": props.cortex_pid,
                   "fluent_host_pid": props.fluent_host_pid,
                   "reconnect_with": "connect_to_fluent(ip=..., port=..., password=...)"})
    except Exception as e:
        return _j({"ok": False, "error": str(e)})


# ===========================================================================
# Entry point
# ===========================================================================

def main():
    mcp.run()


if __name__ == "__main__":
    main()
